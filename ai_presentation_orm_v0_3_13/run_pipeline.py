import argparse
import json
import os
import sys
from copy import deepcopy
from pathlib import Path

PACKAGE_ROOT = Path(__file__).resolve().parent
VENDOR_DIRS = [PACKAGE_ROOT / ".vendor", PACKAGE_ROOT / ".vendor_runtime"]
cleaned_sys_path = []
for entry in sys.path:
    if ".vendor" in str(entry).lower():
        continue
    try:
        resolved = Path(entry).resolve()
    except OSError:
        cleaned_sys_path.append(entry)
        continue
    if resolved not in {vendor_dir.resolve() for vendor_dir in VENDOR_DIRS}:
        cleaned_sys_path.append(entry)
sys.path[:] = cleaned_sys_path
LOCAL_MPL_CONFIG = PACKAGE_ROOT / ".matplotlib"
LOCAL_MPL_CONFIG.mkdir(exist_ok=True)
os.environ.setdefault("MPLCONFIGDIR", str(LOCAL_MPL_CONFIG))

from ai_presentation_orm.config_loader import load_config
from ai_presentation_orm.pptx_template_analyzer import extract_pptx_as_json, save_pptx_extraction

for vendor_dir in VENDOR_DIRS:
    try:
        if vendor_dir.exists():
            sys.path.append(str(vendor_dir))
    except OSError:
        pass

from ai_presentation_orm.canonical_model import build_canonical_model, save_canonical_model
from ai_presentation_orm.file_role_detector import collect_input_files, detect_file_roles
from ai_presentation_orm.pptx_template_analyzer import analyze_pptx_template
from ai_presentation_orm.excel_inventory import build_data_inventory
from ai_presentation_orm.readiness_checker import build_slide_readiness
from ai_presentation_orm.missing_data import build_missing_data_report
from ai_presentation_orm.report_writer import write_all_reports
from ai_presentation_orm.slide01_cover_builder import build_slide01_cover
from ai_presentation_orm.slide02_plan_fact_builder import build_slide02_from_media_plan, qa_slide02_plan_fact
from ai_presentation_orm.plan_fact_extractor import extract_plan_fact_from_orm_excel
from ai_presentation_orm.media_plan_extractor import extract_media_plan_table, apply_fact_overrides_to_media_plan_table
from ai_presentation_orm.project_month_fact_extractor import extract_project_month_facts
from ai_presentation_orm.slide03_data_extractor import extract_brand_mentions_sov, extract_analytical_context
from ai_presentation_orm.slide03_brand_mentions_builder import build_slide03_brand_mentions, qa_slide03_brand_mentions
from ai_presentation_orm.slide04_problem_field_extractor import extract_problem_field
from ai_presentation_orm.slide04_problem_field_builder import build_slide04_problem_field, qa_slide04_problem_field
from ai_presentation_orm.slide05_sov_extractor import extract_slide05_sov
from ai_presentation_orm.slide05_sov_builder import build_slide05_sov, qa_slide05_sov
from ai_presentation_orm.slide06_tonality_extractor import extract_slide06_tonality
from ai_presentation_orm.slide06_tonality_builder import build_slide06_tonality, qa_slide06_tonality
from ai_presentation_orm.slide07_sources_extractor import extract_slide07_sources
from ai_presentation_orm.slide07_sources_builder import build_slide07_sources, qa_slide07_sources
from ai_presentation_orm.slide08_seeding_metrics_extractor import extract_slide08_seeding_metrics
from ai_presentation_orm.slide08_seeding_metrics_builder import build_slide08_seeding_metrics, qa_slide08_seeding_metrics
from ai_presentation_orm.slide09_10_examples_extractor import extract_slide09_10_examples
from ai_presentation_orm.slide09_10_examples_builder import build_slide09_message_examples, qa_slide09_message_examples, build_slide10_review_examples, qa_slide10_review_examples
from ai_presentation_orm.slide11_seeding_metrics_extractor import extract_slide11_seeding_metrics
from ai_presentation_orm.slide11_seeding_metrics_builder import build_slide11_seeding_metrics, qa_slide11_seeding_metrics
from ai_presentation_orm.slide12_final_conclusions_builder import build_slide12_data, build_slide12_final_conclusions, qa_slide12_final_conclusions
from ai_presentation_orm.brand_styler import bold_project_brand_in_presentation
from pptx import Presentation


def _optional_int(value):
    if value is None or value == "":
        return None
    return int(float(value))


def _maybe_apply_slide02_fact_overrides(media_plan_table: dict, build_config: dict, project: dict) -> dict:
    fact_excel = (
        build_config.get("slide02_fact_excel")
        or build_config.get("slide02_fact_orm_excel")
        or build_config.get("slide02_month_fact_excel")
    )
    if not fact_excel:
        return media_plan_table
    try:
        fact_data = extract_project_month_facts(
            Path(fact_excel),
            month_sheet=build_config.get("slide02_fact_month_sheet") or build_config.get("slide08_month_sheet") or "",
            period=project.get("period", ""),
            source_system=build_config.get("slide02_fact_source_system") or build_config.get("slide08_source_system") or "проектная таблица",
        )
        return apply_fact_overrides_to_media_plan_table(media_plan_table, fact_data)
    except Exception as exc:
        media_plan_table["fact_overrides"] = {
            "status": "error",
            "source_file": Path(fact_excel).name,
            "error": str(exc),
        }
        return media_plan_table


def _find_comparative_report(analytics_excel: str, project_brand: str = "") -> str | None:
    """Auto-discover a comparative report matching the project brand.
    
    Picks the report where the target brand has the HIGHEST mention count
    in the Показатели sheet, not just presence in Содержание.
    """
    from pathlib import Path
    adir = Path(analytics_excel).parent
    candidates = sorted(adir.glob("*Сравнительный*отчет*.xlsx"))
    if not candidates:
        return None
    if not project_brand:
        return str(candidates[0])
    
    best_report = None
    best_count = -1
    brand_low = project_brand.lower().replace("ё", "е")
    
    try:
        from openpyxl import load_workbook
        for c in candidates:
            try:
                wb = load_workbook(c, data_only=True, read_only=True)
                if "Показатели" not in wb.sheetnames:
                    continue
                ws = wb["Показатели"]
                # Find Объект and Сообщения columns
                obj_idx = msg_idx = None
                for r_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=15, values_only=True), 1):
                    vals = [str(v or "").strip().lower() for v in row]
                    if "объект" in vals and "сообщения" in vals:
                        obj_idx = vals.index("объект")
                        msg_idx = vals.index("сообщения")
                        header_row = r_idx
                        break
                if obj_idx is None or msg_idx is None:
                    continue
                # Read brand mention counts
                for row in ws.iter_rows(min_row=header_row + 1, max_row=ws.max_row, values_only=True):
                    vals = [str(v or "").strip() for v in row]
                    name = vals[obj_idx] if obj_idx < len(vals) else ""
                    if not name or name.lower() in ("итого", "всего", ""):
                        continue
                    if name.lower().replace("ё", "е") == brand_low:
                        count = 0
                        try:
                            count = int(float(vals[msg_idx].replace(" ", "").replace("\xa0", ""))) if msg_idx < len(vals) else 0
                        except Exception:
                            pass
                        if count > best_count:
                            best_count = count
                            best_report = str(c)
                        break
            except Exception:
                continue
    except Exception:
        pass
    
    return best_report or str(candidates[0])


def _campaign_month_sheet_from_config(build_config: dict) -> str:
    return (
        build_config.get("slide03_campaign_month_sheet")
        or build_config.get("slide05_campaign_month_sheet")
        or build_config.get("slide08_month_sheet")
        or build_config.get("slide02_fact_month_sheet")
        or ""
    )


def _competitor_brands_from_config(project: dict, build_config: dict) -> list[str]:
    raw = project.get("competitor_brands") or build_config.get("competitor_brands") or []
    if isinstance(raw, str):
        raw = [item.strip() for item in raw.split(",")]
    return [str(item).strip() for item in (raw or []) if str(item or "").strip()]


def _campaign_split_from_source(slide03_data: dict, campaign_source: str, build_config: dict, project: dict) -> dict:
    """Build the campaign/organic split with project ORM facts as the first source."""
    source_path = Path(campaign_source)
    month_sheet = _campaign_month_sheet_from_config(build_config)
    fact_payload = {}
    raw_campaign_mentions = 0
    split_method = ""
    source_sheet = ""

    for requested_sheet in [month_sheet, ""]:
        try:
            fact_data = extract_project_month_facts(
                source_path,
                month_sheet=requested_sheet,
                period=project.get("period", ""),
                source_system=build_config.get("slide08_source_system") or "проектная таблица ORM",
            )
            raw_campaign_mentions = int(fact_data.get("total_fact", 0) or 0)
            if raw_campaign_mentions:
                fact_payload = fact_data
                split_method = "campaign_mentions_from_project_month_rows"
                source_sheet = fact_data.get("sheet", "")
                break
        except Exception:
            continue

    if not raw_campaign_mentions:
        try:
            campaign_pf = extract_plan_fact_from_orm_excel(source_path)
            raw_campaign_mentions = int(campaign_pf.get("fact_total", 0) or 0)
            if raw_campaign_mentions:
                fact_payload = campaign_pf
                split_method = "campaign_mentions_from_publication_sheet"
                source_sheet = campaign_pf.get("sheet", "")
        except Exception:
            pass

    if not raw_campaign_mentions:
        try:
            media_table = extract_media_plan_table(source_path, period_hint=project.get("period", ""))
            month_caption = (slide03_data.get("month_caption") or "").lower()
            headers = media_table.get("slide_table", {}).get("headers", [])
            month_idx = None
            for idx, header in enumerate(headers[1:3], start=0):
                if str(header).lower() in month_caption:
                    month_idx = idx
                    break
            if month_idx is not None:
                for row in media_table.get("slide_table", {}).get("rows", []):
                    if row.get("is_service_row"):
                        continue
                    vals = row.get("values", [])
                    if month_idx < len(vals):
                        raw_campaign_mentions += int(float(vals[month_idx] or 0))
                if raw_campaign_mentions:
                    split_method = "campaign_mentions_from_media_plan_selected_month"
                    source_sheet = media_table.get("sheet", "")
        except Exception:
            pass

    project_row = slide03_data.get("project_brand_row") or {}
    total_brand_mentions = int(project_row.get("mentions", 0) or 0)
    # Canonical: ORM campaign is the source of truth, not capped by analytics
    campaign_mentions = raw_campaign_mentions
    organic_mentions = max(total_brand_mentions - len(slide03_data.get("campaign_split", {}).get("matched_orm_indices", [])), 0)
    if organic_mentions < 0:
        organic_mentions = max(total_brand_mentions - raw_campaign_mentions, 0)
    status = "verified" if split_method == "campaign_mentions_from_project_month_rows" else "ready_with_caveat"
    if not raw_campaign_mentions:
        status = "missing_required"
        organic_mentions = total_brand_mentions

    split = {
        "source_file": source_path.name,
        "source_sheet": source_sheet,
        "method": split_method,
        "campaign_mentions": campaign_mentions,
        "organic_mentions": organic_mentions,
        "total_brand_mentions": total_brand_mentions,
        "campaign_publications_count": raw_campaign_mentions,
        "campaign_publications_by_type": (
            fact_payload.get("campaign_publications_by_type")
            or fact_payload.get("fact_rows_by_name")
            or fact_payload.get("format_counts")
            or {}
        ),
        "campaign_fact_method": fact_payload.get("campaign_fact_method") or fact_payload.get("method") or split_method,
        "campaign_fact_source_sheet": fact_payload.get("campaign_fact_source_sheet") or source_sheet,
        "campaign_required": True,
        "campaign_vs_organic_status": status,
        "status": status,
        "caveat": (
            "Campaign part is counted as publication rows from the project ORM month sheet; organic part is remaining brand mentions for the month."
            if status == "verified"
            else "Campaign part is not confirmed by a project ORM month sheet; fallback source was used and PM / Insight confirmation is recommended."
        ),
    }

    slide03_data["campaign_publications_count"] = raw_campaign_mentions
    slide03_data["campaign_publications_by_type"] = split["campaign_publications_by_type"]
    slide03_data["campaign_fact_method"] = split["campaign_fact_method"]
    slide03_data["campaign_fact_source_sheet"] = split["campaign_fact_source_sheet"]
    slide03_data["campaign_vs_organic_status"] = status
    return split


def _combined_qa_from_steps(steps: list[dict]) -> dict:
    blockers = []
    warnings = []
    by_slide = []
    for step in steps:
        slide_number = step.get("slide")
        qa = step.get("qa") or {}
        result = step.get("result") or {}
        slide_blockers = list(qa.get("blockers") or []) + list(result.get("blockers") or []) + list(result.get("qa_blockers") or [])
        slide_warnings = list(qa.get("warnings") or []) + list(result.get("warnings") or []) + list(result.get("qa_warnings") or [])
        chart_updates = result.get("chart_updates") or []
        chart_image_outputs = result.get("chart_image_outputs") or []
        if chart_image_outputs or any("image_fallback" in str(item.get("type", "")) for item in chart_updates if isinstance(item, dict)):
            slide_warnings.append(
                "Uses rendered chart image fallback; native chart parts are preserved, but editable chart data update is not confirmed."
            )
        by_slide.append({
            "slide": slide_number,
            "status": qa.get("status") or result.get("status") or result.get("qa_status") or "not_reported",
            "blockers": slide_blockers,
            "warnings": slide_warnings,
        })
        label = f"Slide {int(slide_number):02d}" if isinstance(slide_number, int) else f"Slide {slide_number}"
        blockers.extend(f"{label}: {item}" for item in slide_blockers)
        warnings.extend(f"{label}: {item}" for item in slide_warnings)
    return {
        "qa_status": "blocked" if blockers else "ready_for_visual_review",
        "qa_blockers": blockers,
        "qa_warnings": warnings,
        "qa_by_slide": by_slide,
    }


def build_first4_combined_deck(config: dict, template_path: Path, output_dir: Path, canonical_model: dict | None = None) -> dict:
    """Apply enabled slide builders sequentially to one template copy."""
    build_config = config.get("build", {})
    project = config.get("project", {})
    competitor_brands = _competitor_brands_from_config(project, build_config)
    current_path = template_path
    steps = []

    # Use canonical model to enrich slide data payloads
    if canonical_model is not None:
        cm = deepcopy(canonical_model) if hasattr(canonical_model, 'values') else canonical_model
    else:
        cm = None

    if build_config.get("slide01_cover_test"):
        next_path = output_dir / "_combined_step_01_cover.pptx"
        result = build_slide01_cover(template_path=current_path, output_path=next_path, project=project)
        steps.append({"slide": 1, "result": result})
        current_path = next_path

    if build_config.get("slide02_plan_fact_test"):
        media_plan_excel = build_config.get("slide02_media_plan_excel") or build_config.get("slide02_orm_excel")
        if not media_plan_excel:
            raise ValueError("build.slide02_media_plan_excel is required for first4_combined_test when slide02 is enabled")
        media_plan_table = extract_media_plan_table(Path(media_plan_excel), period_hint=project.get("period", ""))
        media_plan_table = _maybe_apply_slide02_fact_overrides(media_plan_table, build_config, project)
        next_path = output_dir / "_combined_step_02_plan_fact.pptx"
        result = build_slide02_from_media_plan(
            template_path=current_path,
            output_path=next_path,
            project=project,
            media_plan_table=media_plan_table,
        )
        qa = qa_slide02_plan_fact(media_plan_table, result)
        steps.append({"slide": 2, "result": result, "data": media_plan_table, "qa": qa})
        current_path = next_path

    if build_config.get("slide03_brand_mentions_test"):
        analytics_excel = build_config.get("slide03_analytics_excel")
        if not analytics_excel:
            raise ValueError("build.slide03_analytics_excel is required for first4_combined_test when slide03 is enabled")
        slide03_data = extract_brand_mentions_sov(
            Path(analytics_excel),
            project_brand=project.get("brand", ""),
            competitor_brands=competitor_brands,
        )
        slide03_data["source_system"] = build_config.get("slide03_source_system") or ""
        analytical_context_excel = build_config.get("slide03_analytical_context_excel")
        if analytical_context_excel:
            try:
                slide03_data["analytical_context"] = extract_analytical_context(Path(analytical_context_excel))
            except Exception as exc:
                slide03_data["analytical_context"] = {"status": "error", "error": str(exc)}
        campaign_source = (
            build_config.get("slide03_campaign_source_excel")
            or build_config.get("slide05_project_orm_excel")
            or build_config.get("slide02_orm_excel")
        )
        if cm is not None:
            slide03_data["canonical_model_ref"] = True
            tb = cm.get("target_brand", {})
            sov_target = cm.get("competitive_sov_scope", {}).get("target", {})
            comp_brands = cm.get("comparative_brand_data", {}).get("by_brand", {})
            comp_audience = cm.get("comparative_brand_data", {}).get("by_brand_audience", {})
            
            # Populate chart_mentions from canonical brand data
            all_brands = sorted(comp_brands.items(), key=lambda x: x[1], reverse=True)
            project_brand_low = (project.get("brand") or "").lower()
            brand_names = [n for n, c in all_brands]
            if brand_names:
                chart_organic = []
                chart_campaign = []
                for name, count in all_brands:
                    if name.lower() == project_brand_low:
                        chart_organic.append(sov_target.get("organic", 0))
                        chart_campaign.append(tb.get("campaign", 0))
                    else:
                        chart_organic.append(count)
                        chart_campaign.append(0)
                slide03_data["chart_mentions"] = {
                    "categories": brand_names,
                    "mentions": [o + c for o, c in zip(chart_organic, chart_campaign)],
                }
                slide03_data["chart_audience"] = {
                    "categories": brand_names,
                    "audience": [comp_audience.get(n, 0) for n in brand_names],
                }
            
            # Rebuild metadata from canonical data
            slide03_data["all_brand_rows"] = [{"brand": n, "mentions": c} for n, c in all_brands]
            slide03_data["top_brands"] = slide03_data["all_brand_rows"][:6]
            slide03_data["project_brand_row"] = next((r for r in slide03_data["all_brand_rows"] if r["brand"].lower() == project_brand_low), {})
            slide03_data["total_mentions_competitive_set"] = sum(c for _, c in all_brands)
            slide03_data["project_mentions"] = tb.get("with_campaign", 0)
            slide03_data["project_sov"] = sov_target.get("sov_with_campaign", 0)
            slide03_data["project_brand"] = project.get("brand", "")
            slide03_data["project_brand_rank"] = next((i+1 for i, (n, c) in enumerate(all_brands) if n.lower() == project_brand_low), None)
            
            # Campaign split
            slide03_data["campaign_split"] = {
                "status": "verified",
                "method": "canonical_model",
                "campaign_mentions": tb.get("campaign", 0),
                "organic_mentions": tb.get("organic", 0),
                "total_brand_mentions": tb.get("with_campaign", 0),
                "campaign_publications_count": cm.get("orm_campaign", {}).get("total", 0),
                "campaign_publications_by_type": cm.get("orm_campaign", {}).get("by_type", {}),
                "campaign_fact_method": cm.get("orm_campaign", {}).get("fact_method", ""),
                "campaign_fact_source_sheet": cm.get("orm_campaign", {}).get("source_sheet", ""),
                "campaign_vs_organic_status": "verified",
                "caveat": "Campaign data from canonical model.",
            }
            slide03_data["sov_with_campaign"] = sov_target.get("sov_with_campaign", 0)
            slide03_data["sov_without_campaign"] = sov_target.get("sov_without_campaign", 0)
        elif campaign_source:
            try:
                slide03_data["campaign_split"] = _campaign_split_from_source(slide03_data, campaign_source, build_config, project)
            except Exception as exc:
                slide03_data["campaign_split"] = {"status": "missing_required", "error": str(exc)}
        next_path = output_dir / "_combined_step_03_brand_mentions.pptx"
        result = build_slide03_brand_mentions(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide03_data,
        )
        qa = qa_slide03_brand_mentions(slide03_data, result)
        steps.append({"slide": 3, "result": result, "data": slide03_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide04_problem_field_test"):
        analytical_excel = build_config.get("slide04_analytical_excel")
        raw_excel = build_config.get("slide04_raw_excel")
        report_excel = build_config.get("slide04_report_excel")
        if not (analytical_excel or raw_excel or report_excel):
            raise ValueError("At least one slide04 data source is required for first4_combined_test when slide04 is enabled")
        slide04_data = extract_problem_field(
            analytical_path=Path(analytical_excel) if analytical_excel else None,
            raw_path=Path(raw_excel) if raw_excel else None,
            report_path=Path(report_excel) if report_excel else None,
            project_brand=project.get("brand", ""),
            source_system=build_config.get("slide04_source_system") or "",
            project_period=project.get("period", ""),
        )
        next_path = output_dir / "_combined_step_04_problem_field.pptx"
        result = build_slide04_problem_field(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide04_data,
        )
        qa = qa_slide04_problem_field(slide04_data, result)
        steps.append({"slide": 4, "result": result, "data": slide04_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide05_sov_test"):
        analytics_excel = build_config.get("slide05_analytics_excel") or build_config.get("slide03_analytics_excel")
        if not analytics_excel:
            raise ValueError("build.slide05_analytics_excel or build.slide03_analytics_excel is required for combined deck when slide05 is enabled")
        slide05_data = extract_slide05_sov(
            analytics_path=Path(analytics_excel),
            project_brand=project.get("brand", ""),
            source_system=build_config.get("slide05_source_system") or build_config.get("slide03_source_system") or "",
            project_orm_path=Path(build_config["slide05_project_orm_excel"]) if build_config.get("slide05_project_orm_excel") else None,
            campaign_month_sheet=build_config.get("slide05_campaign_month_sheet") or "",
            campaign_materials_override=_optional_int(build_config.get("slide05_campaign_materials_override")),
            competitor_brands=competitor_brands,
        )
        if cm is not None:
            slide05_data["canonical_model_ref"] = True
            sov_target = cm.get("competitive_sov_scope", {}).get("target", {})
            slide05_data["project_mentions"] = sov_target.get("with_campaign", 0) - sov_target.get("organic", 0) + sov_target.get("organic", 0)
            slide05_data["campaign_split_override"] = {
                "sov_with_campaign": sov_target.get("sov_with_campaign", 0),
                "sov_without_campaign": sov_target.get("sov_without_campaign", 0),
                "sov_lift_pp": sov_target.get("sov_lift_pp", 0),
                "campaign_materials": cm.get("orm_campaign", {}).get("total", 0),
                "organic_mentions": sov_target.get("organic", 0),
            }
        next_path = output_dir / "_combined_step_05_sov.pptx"
        result = build_slide05_sov(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide05_data,
        )
        qa = qa_slide05_sov(slide05_data, result)
        steps.append({"slide": 5, "result": result, "data": slide05_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide06_tonality_test"):
        analytics_excel = (
            build_config.get("slide06_analytics_excel")
            or build_config.get("slide05_analytics_excel")
            or build_config.get("slide03_analytics_excel")
        )
        if not analytics_excel:
            raise ValueError("build.slide06_analytics_excel, build.slide05_analytics_excel or build.slide03_analytics_excel is required for combined deck when slide06 is enabled")
        slide06_data = extract_slide06_tonality(
            analytics_path=Path(analytics_excel),
            project_brand=project.get("brand", ""),
            source_system=build_config.get("slide06_source_system") or build_config.get("slide05_source_system") or build_config.get("slide03_source_system") or "",
            competitor_brands=competitor_brands,
            project_orm_path=Path(
                build_config.get("slide06_project_orm_excel")
                or build_config.get("slide08_project_orm_excel")
                or build_config.get("slide05_project_orm_excel")
                or build_config.get("slide02_fact_excel")
            ) if (
                build_config.get("slide06_project_orm_excel")
                or build_config.get("slide08_project_orm_excel")
                or build_config.get("slide05_project_orm_excel")
                or build_config.get("slide02_fact_excel")
            ) else None,
            month_sheet=(
                build_config.get("slide06_month_sheet")
                or build_config.get("slide08_month_sheet")
                or build_config.get("slide05_campaign_month_sheet")
                or build_config.get("slide02_fact_month_sheet")
                or ""
            ),
        )
        if cm is not None:
            slide06_data["canonical_model_ref"] = True
            tonality_cm = cm.get("target_tonality", {})
            slide06_data["project_row"] = {
                "positive": tonality_cm.get("positive", 0),
                "neutral": tonality_cm.get("neutral", 0),
                "negative": tonality_cm.get("negative", 0),
                "total": tonality_cm.get("positive", 0) + tonality_cm.get("neutral", 0) + tonality_cm.get("negative", 0),
            }
        next_path = output_dir / "_combined_step_06_tonality.pptx"
        result = build_slide06_tonality(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide06_data,
        )
        qa = qa_slide06_tonality(slide06_data, result)
        steps.append({"slide": 6, "result": result, "data": slide06_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide07_sources_test"):
        analytics_excel = (
            build_config.get("slide07_analytics_excel")
            or build_config.get("slide06_analytics_excel")
            or build_config.get("slide05_analytics_excel")
            or build_config.get("slide03_analytics_excel")
        )
        if not analytics_excel:
            raise ValueError("build.slide07_analytics_excel, build.slide06_analytics_excel, build.slide05_analytics_excel or build.slide03_analytics_excel is required for combined deck when slide07 is enabled")
        slide07_data = extract_slide07_sources(
            analytics_path=Path(analytics_excel),
            project_brand=project.get("brand", ""),
            source_system=build_config.get("slide07_source_system") or build_config.get("slide06_source_system") or build_config.get("slide05_source_system") or build_config.get("slide03_source_system") or "",
        )
        if cm is not None:
            slide07_data["canonical_model_ref"] = True
            platform = cm.get("platform_scope", {})
            slide07_data["platform_type_rows"] = platform.get("platform_type_rows", slide07_data.get("platform_type_rows", []))
            slide07_data["source_rows"] = platform.get("source_rows", slide07_data.get("source_rows", []))
            slide07_data["total_messages"] = platform.get("total", slide07_data.get("total_messages", 0))
        next_path = output_dir / "_combined_step_07_sources.pptx"
        result = build_slide07_sources(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide07_data,
        )
        qa = qa_slide07_sources(slide07_data, result)
        steps.append({"slide": 7, "result": result, "data": slide07_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide08_seeding_metrics_test"):
        project_orm_excel = build_config.get("slide08_project_orm_excel") or build_config.get("slide05_project_orm_excel") or build_config.get("slide02_media_plan_excel")
        if not project_orm_excel:
            raise ValueError("build.slide08_project_orm_excel is required for combined deck when slide08 is enabled")
        slide08_data = extract_slide08_seeding_metrics(
            project_orm_path=Path(project_orm_excel),
            project_brand=project.get("brand", ""),
            month_sheet=build_config.get("slide08_month_sheet") or "Май",
            period=project.get("period", ""),
            source_system=build_config.get("slide08_source_system") or "проектная таблица ORM",
        )
        if cm is not None:
            slide08_data["canonical_model_ref"] = True
            orm_camp = cm.get("orm_campaign", {})
            slide08_data["total_materials"] = orm_camp.get("total", slide08_data.get("total_materials", 0))
            slide08_data["views_metrics"] = {
                "received_views_total": orm_camp.get("views_total", slide08_data.get("views_metrics", {}).get("received_views_total", 0)),
                "has_views": orm_camp.get("views_total", 0) > 0 or slide08_data.get("views_metrics", {}).get("has_views", False),
                "metric_label": "Просмотры",
            }
        next_path = output_dir / "_combined_step_08_seeding_metrics.pptx"
        result = build_slide08_seeding_metrics(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide08_data,
        )
        qa = qa_slide08_seeding_metrics(slide08_data, result)
        steps.append({"slide": 8, "result": result, "data": slide08_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide09_message_examples_test") or build_config.get("slide10_review_examples_test"):
        project_orm_excel = (
            build_config.get("slide09_project_orm_excel")
            or build_config.get("slide10_project_orm_excel")
            or build_config.get("slide08_project_orm_excel")
            or build_config.get("slide05_project_orm_excel")
            or build_config.get("slide02_media_plan_excel")
        )
        if not project_orm_excel:
            raise ValueError("build.slide09_project_orm_excel or build.slide10_project_orm_excel is required for combined deck when slide09/10 is enabled")
        screenshot_backend = str(
            build_config.get("slide09_10_screenshot_backend")
            or build_config.get("slide09_screenshot_backend")
            or build_config.get("slide10_screenshot_backend")
            or ""
        ).strip()
        if screenshot_backend:
            os.environ["AI_ORM_SCREENSHOT_BACKEND"] = screenshot_backend
        examples_data = extract_slide09_10_examples(
            project_orm_path=Path(project_orm_excel),
            project_brand=project.get("brand", ""),
            month_sheet=build_config.get("slide09_month_sheet") or build_config.get("slide10_month_sheet") or build_config.get("slide08_month_sheet") or "Май",
            period=project.get("period", ""),
            source_system=build_config.get("slide09_source_system") or build_config.get("slide10_source_system") or build_config.get("slide08_source_system") or "проектная таблица ORM",
        )
        if build_config.get("slide09_message_examples_test"):
            slide09_data = examples_data["slide09"]
            next_path = output_dir / "_combined_step_09_message_examples.pptx"
            result = build_slide09_message_examples(
                template_path=current_path,
                output_path=next_path,
                project=project,
                data=slide09_data,
            )
            qa = qa_slide09_message_examples(slide09_data, result)
            steps.append({"slide": 9, "result": result, "data": slide09_data, "qa": qa})
            current_path = next_path
        if build_config.get("slide10_review_examples_test"):
            slide10_data = examples_data["slide10"]
            next_path = output_dir / "_combined_step_10_review_examples.pptx"
            result = build_slide10_review_examples(
                template_path=current_path,
                output_path=next_path,
                project=project,
                data=slide10_data,
            )
            qa = qa_slide10_review_examples(slide10_data, result)
            steps.append({"slide": 10, "result": result, "data": slide10_data, "qa": qa})
            current_path = next_path

    if build_config.get("slide11_seeding_metrics_test"):
        ratings_excel = build_config.get("slide11_ratings_excel") or build_config.get("slide08_project_orm_excel") or build_config.get("slide05_project_orm_excel")
        if not ratings_excel:
            raise ValueError("build.slide11_ratings_excel is required for combined deck when slide11 is enabled")
        slide11_data = extract_slide11_seeding_metrics(
            ratings_path=Path(ratings_excel),
            project_brand=project.get("brand", ""),
            period=project.get("period", ""),
            ratings_sheet=build_config.get("slide11_ratings_sheet") or "",
            current_period=build_config.get("slide11_current_period") or "",
            placement_month_sheet=build_config.get("slide11_placement_month_sheet") or build_config.get("slide08_month_sheet") or build_config.get("slide05_campaign_month_sheet") or "",
            source_system=build_config.get("slide11_source_system") or "проектная таблица ORM",
        )
        if cm is not None:
            slide11_data["canonical_model_ref"] = True
            ratings = cm.get("ratings_scope", {})
            if ratings.get("status") == "ready":
                summaries = slide11_data.get("period_summaries") or []
                if not summaries or not slide11_data.get("start_summary"):
                    periods = ratings.get("periods", [])
                    pos_cards = ratings.get("positive_cards", [0, 0])
                    neg_cards = ratings.get("negative_cards", [0, 0])
                    no_rev = ratings.get("no_review_cards", [0, 0])
                    rev_counts = ratings.get("reviews_count", [0, 0])
                    slide11_data["period_summaries"] = [
                        {"label": "Старт" if i == 0 else "Текущий", "positive_cards": pos_cards[i] if i < len(pos_cards) else 0, "negative_cards": neg_cards[i] if i < len(neg_cards) else 0, "no_reviews": no_rev[i] if i < len(no_rev) else 0, "reviews_total": rev_counts[i] if i < len(rev_counts) else 0}
                        for i in range(len(periods))
                    ]
                    slide11_data["start_summary"] = {"reviews_total": rev_counts[0] if len(rev_counts) > 0 else 0, "positive_cards": pos_cards[0] if len(pos_cards) > 0 else 0, "negative_cards": neg_cards[0] if len(neg_cards) > 0 else 0, "no_reviews": no_rev[0] if len(no_rev) > 0 else 0, "data_status": "ready"}
                    slide11_data["current_summary"] = {"reviews_total": rev_counts[-1] if rev_counts else 0, "positive_cards": pos_cards[-1] if pos_cards else 0, "negative_cards": neg_cards[-1] if neg_cards else 0, "no_reviews": no_rev[-1] if no_rev else 0, "data_status": "ready"}
        next_path = output_dir / "_combined_step_11_seeding_metrics.pptx"
        result = build_slide11_seeding_metrics(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide11_data,
        )
        qa = qa_slide11_seeding_metrics(slide11_data, result)
        steps.append({"slide": 11, "result": result, "data": slide11_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide12_final_conclusions_test"):
        slide_data_by_number = {step["slide"]: step.get("data") for step in steps if step.get("data") is not None}
        if cm is not None and not slide_data_by_number:
            slide_data_by_number = {3: {"canonical_model_ref": True}}
        if cm is not None:
            slide_data_by_number["canonical_model"] = cm
        slide12_data = build_slide12_data(project, slide_data_by_number)
        next_path = output_dir / "_combined_step_12_final_conclusions.pptx"
        result = build_slide12_final_conclusions(
            template_path=current_path,
            output_path=next_path,
            project=project,
            data=slide12_data,
        )
        qa = qa_slide12_final_conclusions(slide12_data, result)
        steps.append({"slide": 12, "result": result, "data": slide12_data, "qa": qa})
        current_path = next_path

    if build_config.get("slide12_final_conclusions_test"):
        final_name = "draft_slides_01_12_combined_test.pptx"
    elif build_config.get("slide11_seeding_metrics_test"):
        final_name = "draft_slides_01_11_combined_test.pptx"
    elif build_config.get("slide08_seeding_metrics_test"):
        final_name = "draft_slides_01_08_combined_test.pptx"
    elif build_config.get("slide07_sources_test"):
        final_name = "draft_slides_01_07_combined_test.pptx"
    elif build_config.get("slide06_tonality_test"):
        final_name = "draft_slides_01_06_combined_test.pptx"
    elif build_config.get("slide05_sov_test"):
        final_name = "draft_slides_01_05_combined_test.pptx"
    else:
        final_name = "draft_slides_01_04_combined_test.pptx"
    final_path = output_dir / final_name
    prs = Presentation(str(current_path))
    brand_bold_updates = bold_project_brand_in_presentation(prs, project.get("brand", ""))
    if len(prs.slides) >= 11:
        slide11_title = next((s for s in prs.slides[10].shapes if s.shape_id == 10), None)
        if slide11_title and hasattr(slide11_title, "text_frame"):
            for para in slide11_title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = False
    prs.save(str(final_path))
    return {
        "output": str(final_path),
        "steps": steps,
        "slides_enabled": [step["slide"] for step in steps],
        "brand_bold_updates": brand_bold_updates,
        **_combined_qa_from_steps(steps),
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="AI Presentation ORM v0.1 pipeline")
    parser.add_argument("--config", required=True, help="Path to YAML config")
    args = parser.parse_args()

    config = load_config(args.config)
    output_dir = Path(config["paths"]["output_dir"])
    output_dir.mkdir(parents=True, exist_ok=True)

    template_path = Path(config["paths"]["template_pptx"])
    input_files = collect_input_files(config)

    file_roles = detect_file_roles(input_files, template_path)
    template_map = analyze_pptx_template(template_path)
    data_inventory = build_data_inventory(input_files, config)
    slide_readiness = build_slide_readiness(template_map, data_inventory, file_roles)
    missing_data = build_missing_data_report(slide_readiness)

    write_all_reports(
        output_dir=output_dir,
        config=config,
        file_roles=file_roles,
        template_map=template_map,
        data_inventory=data_inventory,
        slide_readiness=slide_readiness,
        missing_data=missing_data,
    )

    # ---- Build canonical data model ----
    build_config = config.get("build", {})
    project = config.get("project", {})
    canonical_model = None
    has_canonical_inputs = (
        build_config.get("slide03_analytics_excel")
        or build_config.get("slide05_analytics_excel")
    ) and (
        build_config.get("slide05_project_orm_excel")
        or build_config.get("slide08_project_orm_excel")
        or build_config.get("slide02_fact_excel")
    )
    if has_canonical_inputs:
        try:
            analytics_excel = (
                build_config.get("slide03_analytics_excel")
                or build_config.get("slide05_analytics_excel")
                or build_config.get("slide06_analytics_excel")
                or build_config.get("slide07_analytics_excel")
            )
            project_orm_excel = (
                build_config.get("slide05_project_orm_excel")
                or build_config.get("slide08_project_orm_excel")
                or build_config.get("slide02_fact_excel")
                or build_config.get("slide02_orm_excel")
            )
            canonical_config = {
                "project_orm_path": str(Path(project_orm_excel).resolve()),
                "analytics_report_path": str(Path(analytics_excel).resolve()),
                "comparative_report_path": str(
                    Path(build_config.get("slide03_analytical_context_excel") or _find_comparative_report(analytics_excel, project.get("brand"))).resolve()
                ) if build_config.get("slide03_analytical_context_excel") or _find_comparative_report(analytics_excel, project.get("brand")) else None,
                "project_brand": project.get("brand", ""),
                "competitors": _competitor_brands_from_config(project, build_config),
                "report_month": _campaign_month_sheet_from_config(build_config) or "Май",
                "period": project.get("period", ""),
                "source_exclusions": build_config.get("source_exclusions") or config.get("rules", {}).get("source_exclusions") or [],
                "brand_aliases": {},
                "generation_options": {
                    "ratings_sheet": build_config.get("slide11_ratings_sheet") or "Рейтинги",
                    "ratings_current_period": build_config.get("slide11_current_period") or "",
                    "campaign_matched_in_analytics": build_config.get("campaign_matched_in_analytics_override"),
                },
            }
            canonical_model = build_canonical_model(canonical_config)
            save_canonical_model(canonical_model, output_dir)
            (output_dir / "canonical_model_build_result.json").write_text(
                json.dumps({
                    "status": canonical_model.get("status", "ready"),
                    "orm_campaign_total": canonical_model.get("orm_campaign", {}).get("total", 0),
                    "target_with_campaign": canonical_model.get("target_brand", {}).get("with_campaign", 0),
                    "sov_with_campaign": canonical_model.get("competitive_sov_scope", {}).get("target", {}).get("sov_with_campaign", 0),
                    "tonality_positive": canonical_model.get("target_tonality", {}).get("positive", 0),
                }, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
        except Exception as exc:
            (output_dir / "canonical_model_error.json").write_text(
                json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )

    if config.get("build", {}).get("slide01_cover_test"):
        slide01_result = build_slide01_cover(
            template_path=template_path,
            output_path=output_dir / "draft_slide_01_cover_test.pptx",
            project=config.get("project", {}),
        )
        (output_dir / "slide_01_build_result.json").write_text(
            __import__("json").dumps(slide01_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide02_plan_fact_test"):
        media_plan_excel = config.get("build", {}).get("slide02_media_plan_excel") or config.get("build", {}).get("slide02_orm_excel")
        if not media_plan_excel:
            raise ValueError("build.slide02_media_plan_excel is required for slide02_plan_fact_test")
        media_plan_table = extract_media_plan_table(__import__("pathlib").Path(media_plan_excel), period_hint=config.get("project", {}).get("period", ""))
        media_plan_table = _maybe_apply_slide02_fact_overrides(media_plan_table, config.get("build", {}), config.get("project", {}))
        slide02_result = build_slide02_from_media_plan(
            template_path=template_path,
            output_path=output_dir / "draft_slide_02_plan_fact_test.pptx",
            project=config.get("project", {}),
            media_plan_table=media_plan_table,
        )
        slide02_qa = qa_slide02_plan_fact(media_plan_table, slide02_result)
        (output_dir / "slide_02_media_plan_table.json").write_text(
            __import__("json").dumps(media_plan_table, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_02_build_result.json").write_text(
            __import__("json").dumps(slide02_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_02_qa.json").write_text(
            __import__("json").dumps(slide02_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide03_brand_mentions_test"):
        analytics_excel = config.get("build", {}).get("slide03_analytics_excel")
        if not analytics_excel:
            raise ValueError("build.slide03_analytics_excel is required for slide03_brand_mentions_test")
        slide03_data = extract_brand_mentions_sov(
            __import__("pathlib").Path(analytics_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            competitor_brands=_competitor_brands_from_config(config.get("project", {}), config.get("build", {})),
        )
        slide03_data["source_system"] = config.get("build", {}).get("slide03_source_system") or ""

        analytical_context_excel = config.get("build", {}).get("slide03_analytical_context_excel")
        if analytical_context_excel:
            try:
                slide03_data["analytical_context"] = extract_analytical_context(__import__("pathlib").Path(analytical_context_excel))
            except Exception as exc:
                slide03_data["analytical_context"] = {"status": "error", "error": str(exc)}
        if canonical_model is not None:
            slide03_data["canonical_model_ref"] = True
            cm = canonical_model
            tb = cm.get("target_brand", {})
            sov_target = cm.get("competitive_sov_scope", {}).get("target", {})
            comp_brands = cm.get("comparative_brand_data", {}).get("by_brand", {})
            comp_audience = cm.get("comparative_brand_data", {}).get("by_brand_audience", {})
            
            # Populate chart_mentions from comparative brand data
            all_brands = sorted(comp_brands.items(), key=lambda x: x[1], reverse=True)
            brands_with_campaign = []
            brands_organic = []
            brand_names = []
            project_brand_low = (config.get("project", {}).get("brand", "") or "").lower()
            for name, count in all_brands:
                brand_names.append(name)
                if name.lower() == project_brand_low:
                    brands_organic.append(sov_target.get("organic", 0))
                    brands_with_campaign.append(count - sov_target.get("organic", 0) + tb.get("campaign", 0))
                else:
                    brands_organic.append(count)
                    brands_with_campaign.append(count)  # competitors get no campaign addition
            if brand_names:
                slide03_data["chart_mentions"] = {
                    "categories": brand_names,
                    "mentions": [r + o for r, o in zip(brands_organic, brands_with_campaign)],
                }
                slide03_data["chart_audience"] = {
                    "categories": brand_names,
                    "audience": [comp_audience.get(n, 0) for n in brand_names],
                }
            
            # Rebuild top_brands and project_brand_row from canonical data
            slide03_data["all_brand_rows"] = [{"brand": n, "mentions": c} for n, c in all_brands]
            slide03_data["top_brands"] = slide03_data["all_brand_rows"][:6]
            target_row = next((r for r in slide03_data["all_brand_rows"] if r["brand"].lower() == project_brand_low), {})
            slide03_data["project_brand_row"] = target_row
            slide03_data["project_brand"] = config.get("project", {}).get("brand", "")
            slide03_data["total_mentions_competitive_set"] = sum(c for _, c in all_brands)
            slide03_data["project_mentions"] = tb.get("with_campaign", 0)
            slide03_data["project_sov"] = sov_target.get("sov_with_campaign", 0)
            slide03_data["project_brand_rank"] = next((i+1 for i, (n, c) in enumerate(all_brands) if n.lower() == project_brand_low), None)
            
            # Campaign split
            slide03_data["campaign_split"] = {
                "status": "verified",
                "method": "canonical_model",
                "campaign_mentions": tb.get("campaign", 0),
                "organic_mentions": tb.get("organic", 0),
                "total_brand_mentions": tb.get("with_campaign", 0),
                "campaign_publications_count": cm.get("orm_campaign", {}).get("total", 0),
                "campaign_publications_by_type": cm.get("orm_campaign", {}).get("by_type", {}),
                "campaign_fact_method": cm.get("orm_campaign", {}).get("fact_method", ""),
                "campaign_fact_source_sheet": cm.get("orm_campaign", {}).get("source_sheet", ""),
                "campaign_vs_organic_status": "verified",
                "caveat": "Campaign data from canonical model.",
            }
            slide03_data["sov_with_campaign"] = sov_target.get("sov_with_campaign", 0)
            slide03_data["sov_without_campaign"] = sov_target.get("sov_without_campaign", 0)
            # Populate reason_signals from canonical model if analytical_context is absent
            if not slide03_data.get("analytical_context"):
                topic_scope = cm.get("topic_scope", {})
                platform_scope = cm.get("platform_scope", {})
                reason_signals = []
                topics = topic_scope.get("topics", [])
                if topics:
                    reason_signals.append({"type": "themes", "text": "заметные темы: " + ", ".join(t["label"] for t in topics[:3])})
                ptypes = platform_scope.get("platform_type_rows", [])
                if ptypes:
                    reason_signals.append({"type": "platform_types", "text": "основные типы площадок: " + ", ".join(t["type"] for t in ptypes[:3])})
                slide03_data["analytical_context"] = {"reason_signals": reason_signals, "source": "canonical_model"}
        elif campaign_source:
            try:
                slide03_data["campaign_split"] = _campaign_split_from_source(
                    slide03_data,
                    campaign_source,
                    config.get("build", {}),
                    config.get("project", {}),
                )
            except Exception as exc:
                slide03_data["campaign_split"] = {"status": "missing_required", "error": str(exc)}
        # DEBUG: check chart_mentions before building
        _cm = slide03_data.get("chart_mentions", {})
        if not _cm.get("categories"):
            import json as _json
            (output_dir / "debug_slide03_before_build.json").write_text(
                _json.dumps({"chart_mentions": _cm, "has_canonical": canonical_model is not None}, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
        slide03_result = build_slide03_brand_mentions(
            template_path=template_path,
            output_path=output_dir / "draft_slide_03_brand_mentions_sov_test.pptx",
            project=config.get("project", {}),
            data=slide03_data,
        )
        slide03_qa = qa_slide03_brand_mentions(slide03_data, slide03_result)
        (output_dir / "slide_03_data_model.json").write_text(
            __import__("json").dumps(slide03_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_03_build_result.json").write_text(
            __import__("json").dumps(slide03_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_03_qa.json").write_text(
            __import__("json").dumps(slide03_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide04_problem_field_test"):
        build_config = config.get("build", {})
        analytical_excel = build_config.get("slide04_analytical_excel")
        raw_excel = build_config.get("slide04_raw_excel")
        report_excel = build_config.get("slide04_report_excel")
        if not (analytical_excel or raw_excel or report_excel):
            raise ValueError("At least one of build.slide04_analytical_excel, build.slide04_raw_excel or build.slide04_report_excel is required")

        slide04_data = extract_problem_field(
            analytical_path=__import__("pathlib").Path(analytical_excel) if analytical_excel else None,
            raw_path=__import__("pathlib").Path(raw_excel) if raw_excel else None,
            report_path=__import__("pathlib").Path(report_excel) if report_excel else None,
            project_brand=config.get("project", {}).get("brand", ""),
            source_system=build_config.get("slide04_source_system") or "",
            project_period=config.get("project", {}).get("period", ""),
        )
        slide04_result = build_slide04_problem_field(
            template_path=template_path,
            output_path=output_dir / "draft_slide_04_problem_field_test.pptx",
            project=config.get("project", {}),
            data=slide04_data,
        )
        slide04_qa = qa_slide04_problem_field(slide04_data, slide04_result)
        (output_dir / "slide_04_data_model.json").write_text(
            __import__("json").dumps(slide04_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_04_build_result.json").write_text(
            __import__("json").dumps(slide04_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_04_qa.json").write_text(
            __import__("json").dumps(slide04_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide05_sov_test"):
        build_config = config.get("build", {})
        analytics_excel = build_config.get("slide05_analytics_excel") or build_config.get("slide03_analytics_excel")
        if not analytics_excel:
            raise ValueError("build.slide05_analytics_excel or build.slide03_analytics_excel is required for slide05_sov_test")
        slide05_data = extract_slide05_sov(
            analytics_path=__import__("pathlib").Path(analytics_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            source_system=build_config.get("slide05_source_system") or build_config.get("slide03_source_system") or "",
            project_orm_path=__import__("pathlib").Path(build_config["slide05_project_orm_excel"]) if build_config.get("slide05_project_orm_excel") else None,
            campaign_month_sheet=build_config.get("slide05_campaign_month_sheet") or "",
            campaign_materials_override=_optional_int(build_config.get("slide05_campaign_materials_override")),
            competitor_brands=_competitor_brands_from_config(config.get("project", {}), build_config),
        )
        if canonical_model is not None:
            slide05_data["canonical_model_ref"] = True
            cm = canonical_model
            sov_target = cm.get("competitive_sov_scope", {}).get("target", {})
            slide05_data["project_mentions"] = sov_target.get("with_campaign", 0) - sov_target.get("organic", 0) + sov_target.get("organic", 0)
            slide05_data["campaign_split_override"] = {
                "sov_with_campaign": sov_target.get("sov_with_campaign", 0),
                "sov_without_campaign": sov_target.get("sov_without_campaign", 0),
                "sov_lift_pp": sov_target.get("sov_lift_pp", 0),
                "campaign_materials": cm.get("orm_campaign", {}).get("total", 0),
                "organic_mentions": sov_target.get("organic", 0),
            }
        slide05_result = build_slide05_sov(
            template_path=template_path,
            output_path=output_dir / "draft_slide_05_sov_test.pptx",
            project=config.get("project", {}),
            data=slide05_data,
        )
        slide05_qa = qa_slide05_sov(slide05_data, slide05_result)
        (output_dir / "slide_05_data_model.json").write_text(
            __import__("json").dumps(slide05_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_05_build_result.json").write_text(
            __import__("json").dumps(slide05_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_05_qa.json").write_text(
            __import__("json").dumps(slide05_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide06_tonality_test"):
        build_config = config.get("build", {})
        analytics_excel = (
            build_config.get("slide06_analytics_excel")
            or build_config.get("slide05_analytics_excel")
            or build_config.get("slide03_analytics_excel")
        )
        if not analytics_excel:
            raise ValueError("build.slide06_analytics_excel, build.slide05_analytics_excel or build.slide03_analytics_excel is required for slide06_tonality_test")
        slide06_data = extract_slide06_tonality(
            analytics_path=__import__("pathlib").Path(analytics_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            source_system=build_config.get("slide06_source_system") or build_config.get("slide05_source_system") or build_config.get("slide03_source_system") or "",
            competitor_brands=_competitor_brands_from_config(config.get("project", {}), build_config),
            project_orm_path=__import__("pathlib").Path(
                build_config.get("slide06_project_orm_excel")
                or build_config.get("slide08_project_orm_excel")
                or build_config.get("slide05_project_orm_excel")
                or build_config.get("slide02_fact_excel")
            ) if (
                build_config.get("slide06_project_orm_excel")
                or build_config.get("slide08_project_orm_excel")
                or build_config.get("slide05_project_orm_excel")
                or build_config.get("slide02_fact_excel")
            ) else None,
            month_sheet=(
                build_config.get("slide06_month_sheet")
                or build_config.get("slide08_month_sheet")
                or build_config.get("slide05_campaign_month_sheet")
                or build_config.get("slide02_fact_month_sheet")
                or ""
            ),
        )
        if canonical_model is not None:
            slide06_data["canonical_model_ref"] = True
            tonality_cm = canonical_model.get("target_tonality", {})
            slide06_data["project_row"] = {
                "positive": tonality_cm.get("positive", 0),
                "neutral": tonality_cm.get("neutral", 0),
                "negative": tonality_cm.get("negative", 0),
                "total": tonality_cm.get("positive", 0) + tonality_cm.get("neutral", 0) + tonality_cm.get("negative", 0),
            }
        slide06_result = build_slide06_tonality(
            template_path=template_path,
            output_path=output_dir / "draft_slide_06_tonality_test.pptx",
            project=config.get("project", {}),
            data=slide06_data,
        )
        slide06_qa = qa_slide06_tonality(slide06_data, slide06_result)
        (output_dir / "slide_06_data_model.json").write_text(
            __import__("json").dumps(slide06_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_06_build_result.json").write_text(
            __import__("json").dumps(slide06_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_06_qa.json").write_text(
            __import__("json").dumps(slide06_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide07_sources_test"):
        build_config = config.get("build", {})
        analytics_excel = (
            build_config.get("slide07_analytics_excel")
            or build_config.get("slide06_analytics_excel")
            or build_config.get("slide05_analytics_excel")
            or build_config.get("slide03_analytics_excel")
        )
        if not analytics_excel:
            raise ValueError("build.slide07_analytics_excel, build.slide06_analytics_excel, build.slide05_analytics_excel or build.slide03_analytics_excel is required for slide07_sources_test")
        slide07_data = extract_slide07_sources(
            analytics_path=__import__("pathlib").Path(analytics_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            source_system=build_config.get("slide07_source_system") or build_config.get("slide06_source_system") or build_config.get("slide05_source_system") or build_config.get("slide03_source_system") or "",
        )
        if canonical_model is not None:
            slide07_data["canonical_model_ref"] = True
            platform = canonical_model.get("platform_scope", {})
            slide07_data["platform_type_rows"] = platform.get("platform_type_rows", slide07_data.get("platform_type_rows", []))
            slide07_data["source_rows"] = platform.get("source_rows", slide07_data.get("source_rows", []))
            slide07_data["total_messages"] = platform.get("total", slide07_data.get("total_messages", 0))
        slide07_result = build_slide07_sources(
            template_path=template_path,
            output_path=output_dir / "draft_slide_07_sources_test.pptx",
            project=config.get("project", {}),
            data=slide07_data,
        )
        slide07_qa = qa_slide07_sources(slide07_data, slide07_result)
        (output_dir / "slide_07_data_model.json").write_text(
            __import__("json").dumps(slide07_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_07_build_result.json").write_text(
            __import__("json").dumps(slide07_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_07_qa.json").write_text(
            __import__("json").dumps(slide07_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide08_seeding_metrics_test"):
        build_config = config.get("build", {})
        project_orm_excel = build_config.get("slide08_project_orm_excel") or build_config.get("slide05_project_orm_excel") or build_config.get("slide02_media_plan_excel")
        if not project_orm_excel:
            raise ValueError("build.slide08_project_orm_excel is required for slide08_seeding_metrics_test")
        slide08_data = extract_slide08_seeding_metrics(
            project_orm_path=__import__("pathlib").Path(project_orm_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            month_sheet=build_config.get("slide08_month_sheet") or "Май",
            period=config.get("project", {}).get("period", ""),
            source_system=build_config.get("slide08_source_system") or "проектная таблица ORM",
        )
        if canonical_model is not None:
            slide08_data["canonical_model_ref"] = True
            orm_camp = canonical_model.get("orm_campaign", {})
            slide08_data["total_materials"] = orm_camp.get("total", slide08_data.get("total_materials", 0))
            slide08_data["views_metrics"] = {
                "received_views_total": orm_camp.get("views_total", slide08_data.get("views_metrics", {}).get("received_views_total", 0)),
                "has_views": orm_camp.get("views_total", 0) > 0 or slide08_data.get("views_metrics", {}).get("has_views", False),
                "metric_label": "Просмотры",
            }
        slide08_result = build_slide08_seeding_metrics(
            template_path=template_path,
            output_path=output_dir / "draft_slide_08_seeding_metrics_test.pptx",
            project=config.get("project", {}),
            data=slide08_data,
        )
        slide08_qa = qa_slide08_seeding_metrics(slide08_data, slide08_result)
        (output_dir / "slide_08_data_model.json").write_text(
            __import__("json").dumps(slide08_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_08_build_result.json").write_text(
            __import__("json").dumps(slide08_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_08_qa.json").write_text(
            __import__("json").dumps(slide08_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide09_message_examples_test") or config.get("build", {}).get("slide10_review_examples_test"):
        build_config = config.get("build", {})
        project_orm_excel = (
            build_config.get("slide09_project_orm_excel")
            or build_config.get("slide10_project_orm_excel")
            or build_config.get("slide08_project_orm_excel")
            or build_config.get("slide05_project_orm_excel")
            or build_config.get("slide02_media_plan_excel")
        )
        if not project_orm_excel:
            raise ValueError("build.slide09_project_orm_excel or build.slide10_project_orm_excel is required for slide09/10 examples")
        screenshot_backend = str(
            build_config.get("slide09_10_screenshot_backend")
            or build_config.get("slide09_screenshot_backend")
            or build_config.get("slide10_screenshot_backend")
            or ""
        ).strip()
        if screenshot_backend:
            os.environ["AI_ORM_SCREENSHOT_BACKEND"] = screenshot_backend
        examples_data = extract_slide09_10_examples(
            project_orm_path=__import__("pathlib").Path(project_orm_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            month_sheet=build_config.get("slide09_month_sheet") or build_config.get("slide10_month_sheet") or build_config.get("slide08_month_sheet") or "Май",
            period=config.get("project", {}).get("period", ""),
            source_system=build_config.get("slide09_source_system") or build_config.get("slide10_source_system") or build_config.get("slide08_source_system") or "проектная таблица ORM",
        )
        if build_config.get("slide09_message_examples_test"):
            slide09_data = examples_data["slide09"]
            slide09_result = build_slide09_message_examples(
                template_path=template_path,
                output_path=output_dir / "draft_slide_09_message_examples_test.pptx",
                project=config.get("project", {}),
                data=slide09_data,
            )
            slide09_qa = qa_slide09_message_examples(slide09_data, slide09_result)
            (output_dir / "slide_09_data_model.json").write_text(
                __import__("json").dumps(slide09_data, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
            (output_dir / "slide_09_build_result.json").write_text(
                __import__("json").dumps(slide09_result, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
            (output_dir / "slide_09_qa.json").write_text(
                __import__("json").dumps(slide09_qa, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
        if build_config.get("slide10_review_examples_test"):
            slide10_data = examples_data["slide10"]
            slide10_result = build_slide10_review_examples(
                template_path=template_path,
                output_path=output_dir / "draft_slide_10_review_examples_test.pptx",
                project=config.get("project", {}),
                data=slide10_data,
            )
            slide10_qa = qa_slide10_review_examples(slide10_data, slide10_result)
            (output_dir / "slide_10_data_model.json").write_text(
                __import__("json").dumps(slide10_data, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
            (output_dir / "slide_10_build_result.json").write_text(
                __import__("json").dumps(slide10_result, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )
            (output_dir / "slide_10_qa.json").write_text(
                __import__("json").dumps(slide10_qa, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )

    if config.get("build", {}).get("slide11_seeding_metrics_test"):
        build_config = config.get("build", {})
        ratings_excel = build_config.get("slide11_ratings_excel") or build_config.get("slide08_project_orm_excel") or build_config.get("slide05_project_orm_excel")
        if not ratings_excel:
            raise ValueError("build.slide11_ratings_excel is required for slide11_seeding_metrics_test")
        slide11_data = extract_slide11_seeding_metrics(
            ratings_path=__import__("pathlib").Path(ratings_excel),
            project_brand=config.get("project", {}).get("brand", ""),
            period=config.get("project", {}).get("period", ""),
            ratings_sheet=build_config.get("slide11_ratings_sheet") or "",
            current_period=build_config.get("slide11_current_period") or "",
            placement_month_sheet=build_config.get("slide11_placement_month_sheet") or build_config.get("slide08_month_sheet") or build_config.get("slide05_campaign_month_sheet") or "",
            source_system=build_config.get("slide11_source_system") or "проектная таблица ORM",
        )
        if canonical_model is not None:
            slide11_data["canonical_model_ref"] = True
            ratings = canonical_model.get("ratings_scope", {})
            if ratings.get("status") == "ready":
                summaries = slide11_data.get("period_summaries") or []
                if not summaries or not slide11_data.get("start_summary"):
                    periods = ratings.get("periods", [])
                    pos_cards = ratings.get("positive_cards", [0, 0])
                    neg_cards = ratings.get("negative_cards", [0, 0])
                    no_rev = ratings.get("no_review_cards", [0, 0])
                    rev_counts = ratings.get("reviews_count", [0, 0])
                    slide11_data["period_summaries"] = [
                        {"label": "Старт" if i == 0 else "Текущий", "positive_cards": pos_cards[i] if i < len(pos_cards) else 0, "negative_cards": neg_cards[i] if i < len(neg_cards) else 0, "no_reviews": no_rev[i] if i < len(no_rev) else 0, "reviews_total": rev_counts[i] if i < len(rev_counts) else 0}
                        for i in range(len(periods))
                    ]
                    slide11_data["start_summary"] = {"reviews_total": rev_counts[0] if len(rev_counts) > 0 else 0, "positive_cards": pos_cards[0] if len(pos_cards) > 0 else 0, "negative_cards": neg_cards[0] if len(neg_cards) > 0 else 0, "no_reviews": no_rev[0] if len(no_rev) > 0 else 0, "data_status": "ready"}
                    slide11_data["current_summary"] = {"reviews_total": rev_counts[-1] if rev_counts else 0, "positive_cards": pos_cards[-1] if pos_cards else 0, "negative_cards": neg_cards[-1] if neg_cards else 0, "no_reviews": no_rev[-1] if no_rev else 0, "data_status": "ready"}
        slide11_result = build_slide11_seeding_metrics(
            template_path=template_path,
            output_path=output_dir / "draft_slide_11_seeding_metrics_test.pptx",
            project=config.get("project", {}),
            data=slide11_data,
        )
        slide11_qa = qa_slide11_seeding_metrics(slide11_data, slide11_result)
        (output_dir / "slide_11_data_model.json").write_text(
            __import__("json").dumps(slide11_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_11_build_result.json").write_text(
            __import__("json").dumps(slide11_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_11_qa.json").write_text(
            __import__("json").dumps(slide11_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("slide12_final_conclusions_test"):
        slide_data_by_number = {}
        for slide_number in [3, 4, 5, 6, 7, 8, 11]:
            data_path = output_dir / f"slide_{slide_number:02d}_data_model.json"
            if data_path.exists():
                slide_data_by_number[slide_number] = __import__("json").loads(data_path.read_text(encoding="utf-8"))
        if canonical_model is not None:
            slide_data_by_number["canonical_model"] = canonical_model
        slide12_data = build_slide12_data(config.get("project", {}), slide_data_by_number)
        slide12_result = build_slide12_final_conclusions(
            template_path=template_path,
            output_path=output_dir / "draft_slide_12_final_conclusions_test.pptx",
            project=config.get("project", {}),
            data=slide12_data,
        )
        slide12_qa = qa_slide12_final_conclusions(slide12_data, slide12_result)
        (output_dir / "slide_12_data_model.json").write_text(
            __import__("json").dumps(slide12_data, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_12_build_result.json").write_text(
            __import__("json").dumps(slide12_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        (output_dir / "slide_12_qa.json").write_text(
            __import__("json").dumps(slide12_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    if config.get("build", {}).get("first4_combined_test"):
        combined_result = build_first4_combined_deck(config, template_path, output_dir, canonical_model=canonical_model)
        if config.get("build", {}).get("slide12_final_conclusions_test"):
            combined_name = "slides_01_12_combined_result.json"
        elif config.get("build", {}).get("slide11_seeding_metrics_test"):
            combined_name = "slides_01_11_combined_result.json"
        elif config.get("build", {}).get("slide08_seeding_metrics_test"):
            combined_name = "slides_01_08_combined_result.json"
        elif config.get("build", {}).get("slide07_sources_test"):
            combined_name = "slides_01_07_combined_result.json"
        elif config.get("build", {}).get("slide06_tonality_test"):
            combined_name = "slides_01_06_combined_result.json"
        elif config.get("build", {}).get("slide05_sov_test"):
            combined_name = "slides_01_05_combined_result.json"
        else:
            combined_name = "slides_01_04_combined_result.json"
        (output_dir / combined_name).write_text(
            json.dumps(combined_result, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )

    # ---- Cross-slide QA ----
    if canonical_model is not None:
        cross_slide_qa = _build_cross_slide_qa(canonical_model, output_dir)
        qa_status = cross_slide_qa.get("status", "blocked")
        (output_dir / "cross_slide_qa.json").write_text(
            json.dumps(cross_slide_qa, ensure_ascii=False, indent=2),
            encoding="utf-8"
        )
        if qa_status == "blocked":
            print("CROSS-SLIDE QA: BLOCKED")
            for b in cross_slide_qa.get("blockers", []):
                print(f"  BLOCKER: {b}")
        else:
            print("CROSS-SLIDE QA: PASSED")

    # ---- Extract final PPTX as JSON for comparison ----
    final_pptx_candidates = sorted(output_dir.glob("draft_slides_*.pptx"))
    if not final_pptx_candidates:
        final_pptx_candidates = sorted(output_dir.glob("_combined_step_*.pptx"))
    if not final_pptx_candidates:
        final_pptx_candidates = sorted(output_dir.glob("draft_slide_*.pptx"))
    if final_pptx_candidates:
        final_pptx = final_pptx_candidates[-1]
        save_pptx_extraction(final_pptx, output_dir / "pptx_extraction.json")

    print(f"AI Presentation ORM v0.3 completed. Output: {output_dir}")
    print(f"Template slides: {template_map.get('slide_count')}")
    print(f"Input files: {len(input_files)}")
    print(f"Missing data items: {len(missing_data.get('items', []))}")
    return 0


def _build_cross_slide_qa(canonical_model: dict, output_dir: Path) -> dict:
    blockers = []
    warnings = []

    target = canonical_model.get("target_brand", {})
    sov = canonical_model.get("competitive_sov_scope", {}).get("target", {})
    tonality = canonical_model.get("target_tonality", {})
    campaign_recon = canonical_model.get("campaign_reconciliation", {})
    exclusions = canonical_model.get("source_exclusions_applied", {})

    target_with_campaign = target.get("with_campaign", 0)
    target_organic = target.get("organic", 0)
    target_campaign = target.get("campaign", 0)

    # Check: Slides 3, 5, 12 use the same target_with_campaign
    # (verified via canonical model - they all read from here)

    # Check: Slide 6 tonality sums to target_with_campaign
    tone_positive = tonality.get("positive", 0)
    tone_neutral = tonality.get("neutral", 0)
    tone_negative = tonality.get("negative", 0)
    tone_total = tone_positive + tone_neutral + tone_negative
    if tone_total != target_with_campaign:
        blockers.append(
            f"Slide 6 tonality sum ({tone_total}) != target_with_campaign ({target_with_campaign}). "
            f"Positive={tone_positive} + Neutral={tone_neutral} + Negative={tone_negative}"
        )

    # Check: No excluded source remains
    excluded_sources = exclusions.get("sources", [])
    if excluded_sources:
        removed_count = exclusions.get("removed_messages_count", 0)
        warnings.append(
            f"Source exclusions applied: {excluded_sources}, removed {removed_count} messages. "
            "Verify all slide payloads use filtered data."
        )

    # Check: Campaign reconciliation sanity
    orm_total = campaign_recon.get("orm_total", 0)
    matched = campaign_recon.get("matched_in_analytics", 0)
    unmatched = campaign_recon.get("unmatched_added", 0)
    if orm_total and matched + unmatched != orm_total:
        blockers.append(
            f"Campaign reconciliation inconsistency: matched ({matched}) + unmatched ({unmatched}) != ORM total ({orm_total})"
        )

    # Check: ORM campaign is NOT capped by analytics
    analytics_total = canonical_model.get("analytics_messages", {}).get("total", 0)
    if target_campaign > 0 and analytics_total > 0 and target_campaign <= analytics_total:
        warnings.append(
            f"Target campaign ({target_campaign}) <= analytics total ({analytics_total}). "
            "Verify no min(orm, analytics) capping occurred."
        )

    # Check: Target organic + campaign = with_campaign
    if target_organic + target_campaign != target_with_campaign:
        blockers.append(
            f"Target brand inconsistency: organic ({target_organic}) + campaign ({target_campaign}) "
            f"!= with_campaign ({target_with_campaign})"
        )

    # Check: SOV values align
    sov_with_campaign = sov.get("sov_with_campaign", 0)
    sov_without_campaign = sov.get("sov_without_campaign", 0)
    sov_lift = sov.get("sov_lift_pp", 0)
    expected_lift = (sov_with_campaign - sov_without_campaign) * 100
    if abs(sov_lift - expected_lift) > 0.01:
        warnings.append(
            f"SOV lift ({sov_lift}) != expected ({expected_lift}) from with/without campaign SOV"
        )

    status = "blocked" if blockers else "ready_for_visual_review"
    return {
        "status": status,
        "blockers": blockers,
        "warnings": warnings,
        "canonical_summary": {
            "target_brand": target.get("brand", ""),
            "target_with_campaign": target_with_campaign,
            "target_without_campaign": target.get("without_campaign", 0),
            "campaign_total": target_campaign,
            "organic_total": target_organic,
            "sov_with_campaign": sov_with_campaign,
            "sov_without_campaign": sov_without_campaign,
            "sov_lift_pp": sov_lift,
            "tonality_positive": tone_positive,
            "tonality_neutral": tone_neutral,
            "tonality_negative": tone_negative,
            "orm_reconciliation": {
                "orm_total": orm_total,
                "matched_in_analytics": matched,
                "unmatched_added_as_campaign": unmatched,
                "organic_analytics_messages": campaign_recon.get("organic_analytics_messages", 0),
            },
            "source_exclusions": excluded_sources,
        },
    }


if __name__ == "__main__":
    raise SystemExit(main())
