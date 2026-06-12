from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

BRAND_BLUE = RGBColor(0, 174, 239)
BLACK = RGBColor(0, 0, 0)

MONTHS = {
    "01": ("январь", "январе", "января", "Январь"),
    "02": ("февраль", "феврале", "февраля", "Февраль"),
    "03": ("март", "марте", "марта", "Март"),
    "04": ("апрель", "апреле", "апреля", "Апрель"),
    "05": ("май", "мае", "мая", "Май"),
    "06": ("июнь", "июне", "июня", "Июнь"),
    "07": ("июль", "июле", "июля", "Июль"),
    "08": ("август", "августе", "августа", "Август"),
    "09": ("сентябрь", "сентябре", "сентября", "Сентябрь"),
    "10": ("октябрь", "октябре", "октября", "Октябрь"),
    "11": ("ноябрь", "ноябре", "ноября", "Ноябрь"),
    "12": ("декабрь", "декабре", "декабря", "Декабрь"),
}


OLD_TEXT_PATTERNS = [
    "БАКТОБЛИС",
    "Baktoblis",
    "Бактоблис",
    "ПИКТОРИД",
    "Пикторид",
    "Pictorid",
    "Декабрь 2025",
]


def _brand_without_registered(brand: str) -> str:
    return re.sub(r"\s*[®Â®]\s*", "", str(brand or "")).strip()


def _set_registered_superscript(run) -> None:
    r_pr = run._r.get_or_add_rPr()
    r_pr.set("baseline", "30000")


def _replace_in_text_frame(shape, replacements: dict[str, str]) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0
    count = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            new_text = text
            for old, new in replacements.items():
                if old in new_text:
                    new_text = new_text.replace(old, new)
            if new_text != text:
                run.text = new_text
                count += 1
    return count


def _period_month_label(period: str) -> str:
    text = str(period or "").strip()
    low = text.lower()
    date_match = re.search(r"\d{1,2}[.](\d{1,2})[.](\d{2,4})", low)
    if date_match:
        month_num, year_raw = date_match.groups()
        month = MONTHS.get(month_num.zfill(2), ("", "", "", ""))[3]
        year = year_raw if len(year_raw) == 4 else "20" + year_raw
        return f"{month} {year}".strip()
    for _, forms in MONTHS.items():
        nominative, prepositional, genitive, title = forms
        if nominative in low or prepositional in low or genitive in low:
            year_match = re.search(r"20\d{2}", low)
            return f"{title} {year_match.group(0) if year_match else ''}".strip()
    return text


def _set_cover_title(shape, brand: str) -> None:
    tf = shape.text_frame
    clean_brand = _brand_without_registered(brand)
    first_size = Pt(36)
    first_name = "Arial"
    reg_size = first_size
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        first_size = run.font.size or first_size
        first_name = run.font.name or first_name
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if "®" in run.text or "Â®" in run.text:
                reg_size = run.font.size or reg_size
    reg_size = first_size
    tf.clear()
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "ИНФОРМАЦИОННАЯ КАМПАНИЯ"
    run.font.size = first_size
    run.font.bold = True
    run.font.name = first_name
    run.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = clean_brand.upper()
    run.font.size = first_size
    run.font.bold = True
    run.font.name = first_name
    run.font.color.rgb = BRAND_BLUE
    reg = p.add_run()
    reg.text = "®"
    reg.font.size = reg_size
    reg.font.bold = True
    reg.font.name = first_name
    reg.font.color.rgb = BRAND_BLUE
    _set_registered_superscript(reg)


def _set_period_text(shape, period: str) -> None:
    shape.text = _period_month_label(period)
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Arial"
                run.font.size = Pt(18)
                run.font.bold = False
                run.font.color.rgb = BLACK


def build_slide01_cover(template_path: Path, output_path: Path, project: dict) -> dict:
    """Create a template copy with slide 1 text replaced.

    This function intentionally edits existing text runs to preserve fonts, sizes,
    positions, colors and all design elements.
    """
    prs = Presentation(str(template_path))
    if not prs.slides:
        raise ValueError("Template has no slides")

    brand = project.get("brand") or ""
    clean_brand = _brand_without_registered(brand)
    period = project.get("period") or ""
    if not brand:
        raise ValueError("project.brand is required for cover slide")
    if not period:
        raise ValueError("project.period is required for cover slide")

    slide = prs.slides[0]

    replacements = {
        "БАКТОБЛИС": clean_brand.upper(),
        "Бактоблис": clean_brand,
        "Baktoblis": clean_brand,
        "ПИКТОРИД®": clean_brand.upper(),
        "ПИКТОРИД": clean_brand.upper(),
        "Пикторид®": clean_brand,
        "Пикторид": clean_brand,
        "Pictorid": clean_brand,
        "Информационная кампания": "ИНФОРМАЦИОННАЯ КАМПАНИЯ",
        "ИНФОРМАЦИОННАЯ КАМПАНИЯ": "ИНФОРМАЦИОННАЯ КАМПАНИЯ",
        "Декабрь 2025": period,
        "декабрь 2025": period,
        "ПРЕДВАРИТЛЬНО": period,
        "ПРЕДВАРИТЕЛЬНО": period,
    }

    replacement_count = 0
    text_after = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            up = (shape.text or "").upper()
            has_old_brand = any(item.upper() in up for item in ["ПИКТОРИД", "БАКТОБЛИС", "PICTORID", "BAKTOBLIS"])
            if "ИНФОРМАЦИОННАЯ КАМПАНИЯ" in up and (has_old_brand or "®" in up or "\n" in shape.text):
                _set_cover_title(shape, brand)
                replacement_count += 1
            elif "ДЕКАБРЬ 2025" in up or "ПРЕДВАР" in up or re.search(r"\d{1,2}[–-]\d{1,2}\s+[А-ЯA-Z]", up):
                _set_period_text(shape, period)
                replacement_count += 1
            else:
                replacement_count += _replace_in_text_frame(shape, replacements)
        if getattr(shape, "has_text_frame", False):
            txt = shape.text.strip()
            if txt:
                text_after.append(txt)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    old_text_left = []
    joined = "\n".join(text_after)
    for old in OLD_TEXT_PATTERNS:
        if old.lower() == brand.lower():
            continue
        if old.lower() in joined.lower():
            old_text_left.append(old)

    return {
        "output": str(output_path),
        "replacement_count": replacement_count,
        "old_text_left": old_text_left,
        "text_after": text_after,
    }
