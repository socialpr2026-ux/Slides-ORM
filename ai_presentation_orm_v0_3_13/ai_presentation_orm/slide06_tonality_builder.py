from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from .chart_style import CHART_FONT, CHART_LABEL_SIZE
from .native_chart_helpers import chart_update, replace_chart_data
from .slide04_problem_field_builder import EMU_PER_INCH, _format_source_period
from .style_rules import BODY_FONT_SIZE, FONT_FAMILY, SOURCE_FONT_SIZE, TITLE_FONT_SIZE, format_int_spaces, qa_slide_copy_principles
from .text_contracts import slide06_tonality_insight, qa_text_contract

TEXT_COLOR = RGBColor(0, 0, 0)
BLUE = RGBColor(1, 162, 255)
POSITIVE = "#00B050"
NEUTRAL = "#D9D9D9"
NEGATIVE = "#C00000"
CHART_TITLE_COLOR = "#666666"


def _fmt_pct(value: float, digits: int = 1) -> str:
    return f"{float(value) * 100:.{digits}f}%".replace(".", ",")


def _fmt_int(value: int | float) -> str:
    return format_int_spaces(value)


def _message_word(value: int | float) -> str:
    number = abs(int(round(float(value or 0))))
    if 11 <= number % 100 <= 14:
        return "сообщений"
    last = number % 10
    if last == 1:
        return "сообщение"
    if 2 <= last <= 4:
        return "сообщения"
    return "сообщений"


def _period_in_text(data: dict) -> str:
    month = str(data.get("month_prepositional") or "отчетном периоде").strip()
    if month.lower().startswith(("в ", "во ")):
        return month
    return f"в {month}"


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _set_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT):
    shape.text = text
    tf = shape.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = TEXT_COLOR


def _add_run(paragraph, text: str, *, size: float = BODY_FONT_SIZE, bold: bool | None = None, color: RGBColor | None = None, hyperlink: str = ""):
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or TEXT_COLOR
    if hyperlink:
        run.hyperlink.address = hyperlink
        run.font.underline = False
    return run


def _set_native_blue_bullet(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    bullet_suffixes = ("}buNone", "}buClr", "}buFont", "}buChar", "}buAutoNum")
    for child in list(p_pr):
        if child.tag.endswith(bullet_suffixes):
            p_pr.remove(child)
    p_pr.set("marL", "171450")
    p_pr.set("indent", "-171450")
    bu_clr = OxmlElement("a:buClr")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", "01A2FF")
    bu_clr.append(srgb)
    bu_font = OxmlElement("a:buFont")
    bu_font.set("typeface", FONT_FAMILY)
    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", "•")
    p_pr.insert(0, bu_char)
    p_pr.insert(0, bu_font)
    p_pr.insert(0, bu_clr)


def _append_highlighted(paragraph, text: str, *, highlight_words: set[str] | None = None, highlight_color: RGBColor | None = None):
    words = highlight_words or set()
    if not words:
        _add_run(paragraph, text)
        return
    cursor = 0
    lowered = text.lower()
    while cursor < len(text):
        matches = [(lowered.find(word.lower(), cursor), word) for word in words]
        matches = [(idx, word) for idx, word in matches if idx >= 0]
        if not matches:
            _add_run(paragraph, text[cursor:])
            break
        idx, word = min(matches, key=lambda item: item[0])
        if idx > cursor:
            _add_run(paragraph, text[cursor:idx])
        _add_run(paragraph, text[idx:idx + len(word)], color=highlight_color or BLUE)
        cursor = idx + len(word)


def _set_conclusions(shape, data: dict):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    header = tf.paragraphs[0]
    header.alignment = PP_ALIGN.LEFT
    header.space_after = Pt(10)
    header.line_spacing = 1.0
    _add_run(header, "ВЫВОДЫ ПО АНАЛИЗУ:", size=BODY_FONT_SIZE, bold=True, color=BLUE)

    row = data.get("project_row") or {}
    brand = data.get("project_brand") or "бренд"
    period_in = _period_in_text(data)
    total = int(row.get("total", 0) or 0)
    positive = int(row.get("positive", 0) or 0)
    neutral = int(row.get("neutral", 0) or 0)
    negative = int(row.get("negative", 0) or 0)
    positive_share = positive / total if total else 0
    neutral_share = neutral / total if total else 0
    negative_share = negative / total if total else 0
    positive_pct = _fmt_pct(positive_share, 0)
    negative_pct = _fmt_pct(negative_share, 1 if 0 < negative_share < 0.01 else 0)
    neutral_pct = _fmt_pct(neutral_share, 0)
    signals = data.get("signals") or {}
    samples = signals.get("samples") or {}

    def sample_link(tone: str, label: str) -> list[tuple[str, dict]]:
        for sample in samples.get(tone) or []:
            if isinstance(sample, dict) and sample.get("link"):
                site = sample.get("site") or sample.get("source_type") or "источник"
                return [(f" Пример: {label} на {site}.", {"link": sample.get("link"), "color": BLUE})]
        return []

    bullets = [
        [
            (f"Доля позитивных упоминаний {period_in} у бренда ", {}),
            (brand, {"bold": True}),
            (" – ", {}),
            (positive_pct, {"bold": True, "color": RGBColor(0, 176, 80)}),
            (
                f" ({_fmt_int(positive)} {_message_word(positive)}). "
                "Позитивные сообщения отражают опыт применения, удобство продукта и общий положительный эффект.",
                {},
            ),
        ],
        [
            ("Доля негативных упоминаний бренда ", {}),
            (brand, {"bold": True}),
            (" – ", {}),
            (negative_pct, {"bold": True, "color": RGBColor(255, 0, 0)}),
            (
                f" ({_fmt_int(negative)} {_message_word(negative)}). "
                "Негативные сообщения нужно разбирать отдельно: они показывают зоны контроля в карточках, ответах и пользовательском опыте.",
                {},
            ),
        ] if negative else [
            ("Негативные упоминания бренда ", {}),
            (brand, {"bold": True}),
            (" не зафиксированы. Важно сохранять мониторинг вопросов и нейтральных сомнений.", {}),
        ],
        [
            (
                f"Нейтральные упоминания составили {neutral_pct} ({_fmt_int(neutral)} {_message_word(neutral)}): "
                "это вопросы, уточнения и обсуждения выбора без выраженной оценки.",
                {},
            ),
        ],
    ]
    for bullet in bullets:
        paragraph = tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = 1.0
        _set_native_blue_bullet(paragraph)
        for text, style in bullet:
            _add_run(
                paragraph,
                text,
                size=BODY_FONT_SIZE,
                bold=style.get("bold"),
                color=style.get("color"),
                hyperlink=style.get("link", ""),
            )


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)
    parent.append(shape._element)


def _source_note(data: dict) -> str:
    source_label = str(data.get("source_label", "сервис мониторинга Brand Analytics")).replace("сервис мониторинга ", "")
    return (
        f"Источник: {source_label}, {_format_source_period(data.get('period_raw'))}. "
        "Исключен спам и нерелевантные упоминания."
    )


def _short_label(value: str, limit: int = 16) -> str:
    text = str(value or "")
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "..."


def _conclusion_text(data: dict) -> str:
    brand = data.get("project_brand") or "бренд"
    month_prep = data.get("month_prepositional") or "периоде"
    row = data.get("project_row") or {}
    total = int(row.get("total", 0) or 0)
    positive = int(row.get("positive", 0) or 0)
    neutral = int(row.get("neutral", 0) or 0)
    negative = int(row.get("negative", 0) or 0)
    return (
        "ВЫВОДЫ ПО АНАЛИЗУ:\n\n"
        f"• Доля позитивных упоминаний в {month_prep} у бренда {brand} – {round((positive / total * 100) if total else 0)}% "
        f"({_fmt_int(positive)} сообщения). В отзывах и комментариях пользователи отмечают удобство применения, качественную упаковку, быструю доставку и поддержку в сезон простуд.\n"
        f"• Доля негативных упоминаний бренда {brand} – {round((negative / total * 100) if total else 0)}% "
        f"({_fmt_int(negative)} сообщений). Негатив связан с доставкой и упаковкой, а также с отдельными сообщениями об отсутствии эффекта в конкретных ситуациях.\n"
        f"• Нейтральные упоминания составили {round((neutral / total * 100) if total else 0)}% "
        f"({_fmt_int(neutral)} сообщений). Они представлены вопросами о назначениях, микрофлоре и опыте применения без яркой эмоциональной окраски."
    )


def _render_brand_stack(data: dict, out_path: Path, *, width_in: float, height_in: float) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    rows = (data.get("brand_rows") or [])[:6]
    project = (data.get("project_brand") or "").lower()
    labels = [_short_label(row.get("brand", ""), 16) for row in rows]
    positive = [int(row.get("positive", 0) or 0) for row in rows]
    neutral = [int(row.get("neutral", 0) or 0) for row in rows]
    negative = [int(row.get("negative", 0) or 0) for row in rows]
    totals = [p + n + g for p, n, g in zip(positive, neutral, negative)]

    x = np.arange(len(rows))
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    bars_pos = ax.bar(x, positive, color=POSITIVE, width=0.48)
    bars_neu = ax.bar(x, neutral, bottom=positive, color=NEUTRAL, width=0.48)
    bottoms_neg = [p + n for p, n in zip(positive, neutral)]
    bars_neg = ax.bar(x, negative, bottom=bottoms_neg, color=NEGATIVE, width=0.48)
    ax.set_title(
        "Распределение тональности упоминаний брендов",
        fontsize=12.5,
        fontweight="bold",
        fontname=CHART_FONT,
        color=CHART_TITLE_COLOR,
        pad=6,
    )

    for bars, values, bottoms, text_color, min_inside in [
        (bars_pos, positive, [0] * len(rows), "white", 2),
        (bars_neu, neutral, positive, "#333333", 2),
    ]:
        for bar, value, bottom in zip(bars, values, bottoms):
            if value < min_inside:
                continue
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                bottom + value / 2,
                f"{value:.0f}",
                ha="center",
                va="center",
                fontsize=CHART_LABEL_SIZE,
                fontname=CHART_FONT,
                color=text_color,
            )
    for bar, value, total in zip(bars_neg, negative, totals):
        top = total
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            top + max(max(totals or [1]) * 0.018, 1.0),
            f"{value}",
            ha="center",
            va="bottom",
            fontsize=CHART_LABEL_SIZE,
            fontname=CHART_FONT,
            color="#000000",
        )

    ax.set_ylim(0, max(totals or [1]) * 1.16)
    ax.set_xticks(x, labels=labels)
    ax.tick_params(axis="x", length=0, labelsize=CHART_LABEL_SIZE, colors="#333333", pad=3)
    ax.tick_params(axis="y", left=False, labelleft=False)
    ax.set_yticks([])
    for label, row in zip(ax.get_xticklabels(), rows):
        label.set_fontname(CHART_FONT)
        if str(row.get("brand", "")).lower() == project:
            label.set_fontweight("bold")
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.grid(False)
    plt.subplots_adjust(left=0.035, right=0.985, top=0.86, bottom=0.15)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _render_daily_lines(data: dict, out_path: Path, *, width_in: float, height_in: float) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    rows = data.get("daily_rows") or []
    x = np.arange(len(rows))
    labels = []
    for row in rows:
        date = str(row.get("date", ""))
        parts = date.split(".")
        labels.append(f"{int(parts[0])} мая" if len(parts) >= 2 and parts[0].isdigit() else date)
    series = [
        ("Позитив", [row.get("positive", 0) for row in rows], POSITIVE),
        ("Нейтрал", [row.get("neutral", 0) for row in rows], NEUTRAL),
        ("Негатив", [row.get("negative", 0) for row in rows], NEGATIVE),
    ]
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_title(
        "Динамика публикаций сообщений бренда",
        fontsize=12.5,
        fontweight="bold",
        fontname="DejaVu Sans",
        color=CHART_TITLE_COLOR,
        pad=6,
    )
    for name, values, color in series:
        ax.plot(x, values, color=color, linewidth=2.1, label=name)
    max_v = max([max(values or [0]) for _, values, _ in series] + [1])
    ax.set_ylim(0, max(6, max_v + 1))
    tick_idx = list(range(len(rows)))
    ax.set_xticks(tick_idx, labels=[labels[idx] for idx in tick_idx])
    ax.tick_params(axis="x", length=0, labelsize=5.8, colors="#333333", pad=1)
    ax.tick_params(axis="y", length=0, labelsize=8.4, colors="#333333")
    for label in ax.get_xticklabels():
        label.set_fontname("DejaVu Sans")
        label.set_rotation(90)
        label.set_ha("center")
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.grid(False)
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.42), ncol=3, frameon=False, prop={"family": "DejaVu Sans", "size": 8.8})
    plt.subplots_adjust(left=0.07, right=0.99, top=0.82, bottom=0.40)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def build_slide06_tonality(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    slide = prs.slides[5]
    image_dir = output_path.parent / "slide06_chart_images"

    left_chart = _shape_by_id(slide, 13)
    right_chart = _shape_by_id(slide, 3)
    title = _shape_by_id(slide, 14)
    conclusions = _shape_by_id(slide, 5)
    source = _shape_by_id(slide, 2)

    if title:
        _set_text(title, "ТОНАЛЬНОСТЬ СООБЩЕНИЙ О БРЕНДАХ КАТЕГОРИИ", font_size=TITLE_FONT_SIZE, bold=False)
    if conclusions:
        _set_conclusions(conclusions, data)
    if source:
        _set_text(source, _source_note(data), font_size=SOURCE_FONT_SIZE)

    chart_updates = []
    image_outputs = []
    chart_shapes = [shape for shape in [left_chart, right_chart] if shape]
    common_top = min((shape.top for shape in chart_shapes), default=None)
    common_height = min((shape.height for shape in chart_shapes), default=None)
    if left_chart and right_chart:
        total_space = (right_chart.left + right_chart.width) - left_chart.left
        left_width = int(total_space * 0.38)
        right_width = total_space - left_width
        left_pos = (left_chart.left, common_top or left_chart.top, left_width, common_height or left_chart.height)
        right_pos = (left_chart.left + left_width, common_top or right_chart.top, right_width, common_height or right_chart.height)
    elif left_chart:
        left_pos = (left_chart.left, common_top or left_chart.top, left_chart.width, common_height or left_chart.height)
    elif right_chart:
        right_pos = (right_chart.left, common_top or right_chart.top, right_chart.width, common_height or right_chart.height)
    if left_chart:
        left_chart.left, left_chart.top, left_chart.width, left_chart.height = [int(value) for value in left_pos]
        rows = (data.get("brand_rows") or [])[:6]
        labels = [_short_label(row.get("brand", ""), 16) for row in rows] or ["Нет данных"]
        replace_chart_data(
            left_chart,
            categories=labels,
            series=[
                ("Позитив", [int(row.get("positive", 0) or 0) for row in rows] or [0], POSITIVE),
                ("Нейтрал", [int(row.get("neutral", 0) or 0) for row in rows] or [0], NEUTRAL),
                ("Негатив", [int(row.get("negative", 0) or 0) for row in rows] or [0], NEGATIVE),
            ],
            show_legend=True,
            label_position=XL_LABEL_POSITION.CENTER,
            bar_grouping="percentStacked",
        )
        chart_updates.append(chart_update("brand_tonality_stack", "template_stacked_columns_replace_data", categories=labels))
    if right_chart:
        right_chart.left, right_chart.top, right_chart.width, right_chart.height = [int(value) for value in right_pos]
        rows = data.get("daily_rows") or []
        labels = []
        for row in rows:
            date = str(row.get("date", ""))
            parts = date.split(".")
            labels.append(f"{int(parts[0])} мая" if len(parts) >= 2 and parts[0].isdigit() else date)
        replace_chart_data(
            right_chart,
            categories=labels or ["Период"],
            series=[
                ("Позитив", [int(row.get("positive", 0) or 0) for row in rows] or [0], POSITIVE),
                ("Нейтрал", [int(row.get("neutral", 0) or 0) for row in rows] or [0], NEUTRAL),
                ("Негатив", [int(row.get("negative", 0) or 0) for row in rows] or [0], NEGATIVE),
            ],
            show_legend=True,
            show_data_labels=False,
            hide_value_axis=False,
        )
        chart_updates.append(chart_update("project_daily_tonality", "template_line_replace_data", categories=labels))

    for shape in [title, conclusions, source]:
        if shape:
            _bring_to_front(shape)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    text_after = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
    joined = "\n".join(text_after)
    forbidden = ["Пикторид", "декабрь 2025", "В декабре", "48% (158"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]
    return {
        "output": str(output_path),
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "old_text_left": old_text_left,
        "conclusion_text": "\n".join(text_after),
        "source_note": source.text if source else "",
        "text_after": text_after,
        "visual_guardrails": {
            "left_chart_geometry": "uses inherited reference chart slot",
            "right_chart_geometry": "uses inherited reference chart slot",
            "chart_alignment": {
                "common_top": int(common_top) if common_top is not None else None,
                "common_height": int(common_height) if common_height is not None else None,
            },
            "tone_colors": {"positive": POSITIVE, "neutral": NEUTRAL, "negative": NEGATIVE},
        },
    }


def qa_slide06_tonality(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    source_note = result.get("source_note", "")
    joined = "\n".join(result.get("text_after") or [])
    totals = data.get("totals_check") or {}
    if not data.get("project_row"):
        blockers.append("Project brand tonality row is missing.")
    if not data.get("daily_rows"):
        blockers.append("Daily tonality rows are missing.")
    if not totals.get("matches"):
        blockers.append(
            "Project tonality totals mismatch: "
            f"tags={totals.get('project_total_from_tags')} daily={totals.get('project_total_from_messages_daily')}"
        )
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 06: " + ", ".join(result["old_text_left"]))
    if ";" in joined:
        blockers.append("Slide 06 text must not use semicolons.")
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    if any(time in source_note for time in ["00:00", "23:59"]):
        blockers.append("Source note must not include hours/minutes.")
    if (data.get("daily_method") or {}).get("method") != "messages_sheet_project_brand_filter":
        warnings.append("Daily tonality uses fallback, not project-brand message filter.")
    project_row = data.get("project_row") or {}
    brand_rows = (data.get("brand_rows") or [])
    project_brand_lower = (data.get("project_brand") or "").lower()
    chart_project_row = {}
    for row in brand_rows:
        if str(row.get("brand", "")).strip().lower() == project_brand_lower:
            chart_project_row = row
            break
    text_pos = int(project_row.get("positive", 0) or 0)
    text_neu = int(project_row.get("neutral", 0) or 0)
    text_neg = int(project_row.get("negative", 0) or 0)
    if chart_project_row:
        chart_pos = int(chart_project_row.get("positive", 0) or 0)
        chart_neu = int(chart_project_row.get("neutral", 0) or 0)
        chart_neg = int(chart_project_row.get("negative", 0) or 0)
        if abs(text_pos - chart_pos) > 1 or abs(text_neu - chart_neu) > 1 or abs(text_neg - chart_neg) > 1:
            blockers.append(
                f"Slide 06 conclusion values ({text_pos}/{text_neu}/{text_neg}) "
                f"mismatch chart values ({chart_pos}/{chart_neu}/{chart_neg}) for project brand."
            )

    if "нейтрализ" in joined.lower() and not data.get("negative_neutralized_by_agency"):
        blockers.append("Slide 06 must mention negative neutralization only when it is confirmed by data.")
    alignment = (result.get("visual_guardrails") or {}).get("chart_alignment") or {}
    if alignment.get("common_top") is None or alignment.get("common_height") is None:
        warnings.append("Slide 06 chart alignment metadata is missing.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=6)
    text_qa = qa_text_contract(joined, slide_number=6)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "project_row_present": bool(data.get("project_row")),
            "daily_rows_present": bool(data.get("daily_rows")),
            "project_total_matches_daily": bool(totals.get("matches")),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_semicolon": ";" not in joined,
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "source_note_without_time": not any(time in source_note for time in ["00:00", "23:59"]),
            "negative_neutralization_guard": "нейтрализ" not in joined.lower() or bool(data.get("negative_neutralized_by_agency")),
            "charts_aligned": alignment.get("common_top") is not None and alignment.get("common_height") is not None,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
    }
