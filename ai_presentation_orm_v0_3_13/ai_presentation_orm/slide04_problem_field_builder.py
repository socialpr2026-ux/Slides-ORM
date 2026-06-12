from __future__ import annotations

from pathlib import Path
import re
import textwrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .chart_style import CHART_COMPACT_SIZE, CHART_FONT
from .native_chart_helpers import add_native_bar_chart, chart_update
from .style_rules import BODY_FONT_SIZE, SOURCE_FONT_SIZE, TITLE_FONT_SIZE, qa_slide_copy_principles
from .text_contracts import slide04_problem_insights, qa_text_contract

EMU_PER_INCH = 914400
DARK_BLUE = "#0070C0"
LIGHT_BLUE = "#1F9BEA"
BAR_GRAY = "#D9D9D9"
BLUE_MARKER = RGBColor(1, 162, 255)
TEXT_COLOR = RGBColor(0, 0, 0)
MONTH_GENITIVE = {
    1: "января",
    2: "февраля",
    3: "марта",
    4: "апреля",
    5: "мая",
    6: "июня",
    7: "июля",
    8: "августа",
    9: "сентября",
    10: "октября",
    11: "ноября",
    12: "декабря",
}


def _fmt_int(value: int | float) -> str:
    try:
        return f"{int(round(float(value))):,}".replace(",", " ")
    except Exception:
        return str(value)


def _fmt_chart_int(value: int | float) -> str:
    try:
        return f"{int(round(float(value))):,}".replace(",", " ")
    except Exception:
        return str(value)


def _messages_word(value: int | float) -> str:
    try:
        number = abs(int(round(float(value))))
    except Exception:
        return "сообщений"
    if 11 <= number % 100 <= 14:
        return "сообщений"
    last = number % 10
    if last == 1:
        return "сообщение"
    if 2 <= last <= 4:
        return "сообщения"
    return "сообщений"


def _ratings_word(value: int | float) -> str:
    try:
        number = abs(int(round(float(value))))
    except Exception:
        return "оценок"
    if 11 <= number % 100 <= 14:
        return "оценок"
    last = number % 10
    if last == 1:
        return "оценка"
    if 2 <= last <= 4:
        return "оценки"
    return "оценок"


def _fmt_pct(value: int | float, digits: int = 1) -> str:
    try:
        return f"{float(value) * 100:.{digits}f}".replace(".", ",") + "%"
    except Exception:
        return "0%"


def _top_names(items: list[dict], count: int = 3) -> str:
    names = [str(item.get("name") or "").strip() for item in items if str(item.get("name") or "").strip()]
    return ", ".join(names[:count])


def _display_label(value: object) -> str:
    text = str(value or "").strip()
    return text[:1].upper() + text[1:] if text else ""


def _sentence_item(value: object) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    if text.startswith("ЛОР"):
        return text
    return text[:1].lower() + text[1:]


def _label_word_count(value: object) -> int:
    return len([token for token in re.split(r"[\s,;/()«»\"®+._-]+", str(value or "").strip().lower()) if token])


def _about_brand(brand: str) -> str:
    text = str(brand or "").strip()
    if not text:
        return "о бренде"
    return ("об " if text[0].lower() in "аоэеиыуёюя" else "о ") + text


def _content_next_step(data: dict, chart_names: str) -> str:
    terms = [str(item.get("name", "")).strip() for item in data.get("chart_terms", [])[:3] if str(item.get("name", "")).strip()]
    if not terms:
        return ""
    term_text = ", ".join(terms).lower()
    if any(token in chart_names for token in ["ангина", "отит", "аденоид", "тонзиллит", "ларингит", "фарингит", "горло", "носоглот", "голос", "осип", "кашель"]):
        return (
            f"Наиболее сильные точки входа — {term_text}. Их стоит связывать с профилактикой ЛОР-инфекций, "
            f"местным иммунитетом и простым объяснением, когда продукт уместен."
        )
    if any(token in chart_names for token in ["микрофлора", "кишечник", "пробиотик", "капсул"]):
        return (
            f"Наиболее сильные точки входа — {term_text}. Их стоит связывать с понятным курсом приема "
            f"и подтверждением эффекта в повседневных ситуациях."
        )
    return ""


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _clear_tf(shape):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    return tf


def _clear_paragraph_bullets(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag.rsplit("}", 1)[-1].startswith("bu"):
            p_pr.remove(child)


def _set_shape_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT):
    shape.text = str(text)
    if getattr(shape, "has_text_frame", False):
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        shape.text_frame.margin_left = Pt(0)
        shape.text_frame.margin_right = Pt(0)
        shape.text_frame.margin_top = Pt(0)
        shape.text_frame.margin_bottom = Pt(0)
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = align
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 1.0
            for run in paragraph.runs:
                run.font.name = "Arial"
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = TEXT_COLOR


def _add_paragraph(tf, text: str, *, font_size: float, bold: bool = False, space_after: float = 1.5):
    p = tf.add_paragraph() if tf.paragraphs and tf.paragraphs[0].text else tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = TEXT_COLOR


def _format_source_period(period_raw: str) -> str:
    text = (period_raw or "").strip()
    numeric_range = re.search(
        r"(\d{1,2})[.](\d{1,2})[.](\d{2,4})(?:\s+\d{1,2}:\d{2}(?::\d{2})?)?\s*[–-]\s*"
        r"(\d{1,2})[.](\d{1,2})[.](\d{2,4})(?:\s+\d{1,2}:\d{2}(?::\d{2})?)?",
        text,
    )
    if numeric_range:
        d1, m1, y1, d2, m2, y2 = [int(part) for part in numeric_range.groups()]
        y1 = 2000 + y1 if y1 < 100 else y1
        y2 = 2000 + y2 if y2 < 100 else y2
        if m1 == m2 and y1 == y2:
            return f"{d1}–{d2} {MONTH_GENITIVE.get(m1, '')} {y1} года".strip()
        return (
            f"{d1} {MONTH_GENITIVE.get(m1, '')} {y1} – "
            f"{d2} {MONTH_GENITIVE.get(m2, '')} {y2} года"
        ).strip()
    text = re.sub(r"\s+\d{1,2}:\d{2}(?::\d{2})?", "", text)
    text = text.replace(" - ", " – ").replace("-", "–")
    if "202" in text and "года" not in text:
        text += " года"
    return text or "период отчета"


def _insight_text(data: dict) -> list[str]:
    return slide04_problem_insights(data)

    brand = data.get("project_brand") or "бренд"
    month_prep = data.get("month_prepositional") or "периоде"
    total = data.get("total_topic_messages") or 0
    chart_names = " ".join(str(item.get("name", "")).lower() for item in data.get("chart_terms", []))
    if any(token in chart_names for token in ["ангина", "отит", "аденоид", "тонзиллит", "ларингит", "фарингит", "горло", "носоглот", "голос", "осип", "кашель"]):
        topic_phrase = "Темы ЛОР-заболеваний, голоса и горла"
    elif any(token in chart_names for token in ["микрофлора", "кишечник", "пробиотик", "капсул"]):
        topic_phrase = "Темы микрофлоры, пробиотиков и работы кишечника"
    else:
        topic_phrase = "Проблемные темы категории"
    contexts = data.get("right_contexts", [])[:4]
    signals = data.get("insight_signals") or {}
    brand_ctx = signals.get("brand_mentions") or {}
    context_text = ", ".join(_sentence_item(item) for item in contexts if str(item or "").strip())
    if not context_text:
        context_text = "ключевых проблемных сценариев"

    fact_interp = (
        f"{topic_phrase} активно обсуждаются: "
        f"в {month_prep} найдено {_fmt_int(total)} {_messages_word(total)}. "
    )

    brand_mentions = int(brand_ctx.get("brand_mentions", 0) or 0)
    if brand_mentions:
        coverage = brand_mentions / total if total else 0
        fact_interp += (
            f"{brand} присутствует в {_fmt_int(brand_mentions)} из них ({_fmt_pct(coverage)}). "
            f"Это только часть доступного инфополя, поэтому есть потенциал роста знания бренда внутри темы."
        )
    else:
        fact_interp += (
            f"{brand} почти не виден в этих обсуждениях. "
            f"Это показывает свободное поле для роста знания бренда внутри темы."
        )

    content_next_step = _content_next_step(data, chart_names)
    action_caveat = (
        f"Рекомендации – выходить в обсуждения, где люди уже говорят про {context_text}. "
        f"В сообщениях связывать бренд с проблемой, моментом выбора и понятным курсом приема."
    )
    if content_next_step:
        action_caveat += " " + content_next_step

    return [fact_interp, action_caveat]


def _set_insight(shape, data: dict):
    tf = _clear_tf(shape)
    for idx, paragraph_text in enumerate(_insight_text(data)):
        p = tf.add_paragraph() if idx else tf.paragraphs[0]
        _clear_paragraph_bullets(p)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8 if idx == 0 else 0)
        p.line_spacing = 1.0
        if paragraph_text.startswith("Рекомендации"):
            head = "Рекомендации"
            rest = paragraph_text[len(head):]
            run = p.add_run()
            run.text = head
            run.font.name = "Arial"
            run.font.size = Pt(BODY_FONT_SIZE)
            run.font.bold = True
            run.font.color.rgb = TEXT_COLOR
            run = p.add_run()
            run.text = rest
            run.font.name = "Arial"
            run.font.size = Pt(BODY_FONT_SIZE)
            run.font.bold = False
            run.font.color.rgb = TEXT_COLOR
        else:
            run = p.add_run()
            run.text = paragraph_text
            run.font.name = "Arial"
            run.font.size = Pt(BODY_FONT_SIZE)
            run.font.bold = False
            run.font.color.rgb = TEXT_COLOR


def _set_context_rail(shape, data: dict):
    contexts = data.get("right_contexts") or [item.get("name") for item in data.get("chart_terms", [])[:5]]
    contexts = [str(item).strip() for item in contexts if str(item or "").strip()][:5]
    tf = _clear_tf(shape)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for idx, context in enumerate(contexts):
        p = tf.add_paragraph() if idx else tf.paragraphs[0]
        _clear_paragraph_bullets(p)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(11 if idx < len(contexts) - 1 else 0)
        p.line_spacing = 1.0
        plus = p.add_run()
        plus.text = "■ "
        plus.font.name = "Arial"
        plus.font.size = Pt(BODY_FONT_SIZE)
        plus.font.bold = False
        plus.font.color.rgb = BLUE_MARKER
        run = p.add_run()
        run.text = context
        run.font.name = "Arial"
        run.font.size = Pt(BODY_FONT_SIZE)
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR


def _render_context_chart(data: dict, out_path: Path, *, width_in: float, height_in: float) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    terms = sorted(data.get("chart_terms", [])[:8], key=lambda item: int(item.get("messages", 0) or 0), reverse=True)
    labels = [textwrap.shorten(_display_label(item.get("name", "")), width=26, placeholder="…") for item in terms]
    values = [int(item.get("messages", 0) or 0) for item in terms]
    if not values:
        labels = ["нет данных"]
        values = [1]

    dpi = 300
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    y = np.arange(len(labels))
    ax.barh(y, values, color=BAR_GRAY, height=0.50)
    ax.set_yticks(y, labels=labels)
    ax.invert_yaxis()
    ax.set_ylim(len(labels) - 0.5, -0.5)
    max_v = max(values) if values else 1
    ax.set_xlim(0, max_v * 1.24)

    for spine in ["top", "right", "bottom"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#D9D9D9")
    ax.spines["left"].set_linewidth(0.8)
    ax.tick_params(axis="x", bottom=False, labelbottom=False)
    ax.tick_params(axis="y", length=0, labelsize=CHART_COMPACT_SIZE, colors="#000000", pad=6)
    for label in ax.get_yticklabels():
        label.set_fontname(CHART_FONT)

    for i, value in enumerate(values):
        ax.text(value + max_v * 0.015, i, _fmt_chart_int(value), va="center", ha="left", fontsize=CHART_COMPACT_SIZE, color="#000000", fontname=CHART_FONT)

    fig.text(
        0.50,
        0.91,
        "Количество упоминаний тематик в месяц",
        ha="center",
        va="top",
        fontsize=CHART_COMPACT_SIZE,
        fontweight="bold",
        fontname=CHART_FONT,
        color="#000000",
    )
    plt.subplots_adjust(left=0.40, right=0.94, top=0.76, bottom=0.08)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=dpi, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def build_slide04_problem_field(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 4:
        raise ValueError("Template does not have slide 4")

    slide = prs.slides[3]
    text_updates = []
    chart_updates = []
    replacements_done = 0

    title = _shape_by_id(slide, 20)
    source = _shape_by_id(slide, 7)
    right_rail = _shape_by_id(slide, 13)
    context_caption = _shape_by_id(slide, 17)
    insight = _shape_by_id(slide, 18)
    chart = next((shape for shape in slide.shapes if getattr(shape, "has_chart", False)), None)

    if title:
        _set_shape_text(title, "ИНФОПОЛЕ ОБСУЖДЕНИЙ ПРОБЛЕМЫ", font_size=TITLE_FONT_SIZE, bold=False)
        text_updates.append({"shape_id": title.shape_id, "role": "title"})
        replacements_done += 1
    if insight:
        _set_insight(insight, data)
        text_updates.append({"shape_id": insight.shape_id, "role": "insight_text"})
        replacements_done += 1
    if context_caption:
        _set_shape_text(context_caption, "Контексты обсуждений брендов категории", font_size=BODY_FONT_SIZE, bold=True)
        text_updates.append({"shape_id": context_caption.shape_id, "role": "context_caption"})
        replacements_done += 1
    if right_rail:
        _set_context_rail(right_rail, data)
        text_updates.append({"shape_id": right_rail.shape_id, "role": "context_rail"})
        replacements_done += 1
    if source:
        source_label = str(data.get('source_label', 'сервис мониторинга Медиалогия')).replace('сервис мониторинга ', '')
        source_text = (
            f"Источник: {source_label}, {_format_source_period(data.get('period_raw'))}. "
            "Исключен спам и нерелевантные упоминания."
        )
        _set_shape_text(source, source_text, font_size=SOURCE_FONT_SIZE, bold=False)
        text_updates.append({"shape_id": source.shape_id, "role": "source_note"})
        replacements_done += 1

    image_outputs = []
    if chart:
        pos = (chart.left, chart.top, chart.width, chart.height)
        existing_titles = []
        for shape in slide.shapes:
            if shape is chart or not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text or "").strip().lower()
            if "количество упоминаний тематик" in text:
                cx = shape.left + shape.width // 2
                cy = shape.top + shape.height // 2
                if pos[0] - int(0.2 * EMU_PER_INCH) <= cx <= pos[0] + pos[2] + int(0.2 * EMU_PER_INCH) and pos[1] - int(0.5 * EMU_PER_INCH) <= cy <= pos[1] + int(0.7 * EMU_PER_INCH):
                    existing_titles.append(shape)
        _remove_shape(slide, chart)
        chart_title = existing_titles[0] if existing_titles else slide.shapes.add_textbox(pos[0], pos[1], pos[2], int(0.24 * EMU_PER_INCH))
        for duplicate in existing_titles[1:]:
            _remove_shape(slide, duplicate)
        chart_title.left = pos[0]
        chart_title.top = pos[1]
        chart_title.width = pos[2]
        chart_title.height = int(0.24 * EMU_PER_INCH)
        _set_shape_text(chart_title, "Количество упоминаний тематик в месяц", font_size=10, bold=True, align=PP_ALIGN.CENTER)
        chart_box = (pos[0], pos[1] + int(0.30 * EMU_PER_INCH), pos[2], max(int(0.4 * EMU_PER_INCH), pos[3] - int(0.30 * EMU_PER_INCH)))
        terms = sorted(data.get("chart_terms", [])[:8], key=lambda item: int(item.get("messages", 0) or 0), reverse=True)
        def _axis_label(value: str) -> str:
            label = _display_label(value)
            wrapped = textwrap.wrap(label, width=23, break_long_words=False, break_on_hyphens=False)
            return "\n".join(wrapped[:3]) if wrapped else label

        labels = [_axis_label(item.get("name", "")) for item in terms] or ["Нет данных"]
        values = [int(item.get("messages", 0) or 0) for item in terms] or [0]
        add_native_bar_chart(
            slide,
            box=chart_box,
            categories=labels,
            values=values,
            color=BAR_GRAY,
            name="Количество упоминаний",
            reverse=True,
        )
        chart_updates.append(chart_update("problem_contexts", "native_horizontal_bar", categories=labels))
        replacements_done += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    text_after = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text.strip()
            if txt:
                text_after.append(txt)
    joined = "\n".join(text_after)
    current_brand = (project.get("brand") or data.get("project_brand") or "").lower()
    forbidden = ["Пикторид", "Бактоблис", "декабрь 2025", "ДЕКАБР", "ЛОР-органов"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower() and item.lower() != current_brand]

    return {
        "output": str(output_path),
        "brand": project.get("brand") or data.get("project_brand"),
        "month_caption": data.get("month_caption"),
        "replacement_count": replacements_done,
        "text_updates": text_updates,
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "old_text_left": old_text_left,
        "source_note": source.text if source else "",
        "insight_text": "\n".join(_insight_text(data)),
        "chart_term_count": len(data.get("chart_terms", [])),
        "text_after": text_after,
        "risk_note": "Slide 04 is a context slice. It does not claim campaign causality.",
        "visual_guardrails": {
            "chart_geometry": "uses inherited reference chart slot",
            "bar_colors": {"default": BAR_GRAY},
            "right_rail_marker": "blue square markers follow final-slide brand grammar",
        },
    }


def qa_slide04_problem_field(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    terms = data.get("chart_terms", [])
    chart_value_sum = sum(int(item.get("messages", 0) or 0) for item in terms)
    unique_message_count = int(data.get("chart_unique_message_count") or 0)
    uses_exclusive_rows = data.get("chart_method") == "mutually_exclusive_message_rows"
    if not (6 <= len(terms) <= 8):
        blockers.append(f"Chart must contain 6-8 readable context terms; got {len(terms)}.")
    if any(int(item.get("messages", 0) or 0) <= 0 for item in terms):
        blockers.append("Chart contains zero or missing context values.")
    if uses_exclusive_rows and unique_message_count and chart_value_sum > unique_message_count:
        blockers.append("Slide 04 exclusive context chart double-counts messages.")
    short_labels = [str(item.get("name", "")).strip() for item in terms if _label_word_count(item.get("name", "")) < 2]
    if short_labels:
        blockers.append("Chart contains one-word labels that are not analytical contexts: " + ", ".join(short_labels[:4]))
    lowercase_labels = [
        str(item.get("name", "")).strip()
        for item in terms
        if str(item.get("name", "")).strip()[:1].islower()
    ]
    if lowercase_labels:
        blockers.append("Chart labels must start with a capital letter: " + ", ".join(lowercase_labels[:4]))
    context_text = " | ".join(str(item or "") for item in data.get("right_contexts", []))
    if "ребенок, дети" in context_text.lower() or "ребёнок, дети" in context_text.lower():
        blockers.append("Context rail contains a keyword list instead of a semantic phrase.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 04: " + ", ".join(result["old_text_left"]))
    source_note = result.get("source_note", "")
    insight_low = result.get("insight_text", "").lower()
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    joined = "\n".join(result.get("text_after") or [])
    if "????" in joined:
        blockers.append("Broken Cyrillic marker found on slide 04.")
    if ";" in joined:
        blockers.append("Slide 04 text must not use semicolons.")
    causal_fragments = ["кампания дала", "кампания привела", "кампания обеспечила", "благодаря кампании"]
    if any(fragment in insight_low for fragment in causal_fragments):
        blockers.append("Slide 04 must not present campaign presence as proven causality.")
    if not data.get("total_topic_messages"):
        warnings.append("Total topic message count is missing.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=4)
    text_qa = qa_text_contract(joined, slide_number=4)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])

    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "chart_terms_6_to_8": 6 <= len(terms) <= 8,
            "chart_values_nonzero": all(int(item.get("messages", 0) or 0) > 0 for item in terms),
            "chart_uses_exclusive_rows": uses_exclusive_rows,
            "chart_values_within_unique_messages": not (uses_exclusive_rows and unique_message_count and chart_value_sum > unique_message_count),
            "chart_labels_are_contexts": not short_labels,
            "chart_labels_capitalized": not lowercase_labels,
            "right_contexts_are_semantic": "ребенок, дети" not in context_text.lower() and "ребёнок, дети" not in context_text.lower(),
            "no_stale_template_text": not result.get("old_text_left"),
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "campaign_causality_guard": not any(fragment in insight_low for fragment in causal_fragments),
            "no_broken_cyrillic": "????" not in joined,
            "no_semicolon": ";" not in joined,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
    }
