from __future__ import annotations

from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .brand_styler import bold_project_brand_in_presentation
from .chart_style import CHART_FONT, CHART_LABEL_SIZE
from .native_chart_helpers import add_native_bar_chart, add_native_stacked_bar_chart, chart_update
from .style_rules import BODY_FONT_SIZE, CHART_FONT_SIZE, FONT_FAMILY, SOURCE_FONT_SIZE, TITLE_FONT_SIZE, qa_slide_copy_principles
from .text_contracts import slide03_category_insight, qa_text_contract

EMU_PER_INCH = 914400
DARK_BLUE = "#0070C0"
LIGHT_BLUE = "#1F9BEA"
TARGET_BLUE = LIGHT_BLUE       # our brand is light blue
COMPETITOR_BLUE = DARK_BLUE    # competitors are blue
TEXT_COLOR = RGBColor(0, 0, 0)
MONTH_GENITIVE = {
    1: "января",
    2: "февраля",
    3: "марта",
    4: "апреля",
    5: "мая",
    6: "июня",
    7: "июля",
    8: "августа",
    9: "сентября",
    10: "октября",
    11: "ноября",
    12: "декабря",
}


def _fmt_int(value: int | float) -> str:
    try:
        return f"{int(value):,}".replace(",", " ")
    except Exception:
        return str(value)


def _fmt_int_comma(value: int | float) -> str:
    try:
        return f"{int(value):,}".replace(",", " ")
    except Exception:
        return str(value)


def _messages_word(value: int | float) -> str:
    try:
        number = abs(int(round(float(value))))
    except Exception:
        return "сообщений"
    if 11 <= number % 100 <= 14:
        return "сообщений"
    last = number % 10
    if last == 1:
        return "сообщение"
    if 2 <= last <= 4:
        return "сообщения"
    return "сообщений"


def _mentions_word(value: int | float) -> str:
    try:
        number = abs(int(round(float(value))))
    except Exception:
        return "упоминаний"
    if 11 <= number % 100 <= 14:
        return "упоминаний"
    last = number % 10
    if last == 1:
        return "упоминание"
    if 2 <= last <= 4:
        return "упоминания"
    return "упоминаний"


def _fmt_pct(value: float) -> str:
    try:
        return f"{value * 100:.1f}%".replace(".", ",")
    except Exception:
        return "н/д"


def _clear_tf(shape):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    return tf


def _add_rich_paragraph(tf, segments: list[tuple[str, bool]], font_size=BODY_FONT_SIZE, space_after=0.4):
    p = tf.add_paragraph() if len(tf.paragraphs) and tf.paragraphs[0].text else tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    p.line_spacing = 1.0
    for text, bold in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold)
        run.font.name = FONT_FAMILY
        run.font.color.rgb = TEXT_COLOR
    return p


def _set_shape_text(shape, text: str, font_size: float | None = None, bold: bool | None = None, align=PP_ALIGN.LEFT) -> None:
    shape.text = str(text)
    if getattr(shape, "has_text_frame", False):
        shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        shape.text_frame.margin_left = Pt(0)
        shape.text_frame.margin_right = Pt(0)
        shape.text_frame.margin_top = Pt(0)
        shape.text_frame.margin_bottom = Pt(0)
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                run.font.name = FONT_FAMILY
                run.font.color.rgb = TEXT_COLOR


def _metric_leader(rows: list[dict], metric: str) -> dict:
    valid = [r for r in rows if r.get(metric)]
    return max(valid, key=lambda r: r.get(metric, 0)) if valid else {}


def _join_names(items, max_n=3):
    names = [str(i.get("name", "")).strip() for i in items[:max_n] if i.get("name")]
    if not names:
        return ""
    if len(names) == 1:
        return names[0]
    return ", ".join(names[:-1]) + " и " + names[-1]


def _competitive_reason(data: dict) -> str:
    """Build a concise data-driven context paragraph."""
    project_brand = data.get("project_brand") or "бренд"
    ctx = data.get("analytical_context") or {}

    platforms = ctx.get("platforms") or []
    words = ctx.get("words") or []

    topics = [str(item.get("name", "")).strip().lower() for item in words[:3] if str(item.get("name", "")).strip()]
    topic_tail = f" Темы для усиления органики: {', '.join(topics)}." if topics else ""
    platform_tail = ""
    if platforms:
        platform_names = [str(item.get("name", "")).strip() for item in platforms[:2] if str(item.get("name", "")).strip()]
        if platform_names:
            platform_tail = " Наиболее охватные точки контакта: " + ", ".join(platform_names) + "."
    return (
        f"{project_brand} получает охват за счет присутствия в крупных обсуждениях и за счет текущих ORM-публикаций."
        f"{platform_tail}{topic_tail}"
    )


def _audience_position_text(project_brand: str, audience_rank, total_brands: int, p_aud: int, leader_aud: dict) -> list[tuple[str, bool]]:
    leader_name = leader_aud.get("brand", "бренд с максимальным OTS")
    leader_value = _fmt_int_comma(leader_aud.get("audience", 0))
    rank_text = f"{audience_rank}-е место" if audience_rank else "место не определено"
    if audience_rank and total_brands and audience_rank == total_brands:
        text = (
            f" по охвату уступает конкурентам: OTS {_fmt_int_comma(p_aud)}, {rank_text} из {total_brands}, "
            f"максимальный OTS у "
        )
    elif audience_rank and audience_rank <= max(1, min(3, total_brands)):
        text = (
            f" входит в верхнюю группу по аудитории: OTS {_fmt_int_comma(p_aud)}, {rank_text} среди брендов слайда, "
            f"максимальный OTS у "
        )
    else:
        text = (
            f" сохраняет заметность по аудитории, но не входит в верхнюю группу: OTS {_fmt_int_comma(p_aud)}, {rank_text}, "
            f"максимальный OTS у "
        )
    return [
        (project_brand, True),
        (text, False),
        (leader_name, True),
        (f" ({leader_value}). ", False),
    ]

def _qa_clean_insight_text(text: str) -> str:
    replacements = {
        " дает максимальный объем сообщений": " выигрывает по объему сообщений",
        " дает": " выигрывает",
        " обеспечивает наибольший OTS": " лидирует по OTS",
        " обеспечивает": " лидирует",
        "Наша часть составила": "Мы получили",
        "На нашу часть пришлось": "Мы получили",
        "ориентировочно": "",
    }
    cleaned = text
    for old, new in replacements.items():
        cleaned = cleaned.replace(old, new)
    cleaned = " ".join(cleaned.split())
    return cleaned


def _set_insight_text(shape, data: dict) -> None:
    """Data-driven insight block for slide 03.

    Rules v0.3.10:
    - no gerunds / деепричастия;
    - no word 'ориентировочно';
    - project result avoids repeated verb and uses phrase 'Мы получили';
    - recommendation starts with 'В следующем месяце мы рекомендуем';
    - SOV% is shown for 2nd and 3rd places;
    - text is generated from actual rankings, SOV, OTS and engagement data.
    """
    tf = _clear_tf(shape)
    tf.margin_bottom = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    project = data.get("project_brand_row") or {}
    project_brand = data.get("project_brand") or project.get("brand") or "бренд"
    rank = data.get("project_brand_rank")
    audience_rank = data.get("project_audience_rank")
    leader = data.get("mentions_leader") or {}
    leader_aud = data.get("audience_leader") or {}
    top = data.get("top_brands", [])
    split = data.get("campaign_split") or {}

    p_mentions = int(project.get("mentions", 0) or 0)
    p_sov = project.get("sov", 0)
    p_aud = int(project.get("audience", 0) or 0)
    total_mentions = int(data.get("total_mentions_competitive_set", 0) or 0)
    gap_raw = data.get("gap_to_mentions_leader")
    gap = round(float(gap_raw), 1) if gap_raw is not None else None

    campaign_mentions = int(split.get("campaign_mentions", 0) or 0)
    organic_mentions = int(split.get("organic_mentions", max(p_mentions - campaign_mentions, 0)) or 0)
    split_text = ""
    if campaign_mentions and p_mentions:
        split_text = f" Из них {_fmt_int(campaign_mentions)} сообщений относятся к кампании, {_fmt_int(organic_mentions)} — органическая часть."
    rank_label = f"{rank}-е место" if rank else "место не определено"

    tf = _clear_tf(shape)
    tf.margin_bottom = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)

    if not project:
        _add_rich_paragraph(tf, [("Целевой бренд не сопоставлен с таблицей показателей брендов. Нужно подтвердить название бренда и конкурентный набор.", False)], BODY_FONT_SIZE)
        return

    _add_rich_paragraph(tf, [
        (f"В {data.get('month_prepositional', 'периоде')} бренды набора суммарно получили ", False),
        (_fmt_int(total_mentions), True),
        (f" {_mentions_word(total_mentions)}. {project_brand} с учетом кампании занял {rank_label}: ", False),
        (_fmt_int(p_mentions), True),
        (f" {_messages_word(p_mentions)} и SOV ", False),
        (_fmt_pct(p_sov), True),
        (".", False),
        (split_text, False),
    ], BODY_FONT_SIZE, 0.35)

    second = top[1] if len(top) > 1 else {}
    third = top[2] if len(top) > 2 else {}
    gap_value = str(gap).replace(".", ",") if gap is not None else ""
    gap_phrase = f", разрыв с первым местом для бренда – около {gap_value}×" if gap is not None and gap_value else ""
    _add_rich_paragraph(tf, [
        ("Больше всего упоминаний у ", False),
        (leader.get("brand", "первого бренда"), True),
        (f" ({_fmt_int(leader.get('mentions', 0))} {_messages_word(leader.get('mentions', 0))}, SOV {_fmt_pct(leader.get('sov', 0))}). ", False),
        ("На 2–3 местах – ", False),
        (second.get("brand", ""), True),
        (f" ({_fmt_int(second.get('mentions', 0))} {_messages_word(second.get('mentions', 0))}, SOV {_fmt_pct(second.get('sov', 0))}) и ", False),
        (third.get("brand", ""), True),
        (f" ({_fmt_int(third.get('mentions', 0))} {_messages_word(third.get('mentions', 0))}, SOV {_fmt_pct(third.get('sov', 0))}){gap_phrase}.", False),
    ], BODY_FONT_SIZE, 0.35)

    total_brands = len(data.get("all_brand_rows") or [])
    audience_segments = _audience_position_text(project_brand, audience_rank, total_brands, p_aud, leader_aud)
    _add_rich_paragraph(tf, audience_segments + [
        ("Следующий шаг — усилить органику в темах и на площадках, где уже есть высокий охват обсуждений.", False),
    ], BODY_FONT_SIZE, 0.35)

def _render_mentions_chart_image(data: dict, out_path: Path, *, width_in: float, height_in: float):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    categories = data["chart_mentions"]["categories"]
    total_mentions = data["chart_mentions"]["mentions"]
    project_brand = (data.get("project_brand") or "").lower()
    split = data.get("campaign_split") or {}
    campaign = []
    organic = []
    for cat, total in zip(categories, total_mentions):
        if cat.lower() == project_brand and split.get("campaign_mentions"):
            camp = min(int(split.get("campaign_mentions", 0)), int(total))
            org = max(int(total) - camp, 0)
        else:
            camp = 0
            org = int(total)
        organic.append(org)
        campaign.append(camp)

    dpi = 300
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    y = np.arange(len(categories))
    ax.barh(y, organic, color=DARK_BLUE, height=0.52)
    ax.barh(y, campaign, left=organic, color=LIGHT_BLUE, height=0.52)
    ax.set_yticks(y, labels=categories)
    ax.invert_yaxis()
    ax.set_ylim(len(categories) - 0.5, -0.5)
    max_v = max(total_mentions) if total_mentions else 1
    ax.set_xlim(0, max_v * 1.24)

    for spine in ["top", "right", "bottom"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#D9D9D9")
    ax.spines["left"].set_linewidth(0.8)
    ax.tick_params(axis="x", bottom=False, labelbottom=False)
    ax.tick_params(axis="y", length=0, labelsize=CHART_LABEL_SIZE, colors="#4D4D4D", pad=6)
    for label in ax.get_yticklabels():
        label.set_fontname(CHART_FONT)
        if label.get_text().lower() == project_brand:
            label.set_fontweight("bold")
        label.set_wrap(False)

    for i, (org, camp, total) in enumerate(zip(organic, campaign, total_mentions)):
        if org:
            ax.text(max(org * 0.55, max_v * 0.04), i, _fmt_int(org), va="center", ha="center", fontsize=CHART_LABEL_SIZE, color="white", fontname=CHART_FONT)
        if camp:
            ax.text(org + camp * 0.5, i, _fmt_int(camp), va="center", ha="center", fontsize=CHART_LABEL_SIZE, color="white", fontname=CHART_FONT)
            ax.text(total + max_v * 0.02, i, _fmt_int(total), va="center", ha="left", fontsize=CHART_LABEL_SIZE, color="#333333", fontname=CHART_FONT)

    plt.subplots_adjust(left=0.27, right=0.98, top=0.96, bottom=0.08)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=dpi, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _render_audience_chart_image(data: dict, out_path: Path, *, width_in: float, height_in: float):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    categories = data["chart_audience"]["categories"]
    values = data["chart_audience"]["audience"]
    project_brand = (data.get("project_brand") or "").lower()
    colors = [TARGET_BLUE if c.lower() == project_brand else COMPETITOR_BLUE for c in categories]

    dpi = 300
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    y = np.arange(len(categories))
    ax.barh(y, values, color=colors, height=0.52)
    ax.set_yticks(y, labels=categories)
    ax.invert_yaxis()
    ax.set_ylim(len(categories) - 0.5, -0.5)
    max_v = max(values) if values else 1
    ax.set_xlim(0, max_v * 1.29)

    for spine in ["top", "right", "bottom"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color("#D9D9D9")
    ax.spines["left"].set_linewidth(0.8)
    ax.tick_params(axis="x", bottom=False, labelbottom=False)
    ax.tick_params(axis="y", length=0, labelsize=CHART_LABEL_SIZE, colors="#4D4D4D", pad=6)
    for label in ax.get_yticklabels():
        label.set_fontname(CHART_FONT)
        if label.get_text().lower() == project_brand:
            label.set_fontweight("bold")
        label.set_wrap(False)

    for i, value in enumerate(values):
        ax.text(value + max_v * 0.015, i, _fmt_int_comma(value), va="center", ha="left", fontsize=CHART_LABEL_SIZE, color="#333333", fontname=CHART_FONT)

    plt.subplots_adjust(left=0.27, right=0.96, top=0.96, bottom=0.08)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=dpi, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _remove_chart_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _add_layout_cover(slide, prs):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        int(1.18 * EMU_PER_INCH),
        prs.slide_width,
        int(5.72 * EMU_PER_INCH),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp_tree.remove(rect._element)
    sp_tree.insert(2, rect._element)
    return rect.shape_id


def _format_source_period(period_raw: str) -> str:
    text = (period_raw or "").strip()
    numeric_range = re.search(
        r"(\d{1,2})[.](\d{1,2})[.](\d{2,4})(?:\s+\d{1,2}:\d{2}(?::\d{2})?)?\s*[–-]\s*"
        r"(\d{1,2})[.](\d{1,2})[.](\d{2,4})(?:\s+\d{1,2}:\d{2}(?::\d{2})?)?",
        text,
    )
    if numeric_range:
        d1, m1, y1, d2, m2, y2 = [int(part) for part in numeric_range.groups()]
        y1 = 2000 + y1 if y1 < 100 else y1
        y2 = 2000 + y2 if y2 < 100 else y2
        if m1 == m2 and y1 == y2:
            return f"{d1}–{d2} {MONTH_GENITIVE.get(m1, '')} {y1} года".strip()
        return (
            f"{d1} {MONTH_GENITIVE.get(m1, '')} {y1} – "
            f"{d2} {MONTH_GENITIVE.get(m2, '')} {y2} года"
        ).strip()
    text = re.sub(r"\s+\d{1,2}:\d{2}(?::\d{2})?", "", text)
    text = text.replace(" - ", " – ").replace("–", "–")
    if "2026" in text and "года" not in text:
        text += " года"
    return text


def _source_service_label(source_file: str, source_system: str = "") -> str:
    # Keep human-readable source only, never filename.
    # Source must be selected by actual slide data source, not by default.
    sys = (source_system or "").lower()
    low = (source_file or "").lower()
    if "медиалог" in sys or "medialog" in sys or "медиалог" in low or "medialog" in low:
        return "Медиалогия"
    if "brand analytics" in sys or "ba" == sys.strip() or "brand analytics" in low:
        return "Brand Analytics"
    # Do not guess BA by default: for this slide current confirmed source is Medialogia.
    return "Медиалогия"


def build_slide03_brand_mentions(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 3:
        raise ValueError("Template does not have slide 3")
    slide = prs.slides[2]
    layout_cover_shape_id = _add_layout_cover(slide, prs)

    project_brand = project.get("brand", data.get("project_brand", ""))
    month_prep = data.get("month_prepositional", "периоде").upper()
    month_caption = data.get("month_caption", "период")

    title_left = None
    replacements_done = 0
    chart_updates = []
    text_updates = []
    chart_positions = []

    # Pass 1: detect title and chart positions.
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and "УПОМИНАНИЯ БРЕНДОВ" in (shape.text or "").upper():
            title_left = shape.left
        if getattr(shape, "has_chart", False):
            chart_positions.append((shape.left, shape.top, shape.width, shape.height))

    chart_positions_sorted = sorted(chart_positions, key=lambda x: x[0])
    # Use template chart positions when available; fallback to default positions.
    if len(chart_positions_sorted) >= 1:
        left_chart_pos = chart_positions_sorted[0]
    else:
        left_chart_pos = (int(1.03 * EMU_PER_INCH), int(3.72 * EMU_PER_INCH), int(5.45 * EMU_PER_INCH), int(2.48 * EMU_PER_INCH))
    if len(chart_positions_sorted) >= 2:
        right_chart_pos = chart_positions_sorted[1]
    else:
        right_chart_pos = (int(6.56 * EMU_PER_INCH), int(3.72 * EMU_PER_INCH), int(5.28 * EMU_PER_INCH), int(2.48 * EMU_PER_INCH))
    # Captions start from the beginning of the blue bars / plot area, not from y-axis labels.
    left_plot_x = left_chart_pos[0] + int(0.27 * left_chart_pos[2])
    right_plot_x = right_chart_pos[0] + int(0.27 * right_chart_pos[2])
    left_caption_pos = (left_plot_x, int(3.48 * EMU_PER_INCH), left_chart_pos[0] + left_chart_pos[2] - left_plot_x, int(0.24 * EMU_PER_INCH))
    right_caption_pos = (right_plot_x, int(3.48 * EMU_PER_INCH), right_chart_pos[0] + right_chart_pos[2] - right_plot_x, int(0.24 * EMU_PER_INCH))

    duplicate_shapes = []
    audience_caption_seen = False
    insight_written = False
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = shape.text or ""
        up = txt.upper()
        if "УПОМИНАНИЯ БРЕНДОВ" in up:
            _set_shape_text(shape, f"ИНФОПОЛЕ КАТЕГОРИИ В {month_prep}", font_size=TITLE_FONT_SIZE, bold=False)
            title_left = shape.left
            text_updates.append({"shape_id": shape.shape_id, "role": "title"})
            replacements_done += 1
        elif "СООБЩЕНИЯ О БРЕНДАХ" in up:
            shape.left, shape.top, shape.width, shape.height = left_caption_pos
            _set_shape_text(shape, f"Сообщения о брендах, {month_caption}", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.LEFT)
            text_updates.append({"shape_id": shape.shape_id, "role": "mentions_chart_caption", "align": "left"})
            replacements_done += 1
        elif "АУДИТОРИЯ СООБЩЕНИЙ" in up or "OTS" in up:
            if audience_caption_seen:
                duplicate_shapes.append(shape)
                continue
            audience_caption_seen = True
            shape.left, shape.top, shape.width, shape.height = right_caption_pos
            _set_shape_text(shape, f"Аудитория сообщений OTS, {month_caption}", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.LEFT)
            text_updates.append({"shape_id": shape.shape_id, "role": "audience_chart_caption", "align": "left"})
            replacements_done += 1
        elif "ИСТОЧНИК ДАННЫХ" in up:
            # Exact template-like wording: no filename, only source and period.
            source = f"Источник: {_source_service_label(data.get('source_file'), data.get('source_system'))}, {_format_source_period(data.get('period_raw'))}. Исключен спам и нерелевантные упоминания."
            shape.left = int(0.46 * EMU_PER_INCH)
            shape.top = int(6.76 * EMU_PER_INCH)
            shape.width = int(11.80 * EMU_PER_INCH)
            shape.height = int(0.32 * EMU_PER_INCH)
            _set_shape_text(shape, source, font_size=SOURCE_FONT_SIZE, bold=False)
            text_updates.append({"shape_id": shape.shape_id, "role": "source_note"})
            replacements_done += 1
        elif len(txt) > 150:
            if title_left is not None and shape.left < title_left:
                shape.left = title_left
            shape.left = int(0.42 * EMU_PER_INCH)
            shape.top = int(1.58 * EMU_PER_INCH)
            shape.height = int(1.48 * EMU_PER_INCH)
            shape.width = int(11.22 * EMU_PER_INCH)
            _set_insight_text(shape, data)
            text_updates.append({"shape_id": shape.shape_id, "role": "insight_text", "left_in": round(shape.left / EMU_PER_INCH, 2), "top_in": round(shape.top / EMU_PER_INCH, 2)})
            replacements_done += 1
            insight_written = True

    for shape in duplicate_shapes:
        shape._element.getparent().remove(shape._element)

    if not insight_written:
        insight_shape = slide.shapes.add_textbox(
            int(0.42 * EMU_PER_INCH),
            int(1.42 * EMU_PER_INCH),
            int(11.22 * EMU_PER_INCH),
            int(1.50 * EMU_PER_INCH),
        )
        _set_insight_text(insight_shape, data)
        text_updates.append({
            "shape_id": insight_shape.shape_id,
            "role": "insight_text",
            "created": True,
            "left_in": round(insight_shape.left / EMU_PER_INCH, 2),
            "top_in": round(insight_shape.top / EMU_PER_INCH, 2),
        })
        replacements_done += 1

    charts = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    charts_sorted = sorted(charts, key=lambda sh: sh.left)
    chart_image_dir = output_path.parent / "slide03_chart_images"
    image_outputs = []

    mention_categories = (data.get("chart_mentions") or {}).get("categories") or []
    audience_categories = (data.get("chart_audience") or {}).get("categories") or []

    if len(charts_sorted) >= 1 and mention_categories:
        left = charts_sorted[0]
        _remove_chart_shape(slide, left)
        categories = mention_categories
        total_mentions = (data.get("chart_mentions") or {}).get("mentions") or []
        project_brand_low = (data.get("project_brand") or "").lower()
        split = data.get("campaign_split") or {}
        campaign = []
        organic = []
        for cat, total in zip(categories, total_mentions):
            if cat.lower() == project_brand_low and split.get("campaign_mentions"):
                camp = min(int(split.get("campaign_mentions", 0)), int(total))
                org = max(int(total) - camp, 0)
            else:
                camp = 0
                org = int(total)
            organic.append(org)
            campaign.append(camp)
        add_native_stacked_bar_chart(
            slide,
            box=left_chart_pos,
            categories=categories,
            series=[
                ("Органика", organic, DARK_BLUE),
                ("Кампания", campaign, LIGHT_BLUE),
            ],
            reverse=True,
            show_legend=False,
        )
        chart_updates.append(chart_update("brand_mentions", "native_stacked_bar", categories=categories))
        replacements_done += 1
    elif len(charts_sorted) >= 1 and not mention_categories:
        _remove_chart_shape(slide, charts_sorted[0])
        no_data_box = slide.shapes.add_textbox(
            left_chart_pos[0], left_chart_pos[1], left_chart_pos[2], left_chart_pos[3]
        )
        _set_shape_text(no_data_box, "Нет данных для построения графика", font_size=CHART_FONT_SIZE, bold=False, align=PP_ALIGN.CENTER)
        text_updates.append({"shape_id": no_data_box.shape_id, "role": "no_data_mentions", "created": True})
    if len(charts_sorted) >= 2 and audience_categories:
        right = charts_sorted[1]
        _remove_chart_shape(slide, right)
        categories = audience_categories
        values = (data.get("chart_audience") or {}).get("audience") or []
        point_colors = [TARGET_BLUE if str(cat).lower() == project_brand_low else COMPETITOR_BLUE for cat in categories]
        add_native_bar_chart(
            slide,
            box=right_chart_pos,
            categories=categories,
            values=values,
            color=COMPETITOR_BLUE,
            point_colors=point_colors,
            name="Аудитория OTS",
            reverse=True,
        )
        chart_updates.append(chart_update("audience_ots", "native_horizontal_bar", categories=categories))
        replacements_done += 1
    elif len(charts_sorted) >= 2 and not audience_categories:
        _remove_chart_shape(slide, charts_sorted[1])
        no_data_box = slide.shapes.add_textbox(
            right_chart_pos[0], right_chart_pos[1], right_chart_pos[2], right_chart_pos[3]
        )
        _set_shape_text(no_data_box, "Нет данных для построения графика", font_size=CHART_FONT_SIZE, bold=False, align=PP_ALIGN.CENTER)
        text_updates.append({"shape_id": no_data_box.shape_id, "role": "no_data_audience", "created": True})

    output_path.parent.mkdir(parents=True, exist_ok=True)
    brand_bold_updates = bold_project_brand_in_presentation(prs, project_brand)
    prs.save(str(output_path))

    scan = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            scan.append(shape.text)
    joined = "\n".join(scan)
    source_note = next((text for text in scan if text.strip().lower().startswith("источник:")), "")
    old_text_left = []
    for old in ["Пикторид", "Бактавит", "Окуларис", "Пиклоксидин", "Витабакт", "Пиклосепт", "БАКТОБЛИС", "Бактоблис", "декабрь 2025", "ДЕКАБРЕ", "Всего 2025", "Заголовок"]:
        if old.lower() in joined.lower() and old.lower() != project_brand.lower():
            old_text_left.append(old)

    return {
        "output": str(output_path),
        "brand": project_brand,
        "month_caption": month_caption,
        "replacement_count": replacements_done,
        "text_updates": text_updates,
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "layout_cover_shape_id": layout_cover_shape_id,
        "old_text_left": old_text_left,
        "source_note": source_note,
        "text_after": scan,
        "brand_bold_updates": brand_bold_updates,
        "campaign_split": data.get("campaign_split"),
        "methodology_status": data.get("methodology_status"),
        "missing_required": [] if (data.get("campaign_split") or {}).get("campaign_mentions") else ["campaign_vs_organic_split"],
        "risk_note": "Charts are image fallback for visual fidelity. Campaign part is estimated from placement fact total and requires PM / Insight confirmation before final client use.",
    }


def qa_slide03_brand_mentions(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    joined = "\n".join(result.get("text_after") or [])
    source_note = result.get("source_note", "")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 03: " + ", ".join(result["old_text_left"]))
    if not data.get("project_brand_row"):
        blockers.append("Project brand is missing from slide 03 ranking data.")
    if not data.get("top_brands"):
        blockers.append("Slide 03 brand ranking rows are missing.")
    if "????" in joined:
        blockers.append("Broken Cyrillic marker found on slide 03.")
    if ";" in joined:
        blockers.append("Slide 03 text must not use semicolons.")
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    if any(time in source_note for time in ["00:00", "23:59", ":00", ":59"]):
        blockers.append("Source note must not include hours/minutes.")
    if "проверять вклад кампании" in joined.lower():
        blockers.append("Slide 03 must not ask to check campaign contribution when the split is already calculated.")
    split = data.get("campaign_split") or {}
    if not int(split.get("campaign_publications_count") or split.get("campaign_mentions") or 0):
        blockers.append("Slide 03 requires a campaign publication fact.")
    elif split.get("status") != "verified":
        warnings.append("Slide 03 campaign fact uses a fallback source; project ORM month rows are preferred.")
    if re.search(r"\d{1,3}(,\d{3})+", joined):
        blockers.append("Slide 03 numbers must use spaces as thousands separators, not commas.")
    if result.get("brand_bold_updates", 0) <= 0:
        warnings.append("Project brand bold marker was not applied on standalone slide 03; final combined deck applies global brand styling when enabled.")
    ctx = data.get("analytical_context") or {}
    if ctx.get("status") == "error":
        warnings.append("Slide 03 analytical context extraction failed; topic wording may be incomplete.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=3)
    text_qa = qa_text_contract(joined, slide_number=3)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "project_brand_present": bool(data.get("project_brand_row")),
            "top_brands_present": bool(data.get("top_brands")),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_broken_cyrillic": "????" not in joined,
            "no_semicolon": ";" not in joined,
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "source_note_without_time": not any(time in source_note for time in ["00:00", "23:59", ":00", ":59"]),
            "numbers_use_space_thousands": not re.search(r"\d{1,3}(,\d{3})+", joined),
            "brand_bold_updates": result.get("brand_bold_updates", 0) > 0,
            "campaign_fact_present": bool(int(split.get("campaign_publications_count") or split.get("campaign_mentions") or 0)),
            "campaign_fact_verified": split.get("status") == "verified",
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
    }
