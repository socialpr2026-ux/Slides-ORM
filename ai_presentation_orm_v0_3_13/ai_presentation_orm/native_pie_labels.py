from __future__ import annotations

import math
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .chart_style import CHART_COMPACT_SIZE, CHART_LABEL_SIZE, CHART_FONT

EMU_PER_PT = 12700


def _fmt_int(value: int | float) -> str:
    return f"{int(round(float(value))):,}".replace(",", " ")


def _fmt_pct(value: float, digits: int = 1) -> str:
    return f"{float(value) * 100:.{digits}f}%".replace(".", ",")


def _rgb(color: str) -> RGBColor:
    text = str(color or "#BFBFBF").strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    try:
        return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except Exception:
        return RGBColor(191, 191, 191)


def _boxes_overlap(a: tuple[int, int, int, int], b: tuple[int, int, int, int], *, gap: int = 0) -> bool:
    return not (a[2] + gap <= b[0] or b[2] + gap <= a[0] or a[3] + gap <= b[1] or b[3] + gap <= a[1])


def _count_overlaps(boxes: list[tuple[int, int, int, int]]) -> int:
    count = 0
    gap = int(2 * EMU_PER_PT)
    for idx, box in enumerate(boxes):
        for other in boxes[idx + 1:]:
            if _boxes_overlap(box, other, gap=gap):
                count += 1
    return count


def _count_out_of_bounds(boxes: list[tuple[int, int, int, int]], bounds: tuple[int, int, int, int]) -> int:
    bx1, by1, bx2, by2 = bounds
    return sum(1 for x1, y1, x2, y2 in boxes if x1 < bx1 or y1 < by1 or x2 > bx2 or y2 > by2)


def _box_intersects_ellipse(box: tuple[int, int, int, int], center: tuple[int, int], rx: int, ry: int, pad: int) -> bool:
    x1, y1, x2, y2 = box
    cx, cy = center
    rx = max(1, rx + pad)
    ry = max(1, ry + pad)
    closest_x = min(max(cx, x1), x2)
    closest_y = min(max(cy, y1), y2)
    return ((closest_x - cx) / rx) ** 2 + ((closest_y - cy) / ry) ** 2 < 1


def _adjust_column(items: list[dict[str, Any]], *, min_y: int, max_y: int, gap: int) -> None:
    if not items:
        return
    items.sort(key=lambda item: item["natural_y"])
    for item in items:
        item["top"] = max(min_y, min(item["natural_y"] - item["height"] // 2, max_y - item["height"]))
    for idx in range(1, len(items)):
        item = items[idx]
        prev = items[idx - 1]
        item["top"] = max(item["top"], prev["top"] + prev["height"] + gap)
    overflow = items[-1]["top"] + items[-1]["height"] - max_y
    if overflow > 0:
        items[-1]["top"] -= overflow
        for idx in range(len(items) - 2, -1, -1):
            item = items[idx]
            next_item = items[idx + 1]
            item["top"] = min(item["top"], next_item["top"] - item["height"] - gap)
        underflow = min_y - items[0]["top"]
        if underflow > 0:
            for item in items:
                item["top"] += underflow


def _estimate_text_width(lines: list[str], font_size_pt: float) -> int:
    max_chars = max((len(line) for line in lines), default=1)
    return int(max_chars * font_size_pt * 0.55 * EMU_PER_PT)


def _layout(
    items: list[dict[str, Any]],
    *,
    box: tuple[int, int, int, int],
    center_rel: tuple[float, float],
    radius_rel: tuple[float, float],
    start_angle: float,
    font_size_pt: float,
    side_padding_pt: float,
    min_gap_pt: float,
    label_y_bounds: tuple[float, float],
) -> dict[str, Any]:
    left, top, width, height = box
    cx = left + int(width * center_rel[0])
    cy = top + int(height * center_rel[1])
    rx = int(width * radius_rel[0])
    ry = int(height * radius_rel[1])
    total = sum(max(float(item.get("value", 0) or 0), 0) for item in items) or 1
    marker = int(max(5.0, font_size_pt * 0.62) * EMU_PER_PT)
    marker_gap = int(3.2 * EMU_PER_PT)
    side_padding = int(side_padding_pt * EMU_PER_PT)
    chart_pad = int(5.5 * EMU_PER_PT)
    min_y = top + int(height * label_y_bounds[0])
    max_y = top + int(height * label_y_bounds[1])
    line_h = int(font_size_pt * 1.12 * EMU_PER_PT)
    gap = int(min_gap_pt * EMU_PER_PT)

    angle = start_angle
    label_items: list[dict[str, Any]] = []
    for idx, item in enumerate(items):
        value = max(float(item.get("value", 0) or 0), 0)
        extent = 360 * value / total
        mid_angle = angle + extent / 2
        theta = math.radians(mid_angle)
        lines = [
            str(item.get("label") or "").strip(),
            _fmt_int(value),
            _fmt_pct(float(item.get("share", value / total) or 0)),
        ]
        text_w = _estimate_text_width(lines, font_size_pt)
        height_box = line_h * len(lines)
        block_w = marker + marker_gap + text_w
        side = "right" if math.cos(theta) >= 0 else "left"
        label_items.append({
            "index": idx,
            "source": item,
            "theta": theta,
            "side": side,
            "lines": lines,
            "width": block_w,
            "text_width": text_w,
            "height": height_box,
            "natural_y": int(cy + ry * 1.16 * math.sin(theta)),
            "font_size_pt": font_size_pt,
        })
        angle += extent

    _adjust_column([item for item in label_items if item["side"] == "left"], min_y=min_y, max_y=max_y, gap=gap)
    _adjust_column([item for item in label_items if item["side"] == "right"], min_y=min_y, max_y=max_y, gap=gap)

    boxes = []
    text_boxes = []
    clipped_text_count = 0
    for item in label_items:
        if item["side"] == "left":
            block_x = left + side_padding
            available = cx - rx - chart_pad - block_x
        else:
            block_x = left + width - side_padding - item["width"]
            available = left + width - side_padding - (cx + rx + chart_pad)
        if item["width"] > available:
            clipped_text_count += 1
        marker_x = block_x
        text_x = marker_x + marker + marker_gap
        item["marker_x"] = marker_x
        item["text_x"] = text_x
        item["marker_size"] = marker
        item["line_h"] = line_h
        item["connector_start"] = (
            int(cx + rx * 1.03 * math.cos(item["theta"])),
            int(cy + ry * 1.03 * math.sin(item["theta"])),
        )
        item["connector_end"] = (
            int(marker_x + marker / 2),
            int(item["top"] + marker / 2),
        )
        box_item = (int(block_x), int(item["top"]), int(block_x + item["width"]), int(item["top"] + item["height"]))
        text_box = (int(text_x), int(item["top"]), int(text_x + item["text_width"]), int(item["top"] + item["height"]))
        boxes.append(box_item)
        text_boxes.append(text_box)
        item["box"] = box_item
        item["text_box"] = text_box

    bounds = (left, top, left + width, top + height)
    diagnostics = {
        "label_count": len(label_items),
        "sector_count": len(items),
        "overlap_count": _count_overlaps(boxes),
        "out_of_bounds_count": _count_out_of_bounds(boxes, bounds),
        "label_chart_overlap_count": sum(1 for box_item in text_boxes if _box_intersects_ellipse(box_item, (cx, cy), rx, ry, chart_pad)),
        "color_mismatch_count": 0,
        "clipped_text_count": clipped_text_count,
        "missing_nonzero_label_count": sum(1 for item in items if float(item.get("value", 0) or 0) > 0) - len(label_items),
        "font_size_pt": font_size_pt,
        "label_boxes": [list(item["box"]) for item in label_items],
    }
    return {"items": label_items, "diagnostics": diagnostics}


def add_native_pie_labels(
    slide,
    *,
    items: list[dict[str, Any]],
    box: tuple[int, int, int, int],
    center_rel: tuple[float, float],
    radius_rel: tuple[float, float],
    start_angle: float = -90.0,
    label_y_bounds: tuple[float, float] = (0.08, 0.92),
    side_padding_pt: float = 14.0,
    min_gap_pt: float = 4.5,
) -> dict[str, Any]:
    if not items:
        items = [{"label": "нет данных", "value": 1, "share": 1, "color": "#D9D9D9"}]

    chosen = None
    font_sizes = [CHART_LABEL_SIZE, CHART_COMPACT_SIZE, CHART_COMPACT_SIZE - 1, CHART_COMPACT_SIZE - 2]
    for font_size_pt in dict.fromkeys(max(7.0, float(size)) for size in font_sizes):
        candidate = _layout(
            items,
            box=box,
            center_rel=center_rel,
            radius_rel=radius_rel,
            start_angle=start_angle,
            font_size_pt=font_size_pt,
            side_padding_pt=side_padding_pt,
            min_gap_pt=min_gap_pt,
            label_y_bounds=label_y_bounds,
        )
        diag = candidate["diagnostics"]
        chosen = candidate
        if not any(diag.get(key, 0) for key in ["overlap_count", "out_of_bounds_count", "label_chart_overlap_count", "clipped_text_count"]):
            break

    assert chosen is not None
    shapes_added = 0
    connector_color = RGBColor(191, 191, 191)
    for item in chosen["items"]:
        color = str(item["source"].get("color") or "#BFBFBF")
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            item["connector_start"][0],
            item["connector_start"][1],
            item["connector_end"][0],
            item["connector_end"][1],
        )
        line.line.color.rgb = connector_color
        line.line.width = Pt(0.5)
        shapes_added += 1

        square = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            item["marker_x"],
            item["top"] + int((item["line_h"] - item["marker_size"]) / 2),
            item["marker_size"],
            item["marker_size"],
        )
        square.fill.solid()
        square.fill.fore_color.rgb = _rgb(color)
        square.line.color.rgb = _rgb(color)
        shapes_added += 1

        text_box = slide.shapes.add_textbox(
            item["text_x"],
            item["top"],
            item["text_width"] + int(2 * EMU_PER_PT),
            item["height"] + int(1 * EMU_PER_PT),
        )
        tf = text_box.text_frame
        tf.clear()
        tf.word_wrap = False
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        for idx, text in enumerate(item["lines"]):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 1.0
            run = paragraph.add_run()
            run.text = text
            run.font.name = CHART_FONT
            run.font.size = Pt(item["font_size_pt"])
            run.font.bold = bool(idx == 0 and item["source"].get("highlight"))
            run.font.color.rgb = RGBColor(0, 0, 0)
        shapes_added += 1

    diagnostics = dict(chosen["diagnostics"])
    diagnostics["shapes_added_count"] = shapes_added
    return diagnostics
