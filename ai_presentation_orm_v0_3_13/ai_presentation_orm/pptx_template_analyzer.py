from __future__ import annotations

from pathlib import Path
from collections import Counter
import json
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from .utils import safe_text, emu_to_cm


def classify_slide(text: str) -> str:
    t = text.lower()
    if "информационная кампания" in t and len(t) < 400:
        return "cover"
    if "текущий отчет" in t or "текущий отчёт" in t or "период:" in t:
        return "current_report_plan_fact"
    if "упоминания брендов" in t or ("sov" in t and "бренд" in t):
        return "brand_mentions_sov"
    if "инфополе" in t or "контексты обсуждений" in t:
        return "problem_field"
    if "тональность" in t:
        return "tonality"
    if "типы площадок" in t or "топ источников" in t or "где обсуждают" in t:
        return "sources"
    if "показатели посева" in t or "посев" in t:
        return "seeding_results"
    if "примеры сообщений" in t:
        return "message_examples"
    if "рейтинги" in t or "карточ" in t or "отзыв" in t:
        return "ratings_reviews_cards"
    if "выводы" in t or "рекомендации" in t:
        return "conclusions"
    return "unknown"


def analyze_pptx_template(template_path: Path) -> dict:
    if not template_path.exists():
        raise FileNotFoundError(f"Template PPTX not found: {template_path}")

    prs = Presentation(str(template_path))
    result = {
        "file": str(template_path),
        "name": template_path.name,
        "slide_width_cm": emu_to_cm(prs.slide_width),
        "slide_height_cm": emu_to_cm(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": []
    }

    total_charts = 0
    total_tables = 0
    total_pictures = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        shapes = []
        for shape in slide.shapes:
            entry = {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "type": str(shape.shape_type),
                "x_cm": emu_to_cm(shape.left),
                "y_cm": emu_to_cm(shape.top),
                "w_cm": emu_to_cm(shape.width),
                "h_cm": emu_to_cm(shape.height),
                "has_text": bool(getattr(shape, "has_text_frame", False)),
                "has_table": bool(getattr(shape, "has_table", False)),
                "has_chart": bool(getattr(shape, "has_chart", False)),
            }
            if getattr(shape, "has_text_frame", False):
                text = safe_text(shape.text, 500)
                if text:
                    slide_texts.append(text)
                    entry["text_preview"] = text

            if getattr(shape, "has_chart", False):
                total_charts += 1
                try:
                    entry["chart_type"] = str(shape.chart.chart_type)
                    entry["chart_series_count"] = len(shape.chart.series)
                    entry["chart_has_legend"] = bool(shape.chart.has_legend)
                except Exception as exc:
                    entry["chart_error"] = str(exc)

            if getattr(shape, "has_table", False):
                total_tables += 1
                try:
                    entry["table_rows"] = len(shape.table.rows)
                    entry["table_cols"] = len(shape.table.columns)
                except Exception as exc:
                    entry["table_error"] = str(exc)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_pictures += 1
                entry["is_picture"] = True

            shapes.append(entry)

        full_text = "\n".join(slide_texts)
        result["slides"].append({
            "slide_number": slide_number,
            "role_guess": classify_slide(full_text),
            "shape_count": len(slide.shapes),
            "text_shape_count": sum(1 for s in shapes if s.get("text_preview")),
            "chart_count": sum(1 for s in shapes if s.get("has_chart")),
            "table_count": sum(1 for s in shapes if s.get("has_table")),
            "picture_count": sum(1 for s in shapes if s.get("is_picture")),
            "text_preview": safe_text(full_text, 1200),
            "shapes": shapes,
        })

    with zipfile.ZipFile(template_path, "r") as z:
        names = z.namelist()
        result["summary"] = {
            "total_charts": total_charts,
            "total_tables": total_tables,
            "total_pictures": total_pictures,
            "roles_count": dict(Counter(slide["role_guess"] for slide in result["slides"])),
            "chart_xml_files": len([n for n in names if n.startswith("ppt/charts/chart") and n.endswith(".xml")]),
            "embedded_xlsx_files": len([n for n in names if n.startswith("ppt/embeddings/") and n.endswith(".xlsx")]),
            "media_files": len([n for n in names if n.startswith("ppt/media/")]),
        }

    return result


def extract_pptx_as_json(pptx_path: Path) -> dict:
    """Extract a PPTX file to a JSON structure matching the reference format.

    Produces: {file, slide_count, slides: [{number, layout, shapes: [{name, shape_id,
    left_in, top_in, width_in, height_in, type, text, paragraphs}]}], total_shapes}
    """
    if not pptx_path.exists():
        return {"file": pptx_path.name, "slide_count": 0, "slides": [], "total_shapes": 0}

    prs = Presentation(str(pptx_path))
    slides_out = []
    total_shapes = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            layout_name = ""

        shapes_out = []
        for shape in slide.shapes:
            total_shapes += 1
            entry = {
                "name": shape.name,
                "shape_id": shape.shape_id,
                "left_in": round(shape.left / 914400, 3),
                "top_in": round(shape.top / 914400, 3),
                "width_in": round(shape.width / 914400, 3),
                "height_in": round(shape.height / 914400, 3),
            }
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                entry["type"] = "text"
                entry["text"] = shape.text
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    para_text = para.text
                    if para_text:
                        paragraphs.append({"text": para_text, "level": para.level})
                if not paragraphs:
                    paragraphs.append({"text": "", "level": 0})
                entry["paragraphs"] = paragraphs
            elif hasattr(shape, "has_table") and shape.has_table:
                entry["type"] = "table"
            elif hasattr(shape, "has_chart") and shape.has_chart:
                entry["type"] = "chart"
                try:
                    chart_obj = shape.chart
                    entry["chart_type"] = str(chart_obj.chart_type)
                    categories = []
                    try:
                        cat_data = chart_obj.plots[0].categories
                        if cat_data is not None:
                            categories = [str(c) for c in cat_data]
                    except Exception:
                        pass
                    entry["chart_categories"] = categories
                    series_list = []
                    try:
                        for plot in chart_obj.plots:
                            for series in plot.series:
                                nm = str(series.name) if series.name else ""
                                vals = [float(v) if v is not None else None for v in series.values]
                                series_list.append({"name": nm, "values": vals})
                    except Exception:
                        pass
                    entry["chart_series"] = series_list
                except Exception:
                    pass
            else:
                entry["type"] = "shape"
            shapes_out.append(entry)

        slides_out.append({
            "number": slide_number,
            "layout": layout_name,
            "shapes": shapes_out,
        })

    return {
        "file": pptx_path.name,
        "slide_count": len(prs.slides),
        "slides": slides_out,
        "total_shapes": total_shapes,
    }


def save_pptx_extraction(pptx_path: Path, output_path: Path) -> dict:
    """Extract PPTX and save as JSON. Returns the extraction dict."""
    extraction = extract_pptx_as_json(pptx_path)
    output_path.write_text(json.dumps(extraction, ensure_ascii=False, indent=2), encoding="utf-8")
    return extraction
