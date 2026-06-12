from __future__ import annotations

from collections.abc import Iterable

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.oxml.ns import qn
from pptx.util import Pt

from .style_rules import CHART_FONT_SIZE, FONT_FAMILY, format_int_spaces


def rgb(value: str | RGBColor) -> RGBColor:
    if isinstance(value, RGBColor):
        return value
    text = str(value or "").strip().lstrip("#")
    if len(text) != 6:
        return RGBColor(128, 128, 128)
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def remove_shape(slide, shape) -> None:
    if shape is not None:
        shape._element.getparent().remove(shape._element)


def _safe_set(fn) -> None:
    try:
        fn()
    except Exception:
        pass


def _font(font, *, size: float = CHART_FONT_SIZE, bold: bool = False, color: str | RGBColor = "#000000") -> None:
    _safe_set(lambda: setattr(font, "name", FONT_FAMILY))
    _safe_set(lambda: setattr(font, "size", Pt(size)))
    _safe_set(lambda: setattr(font, "bold", bold))
    _safe_set(lambda: setattr(font.color, "rgb", rgb(color)))


def _set_axis_font(axis, *, size: float = CHART_FONT_SIZE) -> None:
    _safe_set(lambda: setattr(axis, "major_tick_mark", XL_TICK_MARK.NONE))
    _safe_set(lambda: setattr(axis, "minor_tick_mark", XL_TICK_MARK.NONE))
    _safe_set(lambda: _font(axis.tick_labels.font, size=size, color="#333333"))


def _hide_value_axis(chart) -> None:
    axis = getattr(chart, "value_axis", None)
    if axis is None:
        return
    _safe_set(lambda: setattr(axis, "visible", False))
    _safe_set(lambda: setattr(axis, "has_major_gridlines", False))
    _safe_set(lambda: _font(axis.tick_labels.font, size=CHART_FONT_SIZE, color="#333333"))


def _style_category_axis(chart, *, reverse: bool = False) -> None:
    axis = getattr(chart, "category_axis", None)
    if axis is None:
        return
    _safe_set(lambda: setattr(axis, "reverse_order", reverse))
    _safe_set(lambda: setattr(axis, "has_major_gridlines", False))
    _set_axis_font(axis)


def _style_legend(chart, *, show: bool, position=XL_LEGEND_POSITION.BOTTOM) -> None:
    chart.has_legend = show
    if not show:
        return
    _safe_set(lambda: setattr(chart.legend, "position", position))
    _safe_set(lambda: setattr(chart.legend, "include_in_layout", True))
    _safe_set(lambda: _font(chart.legend.font, size=CHART_FONT_SIZE, color="#333333"))


def _set_plot_labels(chart, *, position=XL_LABEL_POSITION.OUTSIDE_END, number_format: str = "# ##0") -> None:
    for plot in chart.plots:
        _safe_set(lambda plot=plot: setattr(plot, "has_data_labels", True))
        labels = getattr(plot, "data_labels", None)
        if labels is None:
            continue
        _safe_set(lambda labels=labels: setattr(labels, "show_value", True))
        _safe_set(lambda labels=labels: setattr(labels, "number_format", number_format))
        _safe_set(lambda labels=labels: setattr(labels, "position", position))
        _safe_set(lambda labels=labels: _font(labels.font, size=CHART_FONT_SIZE, color="#333333"))


def _hide_plot_labels(chart) -> None:
    for plot in chart.plots:
        _safe_set(lambda plot=plot: setattr(plot, "has_data_labels", False))


def _set_pie_labels(chart, *, position=XL_LABEL_POSITION.BEST_FIT, number_format: str = "# ##0", show_legend_key: bool = False) -> None:
    for plot in chart.plots:
        _safe_set(lambda plot=plot: setattr(plot, "has_data_labels", True))
        labels = getattr(plot, "data_labels", None)
        if labels is None:
            continue
        _safe_set(lambda labels=labels: setattr(labels, "show_category_name", True))
        _safe_set(lambda labels=labels: setattr(labels, "show_value", True))
        _safe_set(lambda labels=labels: setattr(labels, "show_percentage", True))
        _safe_set(lambda labels=labels: setattr(labels, "show_series_name", False))
        _safe_set(lambda labels=labels: setattr(labels, "show_legend_key", show_legend_key))
        _safe_set(lambda labels=labels: setattr(labels, "number_format", number_format))
        _safe_set(lambda labels=labels: setattr(labels, "position", position))
        _safe_set(lambda labels=labels: _font(labels.font, size=CHART_FONT_SIZE, color="#333333"))


def _apply_series_colors(chart, series_colors: Iterable[str | RGBColor] | None = None, point_colors: Iterable[str | RGBColor] | None = None) -> None:
    if series_colors:
        for series, color in zip(chart.series, series_colors):
            _safe_set(lambda series=series, color=color: series.format.fill.solid())
            _safe_set(lambda series=series, color=color: setattr(series.format.fill.fore_color, "rgb", rgb(color)))
            _safe_set(lambda series=series, color=color: setattr(series.format.line.color, "rgb", rgb(color)))
    if point_colors and len(chart.series):
        points = chart.series[0].points
        for idx, color in enumerate(point_colors):
            if idx >= len(points):
                break
            point = points[idx]
            _safe_set(lambda point=point, color=color: point.format.fill.solid())
            _safe_set(lambda point=point, color=color: setattr(point.format.fill.fore_color, "rgb", rgb(color)))
            _safe_set(lambda point=point, color=color: setattr(point.format.line.color, "rgb", rgb(color)))


def _set_bar_grouping(chart, grouping: str | None) -> None:
    if not grouping:
        return
    try:
        for node in chart._chartSpace.xpath(".//c:barGrouping"):
            node.set("val", grouping)
        for node in chart._chartSpace.xpath(".//c:grouping"):
            node.set("val", grouping)
    except Exception:
        try:
            for node in chart._chartSpace.findall(".//" + qn("c:barGrouping")):
                node.set("val", grouping)
            for node in chart._chartSpace.findall(".//" + qn("c:grouping")):
                node.set("val", grouping)
        except Exception:
            pass


def replace_chart_data(
    chart_shape,
    *,
    categories: list[str],
    series: list[tuple[str, list[int | float], str | RGBColor]],
    show_legend: bool = True,
    legend_position=XL_LEGEND_POSITION.BOTTOM,
    show_data_labels: bool = True,
    label_position=XL_LABEL_POSITION.OUTSIDE_END,
    number_format: str = "# ##0",
    reverse_categories: bool = False,
    hide_value_axis: bool = True,
    bar_grouping: str | None = None,
):
    if not categories or not series:
        return chart_shape
    data = CategoryChartData()
    data.categories = categories
    for name, values, _ in series:
        data.add_series(name, values)
    chart = chart_shape.chart
    chart.replace_data(data)
    _set_bar_grouping(chart, bar_grouping)
    chart.has_title = False
    _style_legend(chart, show=show_legend, position=legend_position)
    _style_category_axis(chart, reverse=reverse_categories)
    if hide_value_axis:
        _hide_value_axis(chart)
    else:
        value_axis = getattr(chart, "value_axis", None)
        if value_axis is not None:
            _safe_set(lambda: setattr(value_axis, "has_major_gridlines", False))
            _set_axis_font(value_axis)
    if show_data_labels:
        _set_plot_labels(chart, position=label_position, number_format=number_format)
    else:
        _hide_plot_labels(chart)
    _apply_series_colors(chart, series_colors=[color for _, _, color in series])
    return chart_shape


def replace_pie_chart_data(
    chart_shape,
    *,
    labels: list[str],
    values: list[int | float],
    colors: list[str | RGBColor],
    show_legend: bool = True,
    legend_position=XL_LEGEND_POSITION.RIGHT,
    label_position=XL_LABEL_POSITION.BEST_FIT,
    number_format: str = "# ##0",
    doughnut: bool = False,
    show_legend_key: bool = False,
    show_data_labels: bool = True,
):
    if not labels or not values:
        return chart_shape
    data = CategoryChartData()
    data.categories = labels
    data.add_series("SOV", values)
    chart = chart_shape.chart
    chart.replace_data(data)
    chart.has_title = False
    _style_legend(chart, show=show_legend, position=legend_position)
    _apply_series_colors(chart, point_colors=colors)
    _safe_set(lambda: setattr(chart.plots[0], "vary_by_categories", True))
    if show_data_labels:
        _set_pie_labels(chart, position=label_position, number_format=number_format, show_legend_key=show_legend_key)
    else:
        _hide_plot_labels(chart)
    if doughnut:
        _safe_set(lambda: setattr(chart.plots[0], "hole_size", 55))
    return chart_shape


def add_native_bar_chart(
    slide,
    *,
    box: tuple[int, int, int, int],
    categories: list[str],
    values: list[int | float],
    color: str | RGBColor = "#808080",
    point_colors: list[str | RGBColor] | None = None,
    name: str = "Значение",
    reverse: bool = True,
):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(name, values)
    shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, *box, data)
    chart = shape.chart
    chart.has_title = False
    _style_legend(chart, show=False)
    _style_category_axis(chart, reverse=reverse)
    _hide_value_axis(chart)
    _set_plot_labels(chart, position=XL_LABEL_POSITION.OUTSIDE_END, number_format="# ##0")
    _apply_series_colors(chart, series_colors=[color])
    if point_colors:
        _apply_series_colors(chart, point_colors=point_colors)
    return shape


def add_native_stacked_bar_chart(
    slide,
    *,
    box: tuple[int, int, int, int],
    categories: list[str],
    series: list[tuple[str, list[int | float], str | RGBColor]],
    reverse: bool = True,
    show_legend: bool = True,
):
    data = CategoryChartData()
    data.categories = categories
    for name, values, _ in series:
        data.add_series(name, values)
    shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, *box, data)
    chart = shape.chart
    chart.has_title = False
    _style_legend(chart, show=show_legend)
    _style_category_axis(chart, reverse=reverse)
    _hide_value_axis(chart)
    _set_plot_labels(chart, position=XL_LABEL_POSITION.CENTER, number_format="# ##0")
    _apply_series_colors(chart, series_colors=[color for _, _, color in series])
    return shape


def add_native_column_chart(
    slide,
    *,
    box: tuple[int, int, int, int],
    categories: list[str],
    series: list[tuple[str, list[int | float], str | RGBColor]],
    stacked: bool = False,
    stacked_100: bool = False,
    show_legend: bool = True,
    label_position=XL_LABEL_POSITION.OUTSIDE_END,
    number_format: str = "# ##0",
):
    data = CategoryChartData()
    data.categories = categories
    for name, values, _ in series:
        data.add_series(name, values)
    chart_type = XL_CHART_TYPE.COLUMN_STACKED_100 if stacked_100 else (XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED)
    shape = slide.shapes.add_chart(chart_type, *box, data)
    chart = shape.chart
    chart.has_title = False
    _style_legend(chart, show=show_legend)
    _style_category_axis(chart)
    _hide_value_axis(chart)
    _set_plot_labels(chart, position=label_position, number_format=number_format)
    _apply_series_colors(chart, series_colors=[color for _, _, color in series])
    return shape


def add_native_line_chart(
    slide,
    *,
    box: tuple[int, int, int, int],
    categories: list[str],
    series: list[tuple[str, list[int | float], str | RGBColor]],
    show_legend: bool = True,
):
    data = CategoryChartData()
    data.categories = categories
    for name, values, _ in series:
        data.add_series(name, values)
    shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE, *box, data)
    chart = shape.chart
    chart.has_title = False
    _style_legend(chart, show=show_legend)
    _style_category_axis(chart)
    value_axis = getattr(chart, "value_axis", None)
    if value_axis is not None:
        _safe_set(lambda: setattr(value_axis, "has_major_gridlines", False))
        _set_axis_font(value_axis)
    for chart_series, (_, _, color) in zip(chart.series, series):
        _safe_set(lambda chart_series=chart_series, color=color: setattr(chart_series.format.line.color, "rgb", rgb(color)))
        _safe_set(lambda chart_series=chart_series: setattr(chart_series.format.line, "width", Pt(2.0)))
    return shape


def add_native_pie_chart(
    slide,
    *,
    box: tuple[int, int, int, int],
    labels: list[str],
    values: list[int | float],
    colors: list[str | RGBColor],
    doughnut: bool = False,
):
    data = CategoryChartData()
    data.categories = labels
    data.add_series("SOV", values)
    chart_type = XL_CHART_TYPE.DOUGHNUT if doughnut else XL_CHART_TYPE.PIE
    shape = slide.shapes.add_chart(chart_type, *box, data)
    chart = shape.chart
    chart.has_title = False
    _style_legend(chart, show=False)
    _apply_series_colors(chart, point_colors=colors)
    _safe_set(lambda: setattr(chart.plots[0], "vary_by_categories", True))
    if doughnut:
        _safe_set(lambda: setattr(chart.plots[0], "hole_size", 55))
    return shape


def chart_update(role: str, chart_type: str, **extra) -> dict:
    payload = {"role": role, "type": chart_type, "native_chart": True}
    payload.update(extra)
    return payload


def value_label_text(value: int | float | None) -> str:
    return format_int_spaces(value)
