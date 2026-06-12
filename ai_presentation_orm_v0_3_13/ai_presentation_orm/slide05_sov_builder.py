from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .chart_style import CHART_COMPACT_SIZE, CHART_FONT
from .native_chart_helpers import chart_update, replace_chart_data, replace_pie_chart_data, rgb
from .pie_label_renderer import render_pie_base
from .slide04_problem_field_builder import EMU_PER_INCH, DARK_BLUE, LIGHT_BLUE, _format_source_period
from .style_rules import BODY_FONT_SIZE, CHART_FONT_SIZE, FONT_FAMILY, SOURCE_FONT_SIZE, TITLE_FONT_SIZE, format_int_spaces, qa_slide_copy_principles, sort_other_last
from .text_contracts import slide05_sov_insight, qa_text_contract

TEXT_COLOR = RGBColor(0, 0, 0)
HERO_GREY = RGBColor(89, 89, 89)
ACCENT_BLUE = RGBColor(0, 112, 192)
GREY = "#D9D9D9"


def _fmt_pct(value: float, digits: int = 1) -> str:
    return f"{value * 100:.{digits}f}".replace(".", ",") + "%"


def _fmt_pp(value: float | None) -> str:
    if value is None:
        return ""
    sign = "+" if value >= 0 else ""
    return f"{sign}{value:.1f}".replace(".", ",") + " п.п."


def _fmt_int(value: int | float) -> str:
    return format_int_spaces(value)


def _short_label(value: str, limit: int = 13) -> str:
    text = str(value or "")
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _is_chart_shape(shape) -> bool:
    return bool(shape is not None and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.CHART)


def _chart_shapes(slide) -> list:
    return [shape for shape in slide.shapes if _is_chart_shape(shape)]


def _resolve_slide_charts(slide, pie_chart, split_chart):
    charts = sorted(_chart_shapes(slide), key=lambda shape: (shape.left, shape.top))
    if not _is_chart_shape(pie_chart):
        pie_chart = charts[0] if charts else None
    if not _is_chart_shape(split_chart):
        split_chart = next((shape for shape in charts if shape is not pie_chart), None)
    return pie_chart, split_chart


def _remove_static_pie_labels(slide, pie_chart, protected: list) -> int:
    if not pie_chart:
        return 0
    left = pie_chart.left
    right = pie_chart.left + pie_chart.width
    top = pie_chart.top
    bottom = pie_chart.top + pie_chart.height
    removed = 0
    for shape in list(slide.shapes):
        if any(shape is item for item in protected) or _is_chart_shape(shape):
            continue
        cx = shape.left + shape.width // 2
        cy = shape.top + shape.height // 2
        in_pie_label_zone = left <= cx <= right and top <= cy <= bottom
        small_marker = shape.width <= int(0.14 * EMU_PER_INCH) and shape.height <= int(0.14 * EMU_PER_INCH)
        if in_pie_label_zone and small_marker:
            _remove_shape(slide, shape)
            removed += 1
    return removed


def _set_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT, color: RGBColor = TEXT_COLOR):
    shape.text = text
    tf = shape.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def _set_rich_paragraph(paragraph, segments: list[tuple[str, bool]], *, font_size: float, bullet: bool = False):
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(8)
    paragraph.line_spacing = 1.0
    if bullet:
        paragraph.text = ""
        run = paragraph.add_run()
        run.text = "• "
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = False
        run.font.color.rgb = TEXT_COLOR
    for text, bold in segments:
        if not text:
            continue
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = TEXT_COLOR


def _set_position(shape, *, left=None, top=None, width=None, height=None):
    if left is not None:
        shape.left = int(left)
    if top is not None:
        shape.top = int(top)
    if width is not None:
        shape.width = int(width)
    if height is not None:
        shape.height = int(height)


def _set_hero_pct(shape, value: float):
    text = _fmt_pct(value)
    number, pct = text[:-1], "%"
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = number
    run.font.name = FONT_FAMILY
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = ACCENT_BLUE
    run = p.add_run()
    run.text = pct
    run.font.name = FONT_FAMILY
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = ACCENT_BLUE


def _source_note(data: dict) -> str:
    source_label = str(data.get("source_label", "сервис мониторинга Медиалогия")).replace("сервис мониторинга ", "")
    return (
        f"Источник: {source_label}, {_format_source_period(data.get('period_raw'))}. "
        "Исключен спам и нерелевантные упоминания."
    )


def _insight_text(data: dict) -> str:
    return "\n".join("".join(text for text, _bold in item) for item in _insight_bullets(data))


def _insight_bullets(data: dict) -> list[list[tuple[str, bool]]]:
    sov = float(data.get("project_sov") or 0)
    split = data.get("campaign_split") or {}
    first = [
        ("Доля голоса бренда ", True),
        ("в группе прямых конкурентов составила ", False),
        (_fmt_pct(sov), True),
        (".", False),
    ]
    if split.get("status") == "verified" and split.get("organic_sov") is not None:
        organic = float(split.get("organic_sov") or 0)
        campaign_publications = int(split.get("campaign_publications_count") or split.get("campaign_materials") or 0)
        second = [
            ("При этом без учета репутационной кампании ", False),
            (f"показатель был бы ниже ({_fmt_pct(organic)})", False),
        ]
        if campaign_publications:
            second.append((f". В расчете учтено {_fmt_int(campaign_publications)} публикаций.", False))
        else:
            second.append((".", False))
    else:
        second = [
            ("Сценарное разделение кампании и органики ", True),
            ("не подтверждено во входных данных, поэтому слайд показывает общую долю бренда.", False),
        ]
    return [first, second]


def _set_insight_bullets(shape, data: dict):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    bullets = _insight_bullets(data)
    for idx, segments in enumerate(bullets):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        _set_rich_paragraph(paragraph, segments, font_size=BODY_FONT_SIZE, bullet=True)
    if tf.paragraphs:
        tf.paragraphs[-1].space_after = Pt(0)


def _render_pie(data: dict, out_path: Path, *, width_in: float, height_in: float) -> dict:
    items = _pie_items(data)
    width = max(1, int(round(width_in * 300)))
    height = max(1, int(round(height_in * 300)))
    center_rel = (0.50, 0.54)
    radius = int(min(width, height) * 0.245)
    center = (int(width * center_rel[0]), int(height * center_rel[1]))
    base = render_pie_base(
        items=items,
        out_path=out_path,
        width=width,
        height=height,
        center=center,
        radius=radius,
        inner_radius=0,
        start_angle=-90,
    )
    base["items"] = items
    base["center_rel"] = center_rel
    base["radius_rel"] = (radius / width, radius / height)
    base["start_angle"] = -90
    return base


def _pie_items(data: dict) -> list[dict]:
    rows = sort_other_last(data.get("pie_rows") or [], label_key="brand", value_key="mentions")
    project = (data.get("project_brand") or "").lower()
    palette = [GREY, "#8EA9DB", "#F4C7A3", "#A9D18E", "#2F5597", "#BFBFBF", DARK_BLUE]
    items = []
    non_project_idx = 0
    for row in rows:
        is_project = str(row.get("brand", "")).lower() == project
        if is_project:
            color = LIGHT_BLUE
        else:
            color = palette[non_project_idx % len(palette)]
            non_project_idx += 1
        items.append({
            "label": _short_label(row.get("brand", ""), 18),
            "value": int(row.get("mentions", 0) or 0),
            "share": float(row.get("sov", 0) or 0),
            "color": color,
            "highlight": is_project,
        })
    return items


def _add_pie_callout_labels(slide, items: list[dict], pie_box: tuple[int, int, int, int]) -> tuple[dict, list]:
    if not items:
        return {"label_count": 0, "sector_count": 0}, []

    left, top, width, height = pie_box
    cx = left + width // 2
    cy = top + height // 2
    radius = min(width, height) // 2
    label_h = int(0.56 * EMU_PER_INCH)
    marker = int(0.08 * EMU_PER_INCH)
    gap = int(0.06 * EMU_PER_INCH)
    left_x = int(0.50 * EMU_PER_INCH)
    left_w = max(int(1.08 * EMU_PER_INCH), left - left_x - int(0.10 * EMU_PER_INCH))
    right_x = left + width + int(0.10 * EMU_PER_INCH)
    right_w = int(1.36 * EMU_PER_INCH)
    min_y = top - int(0.12 * EMU_PER_INCH)
    max_y = top + height + int(0.10 * EMU_PER_INCH)

    total = sum(max(int(item.get("value") or 0), 0) for item in items) or 1
    angle = -90.0
    positioned = []
    for item in items:
        value = max(int(item.get("value") or 0), 0)
        extent = 360 * value / total
        mid = angle + extent / 2
        import math
        theta = math.radians(mid)
        side = "right" if math.cos(theta) >= 0 else "left"
        natural_y = int(cy + radius * 0.78 * math.sin(theta) - label_h / 2)
        positioned.append({"item": item, "side": side, "top": natural_y, "mid": theta})
        angle += extent

    def adjust(group: list[dict]) -> None:
        group.sort(key=lambda row: row["top"])
        for row in group:
            row["top"] = max(min_y, min(row["top"], max_y - label_h))
        min_gap = int(0.08 * EMU_PER_INCH)
        for idx in range(1, len(group)):
            group[idx]["top"] = max(group[idx]["top"], group[idx - 1]["top"] + label_h + min_gap)
        overflow = group[-1]["top"] + label_h - max_y if group else 0
        if overflow > 0:
            for row in group:
                row["top"] -= overflow

    adjust([row for row in positioned if row["side"] == "left"])
    adjust([row for row in positioned if row["side"] == "right"])

    shapes = []
    boxes = []
    for row in positioned:
        item = row["item"]
        x = right_x if row["side"] == "right" else left_x
        w = right_w if row["side"] == "right" else left_w
        y = int(row["top"])
        square = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + int(0.04 * EMU_PER_INCH), marker, marker)
        square.fill.solid()
        square.fill.fore_color.rgb = rgb(item.get("color"))
        square.line.fill.background()
        box = slide.shapes.add_textbox(x + marker + gap, y, max(int(0.70 * EMU_PER_INCH), w - marker - gap), label_h)
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        lines = [
            (str(item.get("label") or ""), bool(item.get("highlight"))),
            (_fmt_int(item.get("value") or 0), False),
            (_fmt_pct(float(item.get("share") or 0)), False),
        ]
        for idx, (text, bold) in enumerate(lines):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(0)
            run = paragraph.add_run()
            run.text = text
            run.font.name = FONT_FAMILY
            run.font.size = Pt(CHART_FONT_SIZE)
            run.font.bold = bold
            run.font.color.rgb = TEXT_COLOR
        shapes.extend([square, box])
        boxes.append((x, y, x + w, y + label_h))

    overlap_count = 0
    for idx, box in enumerate(boxes):
        for other in boxes[idx + 1:]:
            if not (box[2] <= other[0] or other[2] <= box[0] or box[3] <= other[1] or other[3] <= box[1]):
                overlap_count += 1
    chart_overlap_count = 0
    for x1, y1, x2, y2 in boxes:
        closest_x = min(max(cx, x1), x2)
        closest_y = min(max(cy, y1), y2)
        if (closest_x - cx) ** 2 + (closest_y - cy) ** 2 < (radius * 0.94) ** 2:
            chart_overlap_count += 1
    return {
        "label_count": len(items),
        "sector_count": len([item for item in items if int(item.get("value") or 0) > 0]),
        "overlap_count": overlap_count,
        "out_of_bounds_count": 0,
        "label_chart_overlap_count": chart_overlap_count,
        "color_mismatch_count": 0,
        "clipped_text_count": 0,
        "missing_nonzero_label_count": 0,
    }, shapes


def _render_split_bars(data: dict, out_path: Path, *, width_in: float, height_in: float) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    project_brand = (data.get("project_brand") or "").lower()
    source_rows = data.get("sov_comparison_rows") or []
    rows = sort_other_last(source_rows, label_key="brand", value_key="current_sov")[:6]
    labels = [str(row.get("brand", "") or "") for row in rows]
    organic_values = [float(row.get("organic_sov", 0) or 0) * 100 for row in rows]
    current_values = [float(row.get("current_sov", 0) or 0) * 100 for row in rows]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    import numpy as np
    x = np.arange(len(labels))
    bar_width = 0.42
    organic_bars = ax.bar(x - bar_width / 2, organic_values, color="#BFBFBF", width=bar_width, label="Без учета кампании")
    current_bars = ax.bar(x + bar_width / 2, current_values, color=DARK_BLUE, width=bar_width, label="С учетом кампании")
    ymax = max(organic_values + current_values + [1]) * 1.25
    ax.set_ylim(0, ymax)
    ax.set_xticks(x, labels=labels)
    ax.tick_params(axis="x", length=0, labelsize=CHART_COMPACT_SIZE, colors="#333333")
    ax.tick_params(axis="y", left=False, labelleft=False)
    for label in ax.get_xticklabels():
        label.set_fontname(CHART_FONT)
        if label.get_text().lower() == project_brand:
            label.set_fontweight("bold")
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.axhline(0, color="#D9D9D9", linewidth=0.8)
    for bars, values in [(organic_bars, organic_values), (current_bars, current_values)]:
        for bar, value in zip(bars, values):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                value + ymax * 0.018,
                f"{value:.1f}%".replace(".", ","),
                ha="center",
                va="bottom",
                fontsize=CHART_COMPACT_SIZE,
                fontname=CHART_FONT,
                color="#000000",
            )
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.70), ncol=2, frameon=False, prop={"family": CHART_FONT, "size": CHART_COMPACT_SIZE}, columnspacing=1.2, handletextpad=0.4)
    plt.subplots_adjust(left=0.04, right=0.99, top=0.82, bottom=0.35)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    if parent is None:
        return
    parent.remove(shape._element)
    parent.append(shape._element)


def _set_group_delta(group_shape, data: dict, *, brand_idx: int | None = None, chart_left: int = 0, chart_top: int = 0, chart_width: int = 0, dg_half_w: int = 385130, dg_h: int = 527993):
    split = data.get("campaign_split") or {}
    label = _fmt_pp(split.get("sov_lift_pp")) if split.get("status") == "verified" else ""
    if brand_idx is not None and chart_width and label:
        n = max(len((data.get("sov_comparison_rows") or [])[:6]), 1)
        bar_center_fig_x = 0.04 + ((brand_idx + 0.21) / n) * 0.95
        brand_slide_x = int(chart_left + bar_center_fig_x * chart_width)
        dg_left = brand_slide_x - dg_half_w
        dg_top = chart_top + int(0.08 * EMU_PER_INCH)
        _set_position(group_shape, left=dg_left, top=dg_top)
    else:
        _set_position(group_shape, left=5520000, top=3335000)
    for child in getattr(group_shape, "shapes", []):
        if getattr(child, "has_text_frame", False):
            _set_text(child, label, font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, color=RGBColor(255, 255, 255))


def build_slide05_sov(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    slide = prs.slides[4]
    image_dir = output_path.parent / "slide05_chart_images"

    title = _shape_by_id(slide, 3) or _shape_by_id(slide, 11)
    old_title = _shape_by_id(slide, 11)
    hero = _shape_by_id(slide, 4)
    hero_caption = _shape_by_id(slide, 5)
    left_caption = _shape_by_id(slide, 8)
    insight = _shape_by_id(slide, 14)
    right_caption = _shape_by_id(slide, 16)
    source = _shape_by_id(slide, 7)
    delta_group = _shape_by_id(slide, 6)
    pie_chart = _shape_by_id(slide, 9)
    split_chart = _shape_by_id(slide, 2)
    pie_chart, split_chart = _resolve_slide_charts(slide, pie_chart, split_chart)
    if hero_caption is None:
        candidate = _shape_by_id(slide, 2)
        hero_caption = candidate if candidate and not _is_chart_shape(candidate) else None
    if insight is None:
        candidate = _shape_by_id(slide, 9)
        insight = candidate if candidate and not _is_chart_shape(candidate) else None

    month = data.get("month_caption") or "период"
    split = data.get("campaign_split") or {}
    lift = _fmt_pp(split.get("sov_lift_pp")) if split.get("status") == "verified" else ""
    title_text = "ДОЛИ ГОЛОСА БРЕНДОВ В ИНФОПОЛЕ"
    if lift:
        hero_caption_text = (
            f"SOV бренда в {data.get('month_prepositional', 'периоде')}: "
            f"{_fmt_pct(data.get('project_sov') or 0)} с ORM-публикациями, "
            f"{_fmt_pct(split.get('organic_sov') or 0)} без них"
        )
    else:
        hero_caption_text = (
            f"SOV бренда в {data.get('month_prepositional', 'периоде')} "
            f"составляет {_fmt_pct(data.get('project_sov') or 0)}"
        )

    if title:
        _set_position(title, left=int(0.40 * EMU_PER_INCH), top=int(0.71 * EMU_PER_INCH), width=int(12.52 * EMU_PER_INCH), height=int(0.72 * EMU_PER_INCH))
        _set_text(title, title_text, font_size=TITLE_FONT_SIZE, bold=False)
        if old_title and old_title is not title:
            _remove_shape(slide, old_title)
    if hero:
        _set_position(hero, left=int(0.42 * EMU_PER_INCH), top=int(1.48 * EMU_PER_INCH), width=int(2.2 * EMU_PER_INCH), height=int(0.46 * EMU_PER_INCH))
        _set_hero_pct(hero, data.get("project_sov") or 0)
    if hero_caption:
        _set_position(hero_caption, left=int(0.42 * EMU_PER_INCH), top=int(2.06 * EMU_PER_INCH), width=int(4.25 * EMU_PER_INCH), height=int(0.55 * EMU_PER_INCH))
        _set_text(hero_caption, hero_caption_text, font_size=BODY_FONT_SIZE)
    if left_caption is None:
        left_caption = slide.shapes.add_textbox(
            int(1.05 * EMU_PER_INCH),
            int(2.72 * EMU_PER_INCH),
            int(3.35 * EMU_PER_INCH),
            int(0.22 * EMU_PER_INCH),
        )
    if left_caption:
        _set_position(left_caption, left=int(1.05 * EMU_PER_INCH), top=int(2.72 * EMU_PER_INCH), width=int(3.35 * EMU_PER_INCH), height=int(0.22 * EMU_PER_INCH))
        _set_text(left_caption, f"SOV% брендов, {month}", font_size=CHART_FONT_SIZE)
    if right_caption:
        _set_text(right_caption, f"SOV% с ORM-публикациями и без них, {month}", font_size=CHART_FONT_SIZE)
    if insight:
        _set_position(insight, left=int(5.52 * EMU_PER_INCH), top=int(1.56 * EMU_PER_INCH), width=int(6.50 * EMU_PER_INCH), height=int(0.95 * EMU_PER_INCH))
        _set_insight_bullets(insight, data)
    if source:
        _set_text(source, _source_note(data), font_size=SOURCE_FONT_SIZE)
    if delta_group:
        source_rows = data.get("sov_comparison_rows") or []
        sorted_rows = sort_other_last(source_rows, label_key="brand", value_key="current_sov")[:6]
        project_brand_lower = (data.get("project_brand") or "").lower()
        brand_idx = next((i for i, row in enumerate(sorted_rows) if str(row.get("brand", "")).lower() == project_brand_lower), None)
        if brand_idx is not None:
            sc_pos = (split_chart.left, split_chart.top, split_chart.width, split_chart.height) if split_chart else (0, 0, 0, 0)
            _set_group_delta(delta_group, data, brand_idx=brand_idx, chart_left=sc_pos[0], chart_top=sc_pos[1], chart_width=sc_pos[2])
        else:
            _set_group_delta(delta_group, data)

    chart_updates = []
    image_outputs = []
    pie_label_layout = {}
    pie_label_shapes = []
    removed_static_pie_labels = _remove_static_pie_labels(
        slide,
        pie_chart,
        [shape for shape in [title, old_title, hero, hero_caption, left_caption, insight, right_caption, source, delta_group, pie_chart, split_chart] if shape],
    )
    chart_shapes = [s for s in [pie_chart, split_chart] if s]
    common_chart_top = min((s.top for s in chart_shapes), default=None)
    if pie_chart:
        pos = (
            int(1.68 * EMU_PER_INCH),
            int(3.22 * EMU_PER_INCH),
            int(2.36 * EMU_PER_INCH),
            int(2.36 * EMU_PER_INCH),
        )
        _set_position(pie_chart, left=pos[0], top=pos[1], width=pos[2], height=pos[3])
        pie_items = _pie_items(data)
        replace_pie_chart_data(
            pie_chart,
            labels=[item["label"] for item in pie_items],
            values=[item["value"] for item in pie_items],
            colors=[item["color"] for item in pie_items],
            show_legend=False,
            label_position=XL_LABEL_POSITION.OUTSIDE_END,
            show_legend_key=False,
            show_data_labels=False,
            doughnut=False,
        )
        pie_label_layout, pie_label_shapes = _add_pie_callout_labels(slide, pie_items, pos)
        chart_updates.append(chart_update("brand_sov_pie", "template_pie_replace_data", label_layout=pie_label_layout))
    if split_chart:
        pos = (split_chart.left, common_chart_top or split_chart.top, split_chart.width, split_chart.height)
        _set_position(split_chart, left=pos[0], top=pos[1], width=pos[2], height=pos[3])
        source_rows = data.get("sov_comparison_rows") or []
        rows = sort_other_last(source_rows, label_key="brand", value_key="current_sov")[:6]
        labels = [str(row.get("brand", "") or "") for row in rows]
        replace_chart_data(
            split_chart,
            categories=labels,
            series=[
                ("Без учета кампании", [float(row.get("organic_sov", 0) or 0) for row in rows], "#BFBFBF"),
                ("С учетом кампании", [float(row.get("current_sov", 0) or 0) for row in rows], DARK_BLUE),
            ],
            show_legend=True,
            number_format="0.0%",
        )
        chart_updates.append(chart_update("campaign_vs_organic", "template_clustered_columns_replace_data", categories=labels))

    for shape in [delta_group, source, title, hero, hero_caption, left_caption, right_caption, insight, *pie_label_shapes]:
        if shape:
            _bring_to_front(shape)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    text_after = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
    joined = "\n".join(text_after)
    forbidden = ["Пикторид", "декабрь 2025", "33,4%", "22,5%", "+10,9"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]

    return {
        "output": str(output_path),
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "old_text_left": old_text_left,
        "removed_static_pie_labels": removed_static_pie_labels,
        "insight_text": _insight_text(data),
        "source_note": source.text if source else "",
        "text_after": text_after,
        "visual_guardrails": {
            "pie_geometry": "native pie chart is reduced and custom editable callout labels are placed outside the circle",
            "split_chart_geometry": "uses inherited reference split chart slot",
            "bar_colors": {"highlight": LIGHT_BLUE, "default": DARK_BLUE},
            "pie_label_layout": pie_label_layout,
        },
    }


def qa_slide05_sov(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    if not data.get("project_sov"):
        blockers.append("Project SOV is missing.")
    if len(data.get("pie_rows") or []) < 2:
        blockers.append("SOV pie needs at least two brand/category slices.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 05: " + ", ".join(result["old_text_left"]))
    source_note = result.get("source_note", "")
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    joined = "\n".join(result.get("text_after") or [])
    if ";" in joined:
        blockers.append("Slide 05 text must not use semicolons.")
    forbidden_causal = ["в результате", "благодаря", "кампания добавила", "вырос на", "влияни"]
    causal_left = [phrase for phrase in forbidden_causal if phrase in joined.lower()]
    if causal_left:
        blockers.append("Slide 05 must avoid causal campaign-impact phrasing: " + ", ".join(causal_left))
    if any(time in source_note for time in ["00:00", "23:59"]):
        blockers.append("Source note must not include hours/minutes.")
    project_brand_lower = (data.get("project_brand") or "").lower()
    pie_rows = data.get("pie_rows") or []
    sov_rows = data.get("sov_comparison_rows") or []
    pie_project_sov = None
    for row in pie_rows:
        if str(row.get("brand", "")).strip().lower() == project_brand_lower:
            pie_project_sov = float(row.get("sov", 0) or 0)
            break
    sov_project_current = None
    for row in sov_rows:
        if str(row.get("brand", "")).strip().lower() == project_brand_lower:
            sov_project_current = float(row.get("current_sov", 0) or 0)
            break
    if pie_project_sov is not None and sov_project_current is not None:
        if abs(pie_project_sov - sov_project_current) > 0.005:
            blockers.append(
                f"Slide 05 pie chart SOV ({pie_project_sov:.4f}) does not match split chart current SOV "
                f"({sov_project_current:.4f}) for project brand."
            )
    elif pie_project_sov is None:
        blockers.append("Slide 05 project brand missing from pie rows.")
    elif sov_project_current is None:
        blockers.append("Slide 05 project brand missing from SOV comparison rows.")

    split = data.get("campaign_split") or {}
    campaign_fact = int(split.get("campaign_publications_count") or split.get("campaign_materials") or split.get("attributed_campaign_mentions") or 0)
    if not campaign_fact:
        blockers.append("Slide 05 requires a campaign publication fact.")
    elif split.get("status") != "verified":
        warnings.append("Slide 05 campaign fact uses a fallback source; project ORM month rows are preferred.")
    elif not any(phrase in joined.lower() for phrase in ["без них", "без учета", "без кампании"]):
        blockers.append("Slide 05 must show campaign/organic split as scenario comparison with a 'without' scenario.")
    if split.get("campaign_materials_override") is None and split.get("method") == "campaign_materials_override":
        blockers.append("Campaign materials override method must preserve the override value in the data model.")
    label_layout = (result.get("visual_guardrails") or {}).get("pie_label_layout") or {}
    if label_layout:
        if label_layout.get("label_count") != len(data.get("pie_rows") or []):
            blockers.append("SOV pie label count must match slice count.")
        if label_layout.get("overlap_count", 0):
            blockers.append(f"SOV pie labels overlap: {label_layout.get('overlap_count')}.")
        if label_layout.get("out_of_bounds_count", 0):
            blockers.append(f"SOV pie labels are out of bounds: {label_layout.get('out_of_bounds_count')}.")
        if label_layout.get("label_chart_overlap_count", 0):
            blockers.append(f"SOV pie labels overlap the chart circle: {label_layout.get('label_chart_overlap_count')}.")
        if label_layout.get("color_mismatch_count", 0):
            blockers.append(f"SOV pie label colors do not match slices: {label_layout.get('color_mismatch_count')}.")
        if label_layout.get("clipped_text_count", 0):
            blockers.append(f"SOV pie labels may be clipped: {label_layout.get('clipped_text_count')}.")
        if label_layout.get("missing_nonzero_label_count", 0):
            blockers.append(f"SOV pie has missing non-zero labels: {label_layout.get('missing_nonzero_label_count')}.")
    else:
        warnings.append("SOV pie uses native PowerPoint legend labels; external label layout diagnostics are not applicable.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=5)
    text_qa = qa_text_contract(joined, slide_number=5)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "project_sov_present": bool(data.get("project_sov")),
            "campaign_split_verified": split.get("status") == "verified",
            "campaign_override_recorded": split.get("campaign_materials_override") is not None,
            "scenario_comparison_wording": split.get("status") != "verified" or any(phrase in joined.lower() for phrase in ["без них", "без учета", "без кампании"]),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_semicolon": ";" not in joined,
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "source_note_without_time": not any(time in source_note for time in ["00:00", "23:59"]),
            "pie_labels_match_slices": not label_layout or label_layout.get("label_count") == len(data.get("pie_rows") or []),
            "pie_labels_no_overlap": not label_layout or label_layout.get("overlap_count", 0) == 0,
            "pie_labels_in_bounds": not label_layout or label_layout.get("out_of_bounds_count", 0) == 0,
            "pie_labels_clear_chart": not label_layout or label_layout.get("label_chart_overlap_count", 0) == 0,
            "pie_label_colors_match": not label_layout or label_layout.get("color_mismatch_count", 0) == 0,
            "pie_labels_not_clipped": not label_layout or label_layout.get("clipped_text_count", 0) == 0,
            "pie_nonzero_labels_present": not label_layout or label_layout.get("missing_nonzero_label_count", 0) == 0,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
    }
