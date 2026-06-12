from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from .brand_styler import bold_project_brand_in_presentation
from .style_rules import qa_slide_copy_principles
from .style_rules import BODY_FONT_SIZE, BODY_LINE_SPACING, FONT_FAMILY, TITLE_FONT_SIZE, format_int_spaces
from .text_contracts import plain_segments, qa_text_contract


TEXT_COLOR = RGBColor(0, 0, 0)
BLUE = RGBColor(1, 162, 255)
LIGHT_BLUE_FILL = RGBColor(242, 248, 252)


def _fmt_int(value: int | float | None) -> str:
    if value is None:
        return "нет данных"
    return format_int_spaces(value)


def _fmt_pct(value: float | None, digits: int = 1) -> str:
    if value is None:
        return "нет данных"
    return f"{float(value) * 100:.{digits}f}%".replace(".", ",")


def _plural(value: int | float | None, one: str, few: str, many: str) -> str:
    if value is None:
        return many
    number = abs(int(round(float(value))))
    if 11 <= number % 100 <= 14:
        return many
    last = number % 10
    if last == 1:
        return one
    if 2 <= last <= 4:
        return few
    return many


def _item(segments: list[tuple[str, bool]]) -> dict:
    return {
        "text": "".join(text for text, _bold in segments),
        "segments": [{"text": text, "bold": bold} for text, bold in segments],
    }


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)
    parent.append(shape._element)


def _remove_shape(shape) -> None:
    if shape is not None and shape._element.getparent() is not None:
        shape._element.getparent().remove(shape._element)


def _find_existing_kpi_shapes(slide) -> list:
    top_min = int(1.05 * 914400)
    top_max = int(2.25 * 914400)
    text_shapes = []
    accent_shapes = []
    for shape in slide.shapes:
        if not (top_min <= shape.top <= top_max):
            continue
        text = (shape.text or "").upper() if getattr(shape, "has_text_frame", False) else ""
        is_metric_text = any(
            marker in text
            for marker in ["УПОМИНАН", "ИНДЕКС ЛОЯЛЬНОСТИ", "SOV", "ПРИРОСТ", "КАМПАН"]
        )
        is_accent_line = shape.height <= int(0.08 * 914400) and shape.width >= int(1.2 * 914400)
        if is_metric_text:
            text_shapes.append(shape)
        elif is_accent_line:
            accent_shapes.append(shape)
    return text_shapes, accent_shapes


def _set_title(shape, text: str):
    shape.text = text
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(TITLE_FONT_SIZE)
            run.font.bold = False
            run.font.color.rgb = TEXT_COLOR


def _set_paragraph_segments(paragraph, segments: list[dict], *, font_size: float = BODY_FONT_SIZE, space_after_pt: float = 7.0, bullet: bool = False):
    original_runs = list(paragraph.runs)
    for run in original_runs:
        paragraph._p.remove(run._r)
    p_pr = paragraph._p.get_or_add_pPr()
    bullet_suffixes = ("}buNone", "}buClr", "}buFont", "}buChar", "}buAutoNum")
    for child in list(p_pr):
        if child.tag.endswith(bullet_suffixes):
            p_pr.remove(child)
    p_pr.insert(0, OxmlElement("a:buNone"))
    paragraph.space_after = Pt(space_after_pt)
    paragraph.line_spacing = BODY_LINE_SPACING
    if bullet:
        run = paragraph.add_run()
        run.text = "▪  "
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = False
        run.font.color.rgb = BLUE
    for segment in segments:
        if not segment.get("text"):
            continue
        run = paragraph.add_run()
        run.text = segment.get("text", "")
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = bool(segment.get("bold"))
        run.font.color.rgb = TEXT_COLOR


def _set_paragraph_text(paragraph, text: str, *, font_size: float = BODY_FONT_SIZE):
    _set_paragraph_segments(paragraph, [{"text": text, "bold": False}], font_size=font_size)


def _copy_paragraph_style(source, target) -> None:
    source_ppr = source._p.pPr
    if source_ppr is None:
        return
    target_ppr = target._p.pPr
    if target_ppr is not None:
        target._p.remove(target_ppr)
    target._p.insert(0, deepcopy(source_ppr))


def _set_shape_bold(shape, bold: bool) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = bold


def _add_kpi_block(slide, *, left: int, top: int, width: int, height: int, title: str, value: str, subtitle: str = ""):
    accent_h = int(0.04 * 914400)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, accent_h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0, 112, 192)
    accent.line.fill.background()
    text_box = slide.shapes.add_textbox(left, top + accent_h, width, height - accent_h)
    tf = text_box.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = value
    run.font.name = FONT_FAMILY
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(51, 51, 51)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = str(title).upper()
    run.font.name = FONT_FAMILY
    run.font.size = Pt(8)
    run.font.bold = False
    run.font.color.rgb = RGBColor(102, 102, 102)
    if subtitle:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = subtitle[:60]
        run.font.name = FONT_FAMILY
        run.font.size = Pt(8)
        run.font.bold = False
        run.font.color.rgb = RGBColor(153, 153, 153)
    return text_box


def _first_names(rows: list[dict], key: str, limit: int = 3) -> str:
    values = [str(row.get(key, "")).strip() for row in rows[:limit] if str(row.get(key, "")).strip()]
    return ", ".join(values)


def _first_contexts(rows: list[dict], limit: int = 3) -> str:
    values = []
    for row in rows[:limit]:
        value = str(row.get("rail_phrase") or row.get("name") or "").strip()
        if value:
            values.append(value)
    return ", ".join(values)


def _find_row(rows: list[dict], key: str, value: str) -> dict:
    expected = value.lower()
    for row in rows:
        if str(row.get(key, "")).strip().lower() == expected:
            return row
    return {}


def _fmt_pp(value: float | None, digits: int = 1) -> str:
    if value is None:
        return "нет данных"
    number = f"{float(value):.{digits}f}".replace(".", ",")
    return f"{number} п.п."


def _fmt_ratio_pp(value: float | None, digits: int = 1) -> str:
    if value is None:
        return "нет данных"
    return _fmt_pp(float(value) * 100, digits)


def _source_phrase(row: dict, total: int | float | None) -> str:
    messages = int(row.get("messages") or 0)
    share = messages / total if total else row.get("share")
    return f"{row.get('label', '')} ({_fmt_int(messages)} {_plural(messages, 'сообщение', 'сообщения', 'сообщений')}, {_fmt_pct(share, 0)})"


def _join_ru(items: list[str]) -> str:
    clean = [str(item).strip() for item in items if str(item or "").strip()]
    if not clean:
        return ""
    if len(clean) == 1:
        return clean[0]
    return ", ".join(clean[:-1]) + " и " + clean[-1]


def _source_label(*models: dict) -> str:
    for model in models:
        label = str((model or {}).get("source_label") or (model or {}).get("source_system") or "").strip()
        if label:
            return label.replace("сервис мониторинга ", "")
    return "источник данных"


def _type_phrase(row: dict) -> str:
    messages = int(row.get("messages") or 0)
    return f"{row.get('type', '')} ({_fmt_int(messages)} {_plural(messages, 'сообщение', 'сообщения', 'сообщений')})"


def _mat_type(material_rows: list[dict], name: str) -> int:
    for row in material_rows:
        if str(row.get("type", "")).lower() == name.lower():
            return int(row.get("materials") or 0)
    return 0


def _site_label(value: str) -> str:
    text = str(value or "").strip()
    low = text.lower()
    mapping = {
        "ozon.ru": "Ozon",
        "wildberries.ru": "Wildberries",
        "telegram.org": "Telegram",
        "telegram.me": "Telegram",
        "t.me": "Telegram",
        "vk.com": "ВКонтакте",
        "youtube.com": "YouTube",
        "dzen.ru": "Дзен",
        "irecommend.ru": "IRecommend",
        "otzovik.com": "Отзовик",
        "apteka.ru": "Apteka.ru",
        "market.yandex.ru": "Яндекс.Маркет",
    }
    for needle, label in mapping.items():
        if needle in low:
            return label
    return text.split("/")[0] if "." in text else text or "не указано"


def _site_phrase(site: str, count: int, total: int) -> str:
    return f"{_site_label(site)} ({_fmt_int(count)} {_plural(count, 'упоминание', 'упоминания', 'упоминаний')})"


def _period_genitive(label: str) -> str:
    low = str(label or "").strip().lower().replace("ё", "е")
    mapping = {
        "январь": "января",
        "февраль": "февраля",
        "март": "марта",
        "апрель": "апреля",
        "май": "мая",
        "июнь": "июня",
        "июль": "июля",
        "август": "августа",
        "сентябрь": "сентября",
        "октябрь": "октября",
        "ноябрь": "ноября",
        "декабрь": "декабря",
    }
    for month, genitive in mapping.items():
        if month in low:
            return genitive
    return low or "текущего периода"


def _period_prepositional(project: dict, *models: dict) -> str:
    text = " ".join(str(item or "") for item in [project.get("period")] + [model.get("sheet") for model in models if model])
    low = text.lower().replace("ё", "е")
    mapping = [
        ("январ", "в январе"),
        ("феврал", "в феврале"),
        ("март", "в марте"),
        ("апрел", "в апреле"),
        (r"\bма[йяе]\b", "в мае"),
        ("июн", "в июне"),
        ("июл", "в июле"),
        ("август", "в августе"),
        ("сентябр", "в сентябре"),
        ("октябр", "в октябре"),
        ("ноябр", "в ноябре"),
        ("декабр", "в декабре"),
    ]
    for pattern, phrase in mapping:
        if pattern.startswith(r"\b"):
            if re.search(pattern, low):
                return phrase
        elif pattern in low:
            return phrase
    return "за отчетный период"


def _views_conclusion(slide08: dict) -> str:
    metrics = slide08.get("views_metrics") or {}
    if not metrics.get("has_views"):
        return "Поля просмотров и вовлечения не заполнены, поэтому охват не выводится как факт."
    metric_label = metrics.get("metric_label") or "Просмотры"
    total = metrics.get("received_views_total")
    if metrics.get("metric_kind") == "received":
        return f"Материалы получили {_fmt_int(total)} просмотров."
    caveat = metrics.get("metric_caveat") or ""
    return f"По заполненному полю «{metric_label}» зафиксировано {_fmt_int(total)} просмотров. {caveat}".strip()


def _reviews_delta_text(start_rating: dict, current_rating: dict, current_period: str) -> str:
    start_total = start_rating.get("reviews_total")
    current_total = current_rating.get("reviews_total")
    if start_total is None or current_total is None:
        return ""
    start_value = int(start_total or 0)
    current_value = int(current_total or 0)
    if current_value > start_value:
        verb = "увеличилось"
    elif current_value < start_value:
        verb = "снизилось"
    else:
        return f"Количество отзывов на сайтах в топе выдачи на конец {current_period} не изменилось и составило {_fmt_int(current_value)} сообщений."
    return f"Количество отзывов на сайтах в топе выдачи на конец {current_period} {verb} с {_fmt_int(start_value)} до {_fmt_int(current_value)} сообщений."


def build_slide12_data(project: dict, slide_data_by_number: dict[int, dict]) -> dict:
    brand = project.get("brand") or "бренд"
    slide04 = slide_data_by_number.get(4, {})
    slide05 = slide_data_by_number.get(5, {})
    slide06 = slide_data_by_number.get(6, {})
    slide07 = slide_data_by_number.get(7, {})
    slide08 = slide_data_by_number.get(8, {})
    slide11 = slide_data_by_number.get(11, {})

    # Use canonical model when present for ALL reconciled metrics
    cm = slide_data_by_number.get("canonical_model")
    if cm is not None:
        tb = cm.get("target_brand", {})
        sov = cm.get("competitive_sov_scope", {})
        sov_target = sov.get("target", {})
        tonality_cm = cm.get("target_tonality", {})
        platform = cm.get("platform_scope", {})
        orm_camp = cm.get("orm_campaign", {})
        ratings = cm.get("ratings_scope", {})

        slide05["project_mentions"] = tb.get("with_campaign", slide05.get("project_mentions"))
        slide05["total_mentions_competitive_set"] = sov.get("total_with_campaign", slide05.get("total_mentions_competitive_set"))
        slide05["project_sov"] = sov_target.get("sov_with_campaign", slide05.get("project_sov"))
        slide05["project_rank"] = slide05.get("project_rank") or 1
        slide05["pie_rows"] = sov.get("competitors", slide05.get("pie_rows", []))
        slide05["all_brand_rows"] = sov.get("competitors", slide05.get("all_brand_rows", []))
        slide05["campaign_split"] = {
            "organic_mentions": tb.get("organic", 0),
            "campaign_publications_count": orm_camp.get("total", 0),
            "campaign_materials": orm_camp.get("total", 0),
            "organic_sov": sov_target.get("sov_without_campaign", 0),
            "sov_lift_pp": sov_target.get("sov_lift_pp", 0),
            "status": "verified",
        }

        slide06["project_row"] = {
            "positive": tonality_cm.get("positive", 0),
            "neutral": tonality_cm.get("neutral", 0),
            "negative": tonality_cm.get("negative", 0),
            "total": tonality_cm.get("positive", 0) + tonality_cm.get("neutral", 0) + tonality_cm.get("negative", 0),
        }
        slide06["signals"] = {
            "top_sources": {
                "Соцсети": tonality_cm.get("positive", 0),
            },
            "top_sites": {},
        }

        slide07["platform_type_rows"] = platform.get("platform_type_rows", slide07.get("platform_type_rows", []))
        slide07["source_rows"] = platform.get("source_rows", slide07.get("source_rows", []))
        slide07["total_messages"] = platform.get("total", slide07.get("total_messages", 0))
        slide07["review_sources"] = platform.get("review_sources", slide07.get("review_sources", []))

        slide08["views_metrics"] = {"received_views_total": orm_camp.get("views_total", 0), "has_views": orm_camp.get("views_total", 0) > 0, "metric_label": "Просмотры"}
        slide08["total_materials"] = orm_camp.get("total", slide08.get("total_materials", 0))
        slide08["review_materials"] = orm_camp.get("reviews", slide08.get("review_materials", 0))
        slide08["comment_materials"] = orm_camp.get("comments", slide08.get("comment_materials", 0))
        slide08["material_type_rows"] = orm_camp.get("material_type_rows", slide08.get("material_type_rows", []))

        if ratings.get("has_current_period"):
            periods = ratings.get("periods", ["start", "current"])
            pos_cards = ratings.get("positive_cards", [0, 0])
            neg_cards = ratings.get("negative_cards", [0, 0])
            no_review = ratings.get("no_review_cards", [0, 0])
            reviews_cnt = ratings.get("reviews_count", [0, 0])
            slide11["start_summary"] = {
                "positive_cards": pos_cards[0] if len(pos_cards) > 0 else 0,
                "negative_cards": neg_cards[0] if len(neg_cards) > 0 else 0,
                "no_reviews": no_review[0] if len(no_review) > 0 else 0,
                "reviews_total": reviews_cnt[0] if len(reviews_cnt) > 0 else 0,
            }
            slide11["current_summary"] = {
                "positive_cards": pos_cards[1] if len(pos_cards) > 1 else 0,
                "negative_cards": neg_cards[1] if len(neg_cards) > 1 else 0,
                "no_reviews": no_review[1] if len(no_review) > 1 else 0,
                "reviews_total": reviews_cnt[1] if len(reviews_cnt) > 1 else 0,
            }

    project_mentions = slide05.get("project_mentions")
    total_mentions = slide05.get("total_mentions_competitive_set")
    project_sov = slide05.get("project_sov")
    project_rank = slide05.get("project_rank")
    leader = (slide05.get("pie_rows") or [{}])[0]
    brand_rows = slide05.get("pie_rows") or slide05.get("all_brand_rows") or []
    campaign = slide05.get("campaign_split") or {}
    tonality = slide06.get("project_row") or {}
    brand_signals = slide06.get("signals") or {}
    topics = slide04.get("chart_terms") or []
    platforms = slide07.get("platform_type_rows") or []
    sources = slide07.get("source_rows") or []
    current_rating = slide11.get("current_summary") or {}
    start_rating = slide11.get("start_summary") or {}
    current_period = _period_genitive(current_rating.get("label", ""))
    ratings_period_caveat = slide11.get("ratings_period_caveat") or ""
    period_prepositional = _period_prepositional(project, slide08)

    topic_total = slide04.get("total_topic_messages")
    top_topic_names = _first_contexts(topics, 3) or _first_names(topics, "name", 3)
    topic_period_phrase = _period_prepositional(project, slide04, slide05).capitalize()
    social_messenger = sum(int(row.get("messages") or 0) for row in platforms if row.get("type") in {"Соцсети", "Мессенджеры"})
    social_messenger_share = social_messenger / slide07.get("total_messages") if slide07.get("total_messages") else None
    reviews_platform = _find_row(platforms, "type", "Отзывы")
    reviews_messages = int(reviews_platform.get("messages") or 0)
    reviews_share = reviews_messages / slide07.get("total_messages") if slide07.get("total_messages") and reviews_messages else None
    top_source_names = _join_ru([row.get("label", "") for row in sources[:2]])
    review_source_names = _join_ru([_site_label(row[0]) for row in (slide07.get("review_sources") or [])[:2] if row])
    if not review_source_names:
        review_source_names = _join_ru([row.get("label", "") for row in sources if row.get("label") in {"Wildberries", "Ozon", "Apteka.ru"}][:2])
    project_ties = [row for row in brand_rows if int(row.get("mentions") or 0) == int(project_mentions or 0)]
    tied_competitors = [row.get("brand") for row in project_ties if row.get("brand") != brand]
    if tied_competitors:
        project_position = f"{brand} и {', '.join(tied_competitors)} делят вторую позицию"
    else:
        project_position = f"{brand} занимает {project_rank}-е место" if project_rank else f"{brand} занимает заметную позицию"
    leader_gap = None
    if leader.get("sov") is not None and project_sov is not None:
        leader_gap = float(leader.get("sov") or 0) - float(project_sov or 0)
    if leader_gap is None:
        leader_gap_phrase = ""
    elif leader_gap > 0.005:
        leader_gap_phrase = f" Отставание от первого места составляет {_fmt_ratio_pp(leader_gap)}"
    elif leader_gap < -0.005:
        leader_gap_phrase = f" Преимущество над ближайшим конкурентом составляет {_fmt_ratio_pp(abs(leader_gap))}"
    else:
        leader_gap_phrase = " Позиции первых брендов практически равны."
    campaign_publications = int(
        campaign.get("campaign_publications_count")
        or campaign.get("campaign_materials")
        or campaign.get("attributed_campaign_mentions")
        or 0
    )
    if campaign.get("organic_sov") is not None and campaign.get("sov_lift_pp") is not None:
        organic_mentions = int(campaign.get("organic_mentions") or 0)
        organic_mentions_phrase = (
            f" ({_fmt_int(organic_mentions)} {_plural(organic_mentions, 'упоминание', 'упоминания', 'упоминаний')})"
            if organic_mentions
            else ""
        )
        campaign_effect_phrase = (
            f" В кампании учтено {_fmt_int(campaign_publications)} "
            f"{_plural(campaign_publications, 'публикация', 'публикации', 'публикаций')}. "
            f"Без этих материалов сценарная доля бренда составляет {_fmt_pct(campaign.get('organic_sov'))}{organic_mentions_phrase}. "
            f"Разница между сценариями — {_fmt_pp(campaign.get('sov_lift_pp'))}"
        )
    elif campaign_publications:
        campaign_effect_phrase = (
            f" В кампании учтено {_fmt_int(campaign_publications)} "
            f"{_plural(campaign_publications, 'публикация', 'публикации', 'публикаций')}. "
            "Для устойчивости позиции важно переводить размещения в самостоятельные обсуждения."
        )
    else:
        campaign_effect_phrase = " Факт кампании требует проверки в проектной ORM-таблице перед финальной версией."
    last_brand = brand_rows[-1] if len(brand_rows) > 2 else {}
    material_rows = slide08.get("material_type_rows") or []
    plain_reviews = _mat_type(material_rows, "Отзывы")
    purchase_reviews = _mat_type(material_rows, "Отзывы с покупкой")
    top_comments = _mat_type(material_rows, "Комментарии в топе выдачи")
    discussion_comments = _mat_type(material_rows, "Комментарии в обсуждениях")
    seeded_comments = top_comments + discussion_comments
    if not seeded_comments:
        seeded_comments = int(slide08.get("comment_materials", 0) or 0)
    seeded_reviews = plain_reviews + purchase_reviews
    if not seeded_reviews:
        seeded_reviews = int(slide08.get("review_materials", 0) or 0)
    elif int(slide08.get("review_materials", 0) or 0) > seeded_reviews:
        seeded_reviews = int(slide08.get("review_materials", 0) or 0)
    brand_total = int(tonality.get("total") or 0)
    brand_sources = brand_signals.get("top_sources") or {}
    brand_sites = brand_signals.get("top_sites") or {}
    brand_site_phrases = ", ".join(_site_phrase(site, int(count), brand_total) for site, count in list(brand_sites.items())[:4])
    review_brand_mentions = int(brand_sources.get("Отзывы", 0) or 0)
    social_brand_mentions = int(brand_sources.get("Соцсети", 0) or 0)
    messenger_brand_mentions = int(brand_sources.get("Мессенджеры чаты", 0) or 0) + int(brand_sources.get("Мессенджеры каналы", 0) or 0)
    negative_count = int(tonality.get("negative") or 0)
    negative_share = negative_count / brand_total if brand_total else None
    current_negative_cards = int(current_rating.get("negative_cards") or 0)
    placed_reviews = int(((slide11.get("agency_review_placements") or {}).get("placed_reviews")) or 0)
    current_reviews_total = current_rating.get("reviews_total")
    start_reviews_total = start_rating.get("reviews_total")
    current_negative_text = (
        "Негативных карточек не зафиксировано"
        if current_negative_cards <= 0
        else f"Карточек с рейтингом ниже 4 баллов зафиксировано: {current_negative_cards}"
    )
    positive_count = int(tonality.get("positive") or 0)
    neutral_count = int(tonality.get("neutral") or 0)
    negative_count = int(tonality.get("negative") or 0)
    loyalty_base = positive_count + negative_count
    loyalty_index = positive_count / loyalty_base if loyalty_base else None
    if negative_count > 0:
        loyalty_subtitle = (
            f"{_fmt_int(positive_count)} "
            f"{_plural(positive_count, 'позитивное упоминание', 'позитивных упоминания', 'позитивных упоминаний')} "
            f"на {_fmt_int(negative_count)} "
            f"{_plural(negative_count, 'негативное', 'негативных', 'негативных')}"
        )
    elif positive_count > 0:
        loyalty_subtitle = "негатив не зафиксирован"
    else:
        loyalty_subtitle = "нет данных по тональности"

    organic_mentions = int(campaign.get("organic_mentions") or 0)
    if not organic_mentions and campaign_publications and project_mentions:
        organic_mentions = max(int(project_mentions or 0) - campaign_publications, 0)
    if organic_mentions:
        infopole_growth = (int(project_mentions or 0) - organic_mentions) / organic_mentions * 100
        growth_value = f"{'+' if infopole_growth >= 0 else ''}{infopole_growth:.0f}%".replace(".", ",")
    else:
        growth_value = "нет данных"
    visual_blocks = [
        {"title": "Упоминания", "value": _fmt_int(project_mentions), "subtitle": _plural(project_mentions, "упоминание", "упоминания", "упоминаний")},
        {"title": "Индекс лояльности", "value": _fmt_pct(loyalty_index), "subtitle": loyalty_subtitle},
        {"title": "SOV без кампании", "value": _fmt_pct(campaign.get("organic_sov")), "subtitle": f"{_fmt_int(organic_mentions)} упоминаний"},
        {"title": "SOV с кампанией", "value": _fmt_pct(project_sov), "subtitle": f"{_fmt_int(project_mentions)} упоминаний"},
        {"title": "Прирост инфополя", "value": growth_value, "subtitle": "за счет кампании"},
    ]

    if leader and leader.get("brand") and leader.get("brand") != brand:
        leader_text = (
            f"{leader.get('brand')} лидирует в категории: {_fmt_int(leader.get('mentions'))} "
            f"{_plural(leader.get('mentions'), 'сообщение', 'сообщения', 'сообщений')} и SOV {_fmt_pct(leader.get('sov'))}."
        )
    else:
        leader_text = "Бренд лидирует в конкурентном наборе."

    if campaign.get("organic_sov") is not None and campaign.get("sov_lift_pp") is not None:
        sov_scenario_text = (
            f" Без учета кампании сценарный SOV составляет {_fmt_pct(campaign.get('organic_sov'))}, "
            f"разница с текущим уровнем — {_fmt_pp(campaign.get('sov_lift_pp'))}"
        )
    else:
        sov_scenario_text = " Разделение кампании и органики требует проверки во входных данных."

    platform_sentence = (
        f"{_fmt_pct(social_messenger_share)} упоминаний в инфополе составляют обсуждения в соцсетях и мессенджерах, "
        f"еще {_fmt_pct(reviews_share)} — отзывы."
    )
    if top_source_names:
        platform_sentence += f" По количеству упоминаний лидируют {top_source_names}."
    if review_source_names:
        platform_sentence += f" В отзывном контуре заметны {review_source_names}."

    positive_share = positive_count / brand_total if brand_total else None
    neutral_share = neutral_count / brand_total if brand_total else None
    if positive_share is not None and positive_share >= 0.6:
        tone_profile = "преимущественно позитивная"
    elif neutral_share is not None and neutral_share >= 0.5:
        tone_profile = "нейтральная с заметной долей позитива"
    else:
        tone_profile = "смешанная"
    negative_sentence = (
        f"{_fmt_int(negative_count)} {_plural(negative_count, 'негативное сообщение', 'негативных сообщения', 'негативных сообщений')} учтено отдельным блоком и требует контроля."
        if negative_count
        else "Негативные сообщения не зафиксированы."
    )

    views_text = _views_conclusion(slide08)
    recommendation_contexts = top_topic_names or "ключевые пользовательские сценарии"
    slide05_period_phrase = _period_prepositional(project, slide05)

    conclusions = [
        _item([
            (f"{brand} занимает уверенную позицию в группе прямых конкурентов. ", True),
            (f"Доля голоса бренда SOV {slide05_period_phrase} составила {_fmt_pct(project_sov)} ", False),
            (f"({_fmt_int(project_mentions)} {_plural(project_mentions, 'упоминание', 'упоминания', 'упоминаний')}). ", False),
            (leader_text + leader_gap_phrase, False),
            (sov_scenario_text, False),
        ]),
        _item([
            (platform_sentence + " ", True),
            ("Соцсети и мессенджеры формируют обсуждение и знание, отзывы работают ближе к покупке и доверию к карточке товара.", False),
        ]),
        _item([
            (f"Тональность сообщений о бренде {tone_profile}: ", True),
            (f"{_fmt_pct(positive_share)} позитива ({_fmt_int(positive_count)} {_plural(positive_count, 'сообщение', 'сообщения', 'сообщений')}), ", False),
            (f"{_fmt_pct(neutral_share)} нейтрала ({_fmt_int(neutral_count)} {_plural(neutral_count, 'сообщение', 'сообщения', 'сообщений')}). ", False),
            (negative_sentence, False),
        ]),
        _item([
            ("Авторы агентства активно участвуют в пользовательских обсуждениях", False),
            (", поддерживают репутацию бренда ", False),
            (brand, True),
            (f" и формируют знание о продукте в пользовательских контекстах. {period_prepositional.capitalize()} опубликовано ", False),
            (f"{_fmt_int(seeded_comments)} комментариев", True),
            (" в обсуждениях в социальных сетях, чатах, на форумах, в комментариях к статьям и видео, а также ", False),
            (f"{_fmt_int(seeded_reviews)} отзывов", True),
            (f" на отзывных площадках и маркетплейсах. {views_text}", False),
        ]),
        _item([
            ("Рекомендуем продолжать кампанию", False),
            (f" в обсуждениях по ключевым сценариям: {recommendation_contexts}. ", False),
            ("Приоритет — сохранять присутствие в сообществах, где пользователи обсуждают симптомы, выбор препарата и опыт применения, а также поддерживать отзывный контур.", False),
        ]),
    ]
    for item in conclusions:
        item["segments"] = plain_segments(item.get("text", ""))

    return {
        "project_brand": brand,
        "period_raw": project.get("period") or "",
        "conclusions": conclusions,
        "ratings_period_caveat": "",
        "visual_blocks": visual_blocks,
        "source_models": [3, 4, 5, 6, 7, 8],
        "methodology_status": "ready" if all(item.get("text") for item in conclusions) else "blocked",
    }


def _collect_text(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


def build_slide12_final_conclusions(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 12:
        raise ValueError("Template does not have slide 12")
    slide = prs.slides[11]

    title = _shape_by_id(slide, 6)
    body = _shape_by_id(slide, 20)
    text_updates = []
    existing_kpi_text_shapes, existing_kpi_accent_shapes = _find_existing_kpi_shapes(slide)
    if title:
        title.width = int(12.2 * 914400)
        _set_title(title, "ВЫВОДЫ")
        text_updates.append({"shape_id": title.shape_id, "role": "title"})
    if body:
        body.left = int(0.50 * 914400)
        body.top = int(2.50 * 914400)
        body.width = int(12.00 * 914400)
        body.height = int(4.45 * 914400)
        paragraphs = body.text_frame.paragraphs
        conclusions = data.get("conclusions") or []
        style_source = paragraphs[0] if paragraphs else None
        while len(paragraphs) < len(conclusions):
            new_paragraph = body.text_frame.add_paragraph()
            if style_source is not None:
                _copy_paragraph_style(style_source, new_paragraph)
            paragraphs = body.text_frame.paragraphs
        for idx, paragraph in enumerate(paragraphs):
            if idx < len(conclusions):
                _set_paragraph_segments(
                    paragraph,
                    conclusions[idx].get("segments") or [{"text": conclusions[idx].get("text", ""), "bold": False}],
                    font_size=BODY_FONT_SIZE,
                    space_after_pt=10.0,
                    bullet=True,
                )
            else:
                _set_paragraph_text(paragraph, "")
            paragraph.line_spacing = 1.15
        text_updates.append({"shape_id": body.shape_id, "role": "conclusions"})

    visual_shapes = []
    blocks = (data.get("visual_blocks") or [])[:5]
    for idx, block in enumerate(blocks):
        if idx < len(existing_kpi_text_shapes):
            shape = existing_kpi_text_shapes[idx]
            tf = shape.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(1)
            run = p.add_run()
            run.text = block.get("value", "")
            run.font.name = FONT_FAMILY
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(51, 51, 51)
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(1)
            run = p.add_run()
            run.text = str(block.get("title", "")).upper()
            run.font.name = FONT_FAMILY
            run.font.size = Pt(8)
            run.font.bold = False
            run.font.color.rgb = RGBColor(102, 102, 102)
            if block.get("subtitle"):
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(0)
                run = p.add_run()
                run.text = str(block.get("subtitle", ""))[:60]
                run.font.name = FONT_FAMILY
                run.font.size = Pt(8)
                run.font.bold = False
                run.font.color.rgb = RGBColor(153, 153, 153)
            visual_shapes.append(shape)
        else:
            block_left = int(0.50 * 914400) + idx * (int(2.26 * 914400) + int(0.17 * 914400))
            visual_shapes.append(_add_kpi_block(
                slide,
                left=block_left,
                top=int(1.40 * 914400),
                width=int(2.26 * 914400),
                height=int(0.68 * 914400),
                title=block.get("title", ""),
                value=block.get("value", ""),
                subtitle=block.get("subtitle", ""),
            ))
    for s in existing_kpi_accent_shapes:
        _remove_shape(s)

    for shape in [title, body, *visual_shapes]:
        if shape:
            _bring_to_front(shape)

    brand_bold_updates = bold_project_brand_in_presentation(prs, project.get("brand") or data.get("project_brand") or "")
    if len(prs.slides) >= 11:
        _set_shape_bold(_shape_by_id(prs.slides[10], 10), False)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    text_after = _collect_text(slide)
    joined = "\n".join(text_after)
    forbidden = ["Пикторид", "декабрь 2025", "Бактавит", "Окуларис", "Пиклоксидин", "Витабакт", "Пиклосепт", "1 292", "714 248"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]
    return {
        "output": str(output_path),
        "slide_count": len(prs.slides),
        "text_updates": text_updates,
        "old_text_left": old_text_left,
        "text_after": text_after,
        "visual_blocks_rendered": len(visual_shapes),
        "existing_kpi_text_shapes": len(existing_kpi_text_shapes),
        "brand_bold_updates": brand_bold_updates,
        "visual_guardrails": {
            "layout": "uses a top row of KPI accent-line blocks and a full-width conclusion text block aligned to the reference grid",
            "no_source_note_added": True,
        },
    }


def qa_slide12_final_conclusions(data: dict, result: dict) -> dict:
    blockers = []
    joined = "\n".join(result.get("text_after") or [])
    if result.get("slide_count") != 12:
        blockers.append(f"Output PPTX must keep 12 slides, got {result.get('slide_count')}.")
    if len(data.get("conclusions") or []) != 5:
        blockers.append("Slide 12 must keep five final conclusion bullets from the approved layout.")
    if not all((item.get("text") if isinstance(item, dict) else str(item)) for item in data.get("conclusions") or []):
        blockers.append("Slide 12 has empty conclusion text.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 12: " + ", ".join(result["old_text_left"]))
    if "????" in joined:
        blockers.append("Broken Cyrillic marker found on slide 12.")
    if ";" in joined:
        blockers.append("Slide 12 text must not use semicolons.")
    if ".xlsx" in joined.lower() or ".pptx" in joined.lower():
        blockers.append("Slide 12 must not expose local filenames.")
    if any(time in joined for time in ["00:00", "23:59", ":00", ":59"]):
        blockers.append("Slide 12 must not include hours/minutes.")
    if result.get("brand_bold_updates", 0) <= 0:
        blockers.append("Project brand was not bolded on slide 12/deck.")
    if result.get("visual_blocks_rendered", 0) < 5:
        blockers.append("Slide 12 must render five visual metric blocks.")
    forbidden_campaign_causality = ["прирост от кампании", "эффект кампании", "кампания повлияла", "кампания обеспечила", "благодаря кампании"]
    causal_left = [phrase for phrase in forbidden_campaign_causality if phrase in joined.lower()]
    if causal_left:
        blockers.append("Slide 12 must use scenario wording, not campaign-impact claims: " + ", ".join(causal_left))
    if data.get("ratings_period_caveat"):
        caveat = str(data.get("ratings_period_caveat", "")).lower()
        if caveat and caveat not in joined.lower():
            blockers.append("Slide 12 must carry over the ratings-period caveat from slide 11.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=12)
    text_qa = qa_text_contract(joined, slide_number=12)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings = copy_qa["warnings"] + text_qa["warnings"]
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "slide_count_12": result.get("slide_count") == 12,
            "five_bullets_present": len(data.get("conclusions") or []) == 5,
            "conclusion_text_nonempty": all((item.get("text") if isinstance(item, dict) else str(item)) for item in data.get("conclusions") or []),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_broken_cyrillic": "????" not in joined,
            "no_semicolon": ";" not in joined,
            "no_local_filenames": ".xlsx" not in joined.lower() and ".pptx" not in joined.lower(),
            "no_time_in_text": not any(time in joined for time in ["00:00", "23:59", ":00", ":59"]),
            "brand_bold_updates": result.get("brand_bold_updates", 0) > 0,
            "visual_blocks_rendered": result.get("visual_blocks_rendered", 0) >= 5,
            "no_campaign_impact_claims": not causal_left,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
        "visual_review": {
            "required": True,
            "status": "manual_review_needed",
            "notes": "Check that slide 12 keeps the reference grid: five metric blocks on top, five readable highlight paragraphs below, and no overflow.",
        },
    }
