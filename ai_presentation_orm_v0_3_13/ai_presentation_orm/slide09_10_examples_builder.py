from __future__ import annotations

from pathlib import Path
import json
import math
import os
import re
import time

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt
from PIL import Image as PILImage

from .brand_styler import bold_project_brand_in_presentation
from .style_rules import qa_slide_copy_principles
from .example_screenshot_capture import capture_browser_screenshot
from .slide04_problem_field_builder import EMU_PER_INCH
from .style_rules import BODY_FONT_SIZE, FONT_FAMILY, TITLE_FONT_SIZE


TEXT_COLOR = RGBColor(0, 0, 0)
BLACK = (0, 0, 0)
GREY = (112, 112, 112)
LIGHT_GREY = (232, 236, 242)
BLUE = (0, 176, 240)
VK_BLUE = (39, 135, 245)
TG_BLUE = (46, 166, 229)
YT_RED = (255, 0, 0)
ZEN_GREEN = (45, 155, 90)
OZON_BLUE = (0, 91, 255)
WB_PURPLE = (154, 59, 157)
STAR = (243, 184, 49)
CHATGPT_SCREENSHOT_MANIFEST = "chatgpt_screenshot_requests.json"
CHATGPT_SCREENSHOT_EXTENSIONS = (".png", ".jpg", ".jpeg", ".webp")


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)
    parent.append(shape._element)


def _set_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT):
    shape.text = str(text)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = TEXT_COLOR


def _font_path(bold: bool = False) -> str:
    candidates = [
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return candidate
    return "arialbd.ttf" if bold else "arial.ttf"


def _font(size: int, bold: bool = False):
    try:
        return ImageFont.truetype(_font_path(bold), size=size)
    except Exception:
        return ImageFont.load_default()


def _text_width(draw: ImageDraw.ImageDraw, text: str, font) -> int:
    if not text:
        return 0
    box = draw.textbbox((0, 0), text, font=font)
    return int(box[2] - box[0])


def _client_text(value: str) -> str:
    return re.sub(r"\s+", " ", str(value or "").replace(";", ",")).strip()


def _wrap_tokens(draw: ImageDraw.ImageDraw, text: str, max_width: int, font_size: int, brand: str) -> list[list[tuple[str, bool]]]:
    tokens = re.findall(r"\S+\s*", _client_text(text))
    lines: list[list[tuple[str, bool]]] = []
    current: list[tuple[str, bool]] = []
    current_width = 0
    brand_lower = str(brand or "").lower()
    for token in tokens:
        is_brand = bool(brand_lower and brand_lower in token.lower())
        token_font = _font(font_size, bold=is_brand)
        token_width = _text_width(draw, token, token_font)
        if current and current_width + token_width > max_width:
            lines.append(current)
            current = []
            current_width = 0
        current.append((token, is_brand))
        current_width += token_width
    if current:
        lines.append(current)
    return lines


def _trim_line_to_width(draw: ImageDraw.ImageDraw, line: list[tuple[str, bool]], max_width: int, font_size: int) -> list[tuple[str, bool]]:
    ellipsis = "..."
    while line:
        width = 0
        for token, is_brand in line:
            width += _text_width(draw, token, _font(font_size, bold=is_brand))
        width += _text_width(draw, ellipsis, _font(font_size))
        if width <= max_width:
            break
        token, is_brand = line[-1]
        shortened = token.rstrip()
        if len(shortened) > 8:
            line[-1] = (shortened[:-4], is_brand)
        else:
            line.pop()
    line.append((ellipsis, False))
    return line


def _draw_rich_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    *,
    max_width: int,
    font_size: int,
    fill=BLACK,
    brand: str = "",
    max_lines: int | None = None,
    line_spacing: float = 1.22,
) -> int:
    x0, y = xy
    lines = _wrap_tokens(draw, text, max_width, font_size, brand)
    if max_lines is not None and len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = _trim_line_to_width(draw, lines[-1], max_width, font_size)
    line_height = int(font_size * line_spacing)
    for line in lines:
        x = x0
        for token, is_brand in line:
            token_font = _font(font_size, bold=is_brand)
            draw.text((x, y), token, font=token_font, fill=fill)
            x += _text_width(draw, token, token_font)
        y += line_height
    return y


def _draw_wrapped(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    *,
    max_width: int,
    font_size: int,
    fill=BLACK,
    bold: bool = False,
    max_lines: int | None = None,
    line_spacing: float = 1.22,
) -> int:
    font = _font(font_size, bold=bold)
    tokens = re.findall(r"\S+\s*", _client_text(text))
    lines: list[str] = []
    current = ""
    for token in tokens:
        if current and _text_width(draw, current + token, font) > max_width:
            lines.append(current.rstrip())
            current = token
        else:
            current += token
    if current:
        lines.append(current.rstrip())
    if max_lines is not None and len(lines) > max_lines:
        lines = lines[:max_lines]
        while lines[-1] and _text_width(draw, lines[-1] + "...", font) > max_width:
            lines[-1] = lines[-1][:-1].rstrip()
        lines[-1] += "..."
    x, y = xy
    line_height = int(font_size * line_spacing)
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        y += line_height
    return y


def _platform_color(platform: str) -> tuple[int, int, int]:
    low = str(platform or "").lower()
    if "vk" in low or "контакте" in low:
        return VK_BLUE
    if "telegram" in low:
        return TG_BLUE
    if "youtube" in low:
        return YT_RED
    if "дзен" in low:
        return ZEN_GREEN
    if "ozon" in low:
        return OZON_BLUE
    if "wildberries" in low:
        return WB_PURPLE
    if "woman" in low:
        return (223, 71, 116)
    return BLUE


def _format_card_date(value: str) -> str:
    text = _client_text(value)
    if not text:
        return ""
    match = re.search(r"(20\d{2})-(\d{1,2})-(\d{1,2})", text)
    if match:
        year, month, day = match.groups()
        return f"{day.zfill(2)}.{month.zfill(2)}.{year}"
    match = re.search(r"(\d{1,2})\.(\d{1,2})\.(\d{2,4})", text)
    if match:
        day, month, year = match.groups()
        if len(year) == 2:
            year = f"20{year}"
        return f"{day.zfill(2)}.{month.zfill(2)}.{year}"
    return re.sub(r"\s+00:00:00$", "", text)[:24]


def _display_author(card: dict) -> str:
    for key in ("nick", "author"):
        value = _client_text(card.get(key, ""))
        if value and value.lower() not in {"none", "nan", "не указано"}:
            return value[:44]
    return ""


def _meta_line(card: dict, *, include_platform: bool = False) -> str:
    parts = []
    author = _display_author(card)
    date = _format_card_date(card.get("date", ""))
    platform = _client_text(card.get("platform", ""))
    if author:
        parts.append(author)
    if date:
        parts.append(date)
    if include_platform and platform:
        parts.append(platform[:32])
    return " · ".join(parts)


def _canvas(width_px: int, height_px: int) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    image = Image.new("RGB", (width_px, height_px), "white")
    draw = ImageDraw.Draw(image)
    return image, draw


def _slot_pixels(width_emu: int, height_emu: int, dpi: int = 240) -> tuple[int, int]:
    return max(640, int(width_emu / EMU_PER_INCH * dpi)), max(280, int(height_emu / EMU_PER_INCH * dpi))


def _screenshot_ok(status: dict) -> bool:
    return status.get("status") in {"captured", "cached", "downloaded", "download_cached"}


def _valid_image_file(path: Path) -> bool:
    try:
        with PILImage.open(path) as image:
            width, height = image.size
        return width >= 120 and height >= 90
    except Exception:
        return False


def _chatgpt_screenshot_candidates(screenshot_dir: Path, stem: str) -> list[Path]:
    direct = [screenshot_dir / f"{stem}_chatgpt{ext}" for ext in CHATGPT_SCREENSHOT_EXTENSIONS]
    nested = [screenshot_dir / "chatgpt" / f"{stem}{ext}" for ext in CHATGPT_SCREENSHOT_EXTENSIONS]
    return direct + nested


def _write_chatgpt_screenshot_request(
    *,
    screenshot_dir: Path,
    stem: str,
    card: dict,
    width_px: int,
    height_px: int,
    reason: str,
) -> None:
    screenshot_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = screenshot_dir / CHATGPT_SCREENSHOT_MANIFEST
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8")) if manifest_path.exists() else {}
    except Exception:
        manifest = {}
    requests = list(manifest.get("requests") or [])
    expected_path = _chatgpt_screenshot_candidates(screenshot_dir, stem)[0]
    request = {
        "id": stem,
        "status": "pending",
        "url": str(card.get("link") or ""),
        "platform": str(card.get("platform") or ""),
        "row": card.get("row"),
        "title": str(card.get("title") or ""),
        "target_text": _client_text(card.get("text", ""))[:600],
        "expected_path": str(expected_path),
        "accepted_alternative_paths": [str(path) for path in _chatgpt_screenshot_candidates(screenshot_dir, stem)[1:]],
        "width_px": int(width_px),
        "height_px": int(height_px),
        "reason": reason,
    }
    by_id = {str(item.get("id")): item for item in requests if item.get("id")}
    by_id[stem] = request
    payload = {
        "schema_version": 1,
        "created_at_epoch": int(time.time()),
        "purpose": "ChatGPT Custom GPT screenshot handoff for slides 09-10.",
        "agent_instructions": [
            "Open each pending URL with ChatGPT browsing/browser tools.",
            "Capture a real screenshot of the source page or visible message/review; do not redraw or synthesize message content.",
            "Save the image exactly to expected_path, or to one of accepted_alternative_paths.",
            "Rerun the PPTX pipeline after at least five usable images per slide are saved.",
        ],
        "requests": list(by_id.values()),
    }
    manifest_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _chatgpt_supplied_screenshot(screenshot_dir: Path, stem: str, card: dict, width_px: int, height_px: int) -> tuple[Path | None, dict | None]:
    for path in _chatgpt_screenshot_candidates(screenshot_dir, stem):
        if path.exists() and _valid_image_file(path):
            return path, {
                "url": str(card.get("link") or ""),
                "image_path": str(path),
                "width_px": int(width_px),
                "height_px": int(height_px),
                "capture_version": "chatgpt-handoff-v1",
                "captured_at_epoch": int(time.time()),
                "status": "captured",
                "reason": "ok",
                "matched_text": False,
                "matched_snippet": "",
                "crop_mode": "external",
                "capture_mode": "chatgpt_handoff",
                "source_kind": "chatgpt_browser_capture",
            }
    return None, None


def _capture_card_screenshot(card: dict, screenshot_dir: Path, stem: str, width_px: int, height_px: int, *, compact_crop: bool = False) -> tuple[Path, dict]:
    screenshot_path = screenshot_dir / f"{stem}_browser.png"
    supplied_path, supplied_status = _chatgpt_supplied_screenshot(screenshot_dir, stem, card, width_px, height_px)
    if supplied_path is not None and supplied_status is not None:
        return supplied_path, supplied_status
    text = _client_text(card.get("text", ""))
    title = _client_text(card.get("title", ""))
    if title and text.lower().startswith(title.lower()):
        text = text[len(title):].lstrip(" .,!?:-—") or text
    target_text = text or title
    if os.environ.get("AI_ORM_SCREENSHOT_BACKEND", "").strip().lower() in {"chatgpt", "external", "manifest"}:
        status = {
            "status": "failed",
            "reason": "chatgpt_screenshot_required",
            "url": str(card.get("link") or ""),
            "image_path": str(screenshot_path),
            "capture_mode": "chatgpt_handoff",
        }
        _write_chatgpt_screenshot_request(
            screenshot_dir=screenshot_dir,
            stem=stem,
            card=card,
            width_px=width_px,
            height_px=height_px,
            reason=status["reason"],
        )
        return screenshot_path, status
    status = capture_browser_screenshot(
        url=str(card.get("link") or ""),
        output_path=screenshot_path,
        target_text=target_text,
        width_px=width_px,
        height_px=height_px,
        compact_crop=compact_crop,
    )
    if not _screenshot_ok(status):
        _write_chatgpt_screenshot_request(
            screenshot_dir=screenshot_dir,
            stem=stem,
            card=card,
            width_px=width_px,
            height_px=height_px,
            reason=str(status.get("reason") or "browser_capture_failed"),
        )
    return screenshot_path, status


def _would_create_platform_run(platforms: list[str], platform: str, max_run: int = 2) -> bool:
    if not platform or len(platforms) < max_run:
        return False
    return all(item == platform for item in platforms[-max_run:])


def _capture_validation(
    cards: list[dict],
    card_index: int,
    cache: dict[tuple[int, int, int, int], tuple[Path, dict]],
    screenshot_dir: Path,
    stem_prefix: str,
    width_px: int,
    height_px: int,
    compact_crop: bool = False,
) -> tuple[Path, dict]:
    key = (card_index, width_px, height_px, int(bool(compact_crop)))
    if key not in cache:
        cache[key] = _capture_card_screenshot(cards[card_index], screenshot_dir, f"{stem_prefix}_candidate_{card_index + 1}", width_px, height_px, compact_crop=compact_crop)
    return cache[key]


def _pick_browser_validated_card(
    cards: list[dict],
    used_indexes: set[int],
    validation_cache: dict[tuple[int, int, int, int], tuple[Path, dict]],
    screenshot_dir: Path,
    stem_prefix: str,
    width_px: int,
    height_px: int,
    placed_platforms: list[str],
    compact_crop: bool = False,
) -> tuple[int | None, Path | None, dict | None]:
    fallback: tuple[int, Path, dict] | None = None
    for strict_platform_run in (True, False):
        for card_index, card in enumerate(cards):
            if card_index in used_indexes:
                continue
            platform = _client_text(card.get("platform", ""))
            if strict_platform_run and _would_create_platform_run(placed_platforms, platform):
                continue
            screenshot_path, status = _capture_validation(cards, card_index, validation_cache, screenshot_dir, stem_prefix, width_px, height_px, compact_crop=compact_crop)
            if _screenshot_ok(status) and bool(status.get("matched_text")):
                return card_index, screenshot_path, status
            if fallback is None and _screenshot_ok(status):
                fallback = (card_index, screenshot_path, status)
        if len(placed_platforms) < 2:
            break
    if fallback is not None:
        return fallback
    return None, None, None


def _browser_capture_status(
    status: dict,
    *,
    role: str,
    card: dict,
    screenshot_path: Path,
    display_path: Path | None = None,
) -> dict:
    display_path = display_path or screenshot_path
    return {
        **dict(status or {}),
        "role": role,
        "source_row": card.get("row"),
        "platform": card.get("platform"),
        "source_kind": status.get("source_kind") or "browser_capture",
        "ready_asset_status": {"status": "skipped", "reason": "browser_capture_preferred"},
        "fallback_used": False,
        "render_mode": status.get("render_mode") or "browser_message_screenshot",
        "image_path": str(display_path),
        "raw_screenshot_path": str(screenshot_path) if Path(display_path) != Path(screenshot_path) else "",
        "postprocess": "real_screenshot_crop_to_slot" if Path(display_path) != Path(screenshot_path) else "none",
    }


def _display_url(url: str, limit: int = 96) -> str:
    text = _client_text(url).replace(";", "%3B")
    if len(text) <= limit:
        return text
    return text[: max(24, limit - 3)] + "..."


def _next_manual_link_card(cards: list[dict], used_indexes: set[int]) -> tuple[int | None, dict | None]:
    for index, card in enumerate(cards):
        if index in used_indexes:
            continue
        if str(card.get("link") or "").strip().lower().startswith(("http://", "https://")):
            return index, card
    return None, None


def _add_manual_link_card(slide, card: dict, pos: tuple[int, int, int, int], role: str, brand: str):
    left, top, width, height = pos
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    try:
        shape.line.fill.background()
    except Exception:
        shape.line.color.rgb = RGBColor(255, 255, 255)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(7)
    tf.margin_right = Pt(7)
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(5)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    platform = _client_text(card.get("platform", "")) or "source"
    title = _client_text(card.get("title", "")) or "Ссылка для ручного скриншота"
    text = _client_text(card.get("text", ""))
    url = str(card.get("link") or "").strip()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = platform[:48]
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*BLUE)

    p = tf.add_paragraph()
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = title[:110]
    run.font.name = "Arial"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = TEXT_COLOR

    if text:
        p = tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text[:170]
        run.font.name = "Arial"
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(*GREY)

    p = tf.add_paragraph()
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = _display_url(url)
    run.font.name = "Arial"
    run.font.size = Pt(7)
    run.font.color.rgb = RGBColor(0, 100, 180)
    try:
        run.hyperlink.address = url
    except Exception:
        pass

    _bring_to_front(shape)
    return shape


def _manual_link_status(*, role: str, card: dict, shape_id: int) -> dict:
    return {
        "status": "manual_link_inserted",
        "reason": "screenshot_unavailable_link_inserted",
        "role": role,
        "source_row": card.get("row"),
        "platform": card.get("platform"),
        "url": str(card.get("link") or ""),
        "source_kind": "manual_screenshot_link",
        "ready_asset_status": {"status": "skipped", "reason": "manual_link_fallback_used"},
        "fallback_used": False,
        "render_mode": "manual_link_for_screenshot",
        "shape_id": shape_id,
    }


def _add_picture_fit(slide, image_path: Path, pos: tuple[int, int, int, int]):
    left, top, slot_width, slot_height = pos
    try:
        with PILImage.open(image_path) as image:
            image_width, image_height = image.size
    except Exception:
        return slide.shapes.add_picture(str(image_path), *pos)
    if image_width <= 0 or image_height <= 0:
        return slide.shapes.add_picture(str(image_path), *pos)
    image_ratio = image_width / image_height
    slot_ratio = slot_width / max(slot_height, 1)
    if image_ratio >= slot_ratio:
        width = slot_width
        height = int(slot_width / image_ratio)
    else:
        height = slot_height
        width = int(slot_height * image_ratio)
    left = int(left + (slot_width - width) / 2)
    top = int(top + (slot_height - height) / 2)
    return slide.shapes.add_picture(str(image_path), left, top, width, height)


def _crop_real_screenshot_to_slot(image_path: Path, output_path: Path, target_ratio: float, *, anchor_y: float = 0.18) -> Path:
    image_path = Path(image_path)
    output_path = Path(output_path)
    try:
        with PILImage.open(image_path) as source:
            image = source.convert("RGB")
    except Exception:
        return image_path
    width, height = image.size
    if width <= 0 or height <= 0 or target_ratio <= 0:
        return image_path
    image_ratio = width / height
    if image_ratio < target_ratio:
        new_height = max(1, int(width / target_ratio))
        if new_height < height:
            top = int(max(0, min(height - new_height, (height - new_height) * anchor_y)))
            image = image.crop((0, top, width, top + new_height))
    elif image_ratio > target_ratio:
        new_width = max(1, int(height * target_ratio))
        if new_width < width:
            left = int(max(0, min(width - new_width, (width - new_width) / 2)))
            image = image.crop((left, 0, left + new_width, height))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path)
    return output_path


def _collect_text(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


def _picture_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == 13]


def _bbox(shapes) -> tuple[int, int, int, int] | None:
    shapes = list(shapes)
    if not shapes:
        return None
    left = min(shape.left for shape in shapes)
    top = min(shape.top for shape in shapes)
    right = max(shape.left + shape.width for shape in shapes)
    bottom = max(shape.top + shape.height for shape in shapes)
    return left, top, right - left, bottom - top


def _rects_intersect(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by


def _remove_empty_frame_shapes(slide, slots: list[tuple[int, int, int, int]]) -> int:
    removed = 0
    for shape in list(slide.shapes):
        if shape.shape_type == 13:
            continue
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            continue
        if shape.width < EMU_PER_INCH or shape.height < EMU_PER_INCH:
            continue
        shape_rect = (shape.left, shape.top, shape.width, shape.height)
        if not any(_rects_intersect(shape_rect, slot) for slot in slots):
            continue
        _remove_shape(slide, shape)
        removed += 1
    return removed


def _fallback_slot(left: float, top: float, width: float, height: float) -> tuple[int, int, int, int]:
    return tuple(int(value * EMU_PER_INCH) for value in (left, top, width, height))


def _example_slots(slide_number: int) -> list[tuple[int, int, int, int]]:
    if slide_number == 9:
        return [
            _fallback_slot(0.62, 2.08, 5.55, 1.28),
            _fallback_slot(0.62, 3.58, 5.55, 1.28),
            _fallback_slot(0.62, 5.08, 5.55, 1.28),
            _fallback_slot(6.96, 1.38, 5.55, 1.28),
            _fallback_slot(6.96, 3.25, 5.55, 1.28),
            _fallback_slot(6.96, 5.12, 5.55, 1.28),
        ]
    return [
        _fallback_slot(0.58, 2.08, 6.05, 1.72),
        _fallback_slot(0.58, 4.16, 6.05, 1.72),
        _fallback_slot(7.03, 1.34, 5.61, 1.55),
        _fallback_slot(7.03, 3.16, 5.61, 1.55),
        _fallback_slot(7.03, 4.98, 5.61, 1.55),
    ]


def _flatten_review_cards(cards: list[dict], brand: str) -> list[dict]:
    flat: list[dict] = []
    for card in cards:
        if card.get("kind") == "marketplace_reviews":
            for item in card.get("items") or []:
                flat.append({
                    "kind": "review_site",
                    "row": item.get("row"),
                    "platform": item.get("platform"),
                    "product": item.get("product") or brand,
                    "title": item.get("title") or "Отзыв о бренде",
                    "text": item.get("text") or "",
                    "link": item.get("link") or "",
                    "screenshot_link": item.get("screenshot_link") or "",
                    "date": item.get("date") or "",
                    "nick": item.get("nick") or "",
                    "author": item.get("author") or "",
                })
        else:
            flat.append(card)
    return flat


def _all_card_text(data: dict) -> str:
    parts = [data.get("summary", "")]
    for card in data.get("cards") or []:
        parts.append(card.get("title", ""))
        parts.append(card.get("text", ""))
        for item in card.get("items") or []:
            parts.append(item.get("text", ""))
            parts.append(item.get("product", ""))
    return "\n".join(str(part) for part in parts if part)


def build_slide09_message_examples(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 9:
        raise ValueError("Template does not have slide 9")
    slide = prs.slides[8]
    screenshot_dir = output_path.parent / "slide09_10_screenshots"
    title = _shape_by_id(slide, 12)
    summary = _shape_by_id(slide, 7)
    slots = _example_slots(9)

    text_updates = []
    if title:
        _set_text(title, str(data.get("title") or "Примеры сообщений о бренде").upper(), font_size=TITLE_FONT_SIZE)
        text_updates.append({"shape_id": title.shape_id, "role": "title"})
    if summary:
        _set_text(summary, str(data.get("summary") or ""), font_size=BODY_FONT_SIZE)
        text_updates.append({"shape_id": summary.shape_id, "role": "summary"})

    removed_pictures = len(_picture_shapes(slide))
    for shape in list(_picture_shapes(slide)):
        _remove_shape(slide, shape)
    removed_frames = _remove_empty_frame_shapes(slide, slots)

    image_outputs = []
    image_shapes = []
    manual_link_shapes = []
    screenshot_statuses = []
    cards = list(data.get("cards") or [])
    brand = project.get("brand") or data.get("project_brand") or ""
    used_card_indexes: set[int] = set()
    placed_platforms: list[str] = []
    validation_cache: dict[tuple[int, int, int, int], tuple[Path, dict]] = {}
    for idx, pos in enumerate(slots):
        if not cards:
            continue
        width_px, height_px = _slot_pixels(pos[2], pos[3])
        role = f"card_{idx + 1}"
        card_index, screenshot_path, screenshot_status = _pick_browser_validated_card(
            cards,
            used_card_indexes,
            validation_cache,
            screenshot_dir,
            f"slide09_{role}",
            width_px,
            height_px,
            placed_platforms,
        )
        if card_index is None:
            card_index, card = _next_manual_link_card(cards, used_card_indexes)
            if card_index is None or card is None:
                continue
            used_card_indexes.add(card_index)
            placed_platforms.append(_client_text(card.get("platform", "")))
            shape = _add_manual_link_card(slide, card, pos, role, brand)
            status = _manual_link_status(role=role, card=card, shape_id=shape.shape_id)
            screenshot_statuses.append(status)
            manual_link_shapes.append({
                "shape_id": shape.shape_id,
                "role": role,
                "source_row": card.get("row"),
                "render_mode": status.get("render_mode"),
            })
            continue
        card = cards[card_index]
        used_card_indexes.add(card_index)
        placed_platforms.append(_client_text(card.get("platform", "")))
        image_path = screenshot_path
        screenshot_status = _browser_capture_status(
            screenshot_status,
            role=role,
            card=card,
            screenshot_path=screenshot_path,
        )
        screenshot_statuses.append(screenshot_status)
        picture = _add_picture_fit(slide, image_path, pos)
        _bring_to_front(picture)
        image_outputs.append(str(image_path))
        image_shapes.append({"shape_id": picture.shape_id, "role": role, "source_row": card.get("row"), "screenshot_status": screenshot_status.get("status"), "render_mode": screenshot_status.get("render_mode")})

    for shape in [title, summary]:
        if shape:
            _bring_to_front(shape)
    brand_bold_updates = bold_project_brand_in_presentation(prs, brand)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    text_after = _collect_text(slide)
    joined = "\n".join(text_after)
    forbidden = ["ПИКТОРИД", "Пикторид", "декабре", "170 комментариев", "конъюнктивит"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]
    return {
        "output": str(output_path),
        "slide_count": len(prs.slides),
        "text_updates": text_updates,
        "image_outputs": image_outputs,
        "image_shapes": image_shapes,
        "manual_link_shapes": manual_link_shapes,
        "screenshot_statuses": screenshot_statuses,
        "removed_picture_count": removed_pictures,
        "removed_frame_count": removed_frames,
        "old_text_left": old_text_left,
        "summary_text": data.get("summary", ""),
        "text_after": text_after,
        "image_card_text": _all_card_text(data),
        "brand_bold_updates": brand_bold_updates,
        "visual_guardrails": {
            "layout": "uses inherited reference slide 9 card slots",
            "screenshots": "uses real browser screenshots by source URL; when screenshots are unavailable, inserts source links for manual capture",
        },
    }


def build_slide10_review_examples(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 10:
        raise ValueError("Template does not have slide 10")
    slide = prs.slides[9]
    screenshot_dir = output_path.parent / "slide09_10_screenshots"
    title = _shape_by_id(slide, 12)
    summary = _shape_by_id(slide, 7)
    slots = _example_slots(10)

    text_updates = []
    if title:
        _set_text(title, str(data.get("title") or "Примеры сообщений о бренде").upper(), font_size=TITLE_FONT_SIZE)
        text_updates.append({"shape_id": title.shape_id, "role": "title"})
    if summary:
        _set_text(summary, str(data.get("summary") or ""), font_size=BODY_FONT_SIZE)
        text_updates.append({"shape_id": summary.shape_id, "role": "summary"})

    removed_pictures = len(_picture_shapes(slide))
    for shape in list(_picture_shapes(slide)):
        _remove_shape(slide, shape)
    removed_frames = _remove_empty_frame_shapes(slide, slots)

    image_outputs = []
    image_shapes = []
    manual_link_shapes = []
    screenshot_statuses = []
    brand = project.get("brand") or data.get("project_brand") or ""
    cards = _flatten_review_cards(list(data.get("cards") or []), brand)
    used_card_indexes: set[int] = set()
    placed_platforms: list[str] = []
    validation_cache: dict[tuple[int, int, int, int], tuple[Path, dict]] = {}
    for idx, pos in enumerate(slots):
        if not cards:
            continue
        role = f"review_{idx + 1}"
        width_px, height_px = _slot_pixels(pos[2], pos[3])
        card_index, screenshot_path, screenshot_status = _pick_browser_validated_card(
            cards,
            used_card_indexes,
            validation_cache,
            screenshot_dir,
            f"slide10_{role}",
            width_px,
            height_px,
            placed_platforms,
            compact_crop=True,
        )
        if card_index is None:
            card_index, card = _next_manual_link_card(cards, used_card_indexes)
            if card_index is None or card is None:
                continue
            used_card_indexes.add(card_index)
            placed_platforms.append(_client_text(card.get("platform", "")))
            shape = _add_manual_link_card(slide, card, pos, role, brand)
            status = _manual_link_status(role=role, card=card, shape_id=shape.shape_id)
            screenshot_statuses.append(status)
            manual_link_shapes.append({
                "shape_id": shape.shape_id,
                "role": role,
                "source_row": card.get("row"),
                "render_mode": status.get("render_mode"),
            })
            continue
        card = cards[card_index]
        used_card_indexes.add(card_index)
        placed_platforms.append(_client_text(card.get("platform", "")))
        image_path = _crop_real_screenshot_to_slot(
            screenshot_path,
            screenshot_dir / f"slide10_{role}_slot_crop.png",
            pos[2] / max(pos[3], 1),
        )
        screenshot_status = _browser_capture_status(
            screenshot_status,
            role=role,
            card=card,
            screenshot_path=screenshot_path,
            display_path=image_path,
        )
        screenshot_statuses.append(screenshot_status)
        picture = _add_picture_fit(slide, image_path, pos)
        _bring_to_front(picture)
        image_outputs.append(str(image_path))
        image_shapes.append({"shape_id": picture.shape_id, "role": role, "source_row": card.get("row"), "screenshot_status": screenshot_status.get("status"), "render_mode": screenshot_status.get("render_mode")})

    for shape in [title, summary]:
        if shape:
            _bring_to_front(shape)
    brand_bold_updates = bold_project_brand_in_presentation(prs, brand)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    text_after = _collect_text(slide)
    joined = "\n".join(text_after)
    forbidden = ["ПИКТОРИД", "Пикторид", "декабре", "капель", "Солофарм"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]
    return {
        "output": str(output_path),
        "slide_count": len(prs.slides),
        "text_updates": text_updates,
        "image_outputs": image_outputs,
        "image_shapes": image_shapes,
        "manual_link_shapes": manual_link_shapes,
        "screenshot_statuses": screenshot_statuses,
        "removed_picture_count": removed_pictures,
        "removed_frame_count": removed_frames,
        "old_text_left": old_text_left,
        "summary_text": data.get("summary", ""),
        "text_after": text_after,
        "image_card_text": _all_card_text(data),
        "brand_bold_updates": brand_bold_updates,
        "visual_guardrails": {
            "layout": "uses inherited reference slide 10 two-column screenshot area",
            "screenshots": "uses real browser screenshots by source URL; when screenshots are unavailable, inserts source links for manual capture",
        },
    }


def _qa_common(data: dict, result: dict, *, slide_number: int, minimum_images: int, minimum_cards: int) -> dict:
    blockers = []
    warnings = []
    joined = "\n".join(result.get("text_after") or [])
    image_text = result.get("image_card_text") or ""
    all_client_text = f"{joined}\n{image_text}"
    screenshot_statuses = result.get("screenshot_statuses") or []
    captured_screenshots = [status for status in screenshot_statuses if _screenshot_ok(status)]
    ready_assets = [status for status in screenshot_statuses if status.get("source_kind") == "ready_screenshot_link"]
    manual_links = [status for status in screenshot_statuses if status.get("source_kind") == "manual_screenshot_link"]
    fallback_assets = [status for status in screenshot_statuses if status.get("fallback_used")]
    allowed_source_kinds = {"browser_capture", "chatgpt_browser_capture", "manual_screenshot_link"}
    invalid_source_kinds = [status.get("source_kind") for status in screenshot_statuses if status.get("source_kind") not in allowed_source_kinds]
    source_rows = [status.get("source_row") for status in screenshot_statuses if status.get("source_row")]
    usable_visual_slots = len(captured_screenshots) + len(manual_links)
    if result.get("slide_count") != 12:
        blockers.append(f"Output PPTX must keep 12 slides, got {result.get('slide_count')}.")
    if data.get("methodology_status") != "ready":
        blockers.append(f"Slide {slide_number:02d} examples are missing.")
    if len(data.get("cards") or []) < minimum_cards:
        blockers.append(f"Slide {slide_number:02d} needs at least {minimum_cards} example cards.")
    if usable_visual_slots < minimum_images:
        blockers.append(f"Slide {slide_number:02d} needs at least {minimum_images} real browser screenshots or manual screenshot links.")
    if fallback_assets:
        blockers.append(f"Slide {slide_number:02d} must not use rendered or generated screenshot fallbacks.")
    if ready_assets:
        blockers.append(f"Slide {slide_number:02d} must not use ready screenshot links as visual assets.")
    if invalid_source_kinds:
        blockers.append(f"Slide {slide_number:02d} has invalid screenshot source kinds: " + ", ".join(str(item) for item in invalid_source_kinds[:4]))
    if len(source_rows) != len(set(source_rows)):
        blockers.append(f"Slide {slide_number:02d} must use unique source rows for examples.")
    small_images = []
    for status in captured_screenshots:
        image_path = status.get("image_path")
        if not image_path:
            continue
        try:
            with PILImage.open(image_path) as image:
                width, height = image.size
            if width < 120 or height < 90:
                small_images.append(image_path)
        except Exception:
            small_images.append(image_path)
    if small_images:
        blockers.append(f"Slide {slide_number:02d} has missing or too small screenshot images.")
    if result.get("removed_picture_count", 0) == 0:
        blockers.append(f"Slide {slide_number:02d} old template pictures were not removed.")
    if result.get("old_text_left"):
        blockers.append(f"Stale template text remains on slide {slide_number:02d}: " + ", ".join(result["old_text_left"]))
    if "????" in all_client_text:
        blockers.append(f"Broken Cyrillic marker found on slide {slide_number:02d}.")
    if ";" in all_client_text:
        blockers.append(f"Slide {slide_number:02d} text must not use semicolons.")
    if ".xlsx" in joined.lower() or ".pptx" in joined.lower():
        blockers.append(f"Slide {slide_number:02d} text must not expose local filenames.")
    if result.get("brand_bold_updates", 0) <= 0:
        warnings.append(f"No text-shape brand bolding detected on slide {slide_number:02d}.")
    if manual_links and len(captured_screenshots) < minimum_images:
        warnings.append(f"Slide {slide_number:02d} uses manual screenshot links because browser screenshots are unavailable.")
    elif screenshot_statuses and not captured_screenshots:
        warnings.append(f"Slide {slide_number:02d} has no successful real browser screenshots.")
    copy_subject = result.get("summary_text") or data.get("summary") or ""
    copy_qa = qa_slide_copy_principles(copy_subject, slide_number=slide_number)
    blockers.extend(copy_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    qa_status = "ready_for_visual_review" if not blockers else "blocked"
    if not blockers and manual_links and len(captured_screenshots) < minimum_images:
        qa_status = "ready_for_manual_screenshot_links"
    return {
        "status": qa_status,
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "slide_count_12": result.get("slide_count") == 12,
            "examples_ready": data.get("methodology_status") == "ready",
            "minimum_cards": len(data.get("cards") or []) >= minimum_cards,
            "images_generated": len(result.get("image_outputs") or []) + len(manual_links) >= minimum_images,
            "manual_links_inserted": len(manual_links),
            "screenshots_or_manual_links_ready": usable_visual_slots >= minimum_images,
            "old_pictures_removed": result.get("removed_picture_count", 0) > 0,
            "no_stale_template_text": not result.get("old_text_left"),
            "no_broken_cyrillic": "????" not in all_client_text,
            "no_semicolon": ";" not in all_client_text,
            "no_local_filenames": ".xlsx" not in joined.lower() and ".pptx" not in joined.lower(),
            "brand_bold_updates": result.get("brand_bold_updates", 0) > 0,
            "screenshot_assets_attempted": len(screenshot_statuses),
            "screenshot_assets_captured": len(captured_screenshots),
            "ready_screenshot_assets_used": len(ready_assets),
            "manual_screenshot_links_used": len(manual_links),
            "rendered_fallback_assets_used": len(fallback_assets),
            "all_images_real_browser_screenshots": len(result.get("image_outputs") or []) == len(captured_screenshots),
            "valid_screenshot_source_kinds": not invalid_source_kinds,
            "unique_source_rows": len(source_rows) == len(set(source_rows)),
            "screenshot_images_min_size": not small_images,
            "copy_principles": copy_qa["checks"],
        },
        "visual_review": {
            "required": True,
            "status": "manual_link_capture_needed" if manual_links and len(captured_screenshots) < minimum_images else "manual_review_needed",
            "notes": (
                f"Use slide {slide_number:02d} links to capture screenshots manually, then rerun or replace cards."
                if manual_links and len(captured_screenshots) < minimum_images
                else f"Check slide {slide_number:02d} real browser message screenshots for readability and preserved reference geometry."
            ),
        },
    }


def qa_slide09_message_examples(data: dict, result: dict) -> dict:
    qa = _qa_common(data, result, slide_number=9, minimum_images=5, minimum_cards=5)
    if int(data.get("comment_count", 0) or 0) <= 0:
        qa["blockers"].append("Slide 09 comment count is missing.")
        qa["status"] = "blocked"
    qa["checks"]["comment_count_present"] = int(data.get("comment_count", 0) or 0) > 0
    return qa


def qa_slide10_review_examples(data: dict, result: dict) -> dict:
    qa = _qa_common(data, result, slide_number=10, minimum_images=5, minimum_cards=5)
    if int(data.get("review_count", 0) or 0) <= 0:
        qa["blockers"].append("Slide 10 review count is missing.")
        qa["status"] = "blocked"
    qa["checks"]["review_count_present"] = int(data.get("review_count", 0) or 0) > 0
    return qa
