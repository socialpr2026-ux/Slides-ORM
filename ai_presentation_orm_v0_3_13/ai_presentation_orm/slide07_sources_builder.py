from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .chart_style import CHART_COMPACT_SIZE, CHART_FONT
from .native_chart_helpers import chart_update, replace_chart_data, replace_pie_chart_data
from .pie_label_renderer import render_pie_base
from .slide04_problem_field_builder import EMU_PER_INCH, _format_source_period
from .style_rules import BODY_FONT_SIZE, CHART_FONT_SIZE, FONT_FAMILY, SOURCE_FONT_SIZE, TABLE_FONT_SIZE, TITLE_FONT_SIZE, format_int_spaces, qa_slide_copy_principles, sort_other_last
from .text_contracts import slide07_sources_insight, qa_text_contract


TEXT_COLOR = RGBColor(0, 0, 0)
CAPTION_COLOR = RGBColor(89, 89, 89)
BLUE = RGBColor(1, 162, 255)
DARK_BLUE = "#0070C0"
LIGHT_BLUE = "#1F9BEA"
PALETTE = {
    "Форумы": "#7030A0",
    "Соцсети": "#00B0F0",
    "Мессенджеры": "#D9D9D9",
    "Блоги": "#FFC000",
    "Отзывы": "#4472C4",
    "Видео": "#ED7D31",
    "Другое": "#70AD47",
}


def _fmt_int(value: int | float) -> str:
    return format_int_spaces(value)


def _fmt_pct(value: float, digits: int = 1) -> str:
    return f"{float(value) * 100:.{digits}f}%".replace(".", ",")


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)
    parent.append(shape._element)


def _set_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT, color: RGBColor = TEXT_COLOR):
    shape.text = str(text)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def _insight_highlights(data: dict) -> list[str]:
    highlights = ["соцсети", "мессенджеры", "отзывы"]
    source_rows = [row for row in data.get("source_rows", []) if row.get("label") != "Другие"]
    highlights.extend(str(row.get("label", "")).strip() for row in source_rows[:6])
    return [item for item in highlights if item]


def _add_rich_runs(paragraph, text: str, *, highlights: list[str], font_size: float):
    lower_text = text.lower()
    terms = sorted({term for term in highlights if term}, key=len, reverse=True)
    cursor = 0
    while cursor < len(text):
        match_pos = len(text)
        match_term = ""
        for term in terms:
            pos = lower_text.find(term.lower(), cursor)
            if pos != -1 and pos < match_pos:
                match_pos = pos
                match_term = term
        if not match_term:
            run = paragraph.add_run()
            run.text = text[cursor:]
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = TEXT_COLOR
            break
        if match_pos > cursor:
            run = paragraph.add_run()
            run.text = text[cursor:match_pos]
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = TEXT_COLOR
        run = paragraph.add_run()
        match_end = match_pos + len(match_term)
        while match_end < len(text) and text[match_end].isspace():
            match_end += 1
        run.text = text[match_pos:match_end]
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = TEXT_COLOR
        cursor = match_end


def _set_insight(shape, text: str, *, highlights: list[str] | None = None):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    paragraphs = [item.strip() for item in str(text or "").splitlines() if item.strip()]
    if not paragraphs:
        paragraphs = ["Данные по источникам не найдены."]
    for idx, paragraph_text in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(6 if idx < len(paragraphs) - 1 else 0)
        paragraph.line_spacing = 1.0
        _add_rich_runs(paragraph, paragraph_text, highlights=highlights or [], font_size=BODY_FONT_SIZE)


def _source_note(data: dict) -> str:
    source_label = str(data.get("source_label", "сервис мониторинга Brand Analytics")).replace("сервис мониторинга ", "")
    return (
        f"Источник: {source_label}, {_format_source_period(data.get('period_raw'))}. "
        "Исключен спам и нерелевантные упоминания."
    )


def _short_label(value: str, limit: int = 18) -> str:
    text = str(value or "")
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _donut_items(data: dict) -> list[dict]:
    rows = sort_other_last(data.get("platform_type_rows") or [], label_key="type", value_key="messages")[:6]
    if not rows:
        rows = [{"type": "нет данных", "messages": 1, "share": 1}]
    total = sum(int(row.get("messages", 0) or 0) for row in rows) or 1
    items = []
    for row in rows:
        label = str(row.get("type", ""))
        value = int(row.get("messages", 0) or 0)
        items.append({
            "label": _short_label(label, 18),
            "value": value,
            "share": value / total if total else 0,
            "color": PALETTE.get(label, "#BFBFBF"),
        })
    return items


def _render_type_donut(data: dict, out_path: Path, *, width_in: float, height_in: float) -> dict:
    items = _donut_items(data)
    width = max(1, int(round(width_in * 300)))
    height = max(1, int(round(height_in * 300)))
    center_rel = (0.55, 0.52)
    radius = int(min(width, height) * 0.225)
    center = (int(width * center_rel[0]), int(height * center_rel[1]))
    base = render_pie_base(
        items=items,
        out_path=out_path,
        width=width,
        height=height,
        center=center,
        radius=radius,
        inner_radius=int(radius * 0.50),
        start_angle=-92,
    )
    base["items"] = items
    base["center_rel"] = center_rel
    base["radius_rel"] = (radius / width, radius / height)
    base["start_angle"] = -92
    return base


def _render_source_bars(data: dict, out_path: Path, *, width_in: float, height_in: float) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    rows = sort_other_last(data.get("source_rows") or [], label_key="label", value_key="messages")[:11]
    if not rows:
        rows = [{"label": "нет данных", "messages": 1}]
    labels = [_short_label(row.get("label", ""), 18) for row in rows]
    values = [int(row.get("messages", 0) or 0) for row in rows]
    max_v = max(values or [1])
    y = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.barh(y, values, color="#808080", height=0.52)
    ax.set_yticks(y, labels=labels)
    ax.invert_yaxis()
    ax.set_xlim(0, max_v * 1.22)
    ax.tick_params(axis="x", bottom=False, labelbottom=False)
    ax.tick_params(axis="y", length=0, labelsize=CHART_COMPACT_SIZE, colors="#000000", pad=5)
    for label in ax.get_yticklabels():
        label.set_fontname(CHART_FONT)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.grid(False)
    for idx, value in enumerate(values):
        ax.text(
            value + max_v * 0.018,
            idx,
            _fmt_int(value),
            ha="left",
            va="center",
            fontsize=CHART_COMPACT_SIZE,
            fontname=CHART_FONT,
            color="#000000",
        )
    plt.subplots_adjust(left=0.36, right=0.94, top=0.96, bottom=0.04)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _style_cell(cell, *, text: str, font_size: float = TABLE_FONT_SIZE, bold: bool = False, align=PP_ALIGN.LEFT, hyperlink: str = ""):
    cell.text = str(text)
    cell.margin_left = Pt(2)
    cell.margin_right = Pt(2)
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = TEXT_COLOR
            if hyperlink:
                run.hyperlink.address = hyperlink
                run.font.underline = False


def _ensure_communities_table_rows(slide, table_shape, *, body_rows: int = 10):
    required_rows = body_rows + 1
    tbl = table_shape.table
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    while len(tbl.rows) > required_rows:
        tr_elements = tbl._tbl.findall(f"{ns}tr")
        if not tr_elements:
            break
        tbl._tbl.remove(tr_elements[-1])
    if len(tbl.rows) >= required_rows:
        return table_shape
    left, top, width, height = table_shape.left, table_shape.top, table_shape.width, table_shape.height
    col_widths = [column.width for column in tbl.columns]
    _remove_shape(slide, table_shape)
    new_table_shape = slide.shapes.add_table(required_rows, 4, left, top, width, height)
    for idx, col_width in enumerate(col_widths[:4]):
        new_table_shape.table.columns[idx].width = col_width
    return new_table_shape


def _set_communities_table(table_shape, rows: list[dict]):
    table = table_shape.table
    headers = ["Топ сообществ-источников", "Площадка", "Сообщения", "Аудитория"]
    for col, header in enumerate(headers):
        _style_cell(table.cell(0, col), text=header, font_size=TABLE_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER)
    for idx in range(1, len(table.rows)):
        item = rows[idx - 1] if idx - 1 < len(rows) else {}
        values = [
            item.get("community", ""),
            item.get("platform", ""),
            _fmt_int(item.get("messages", 0)) if item else "",
            _fmt_int(item.get("audience", 0)) if item else "",
        ]
        for col, value in enumerate(values):
            _style_cell(
                table.cell(idx, col),
                text=value,
                font_size=TABLE_FONT_SIZE,
                bold=False,
                align=PP_ALIGN.CENTER,
                hyperlink=item.get("url", "") if col == 0 else "",
            )


def build_slide07_sources(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    slide = prs.slides[6]
    image_dir = output_path.parent / "slide07_chart_images"

    title = _shape_by_id(slide, 3)
    insight = _shape_by_id(slide, 8)
    type_caption = _shape_by_id(slide, 17)
    source_caption = _shape_by_id(slide, 5)
    source_note = _shape_by_id(slide, 10)
    donut = _shape_by_id(slide, 6)
    bars = _shape_by_id(slide, 7)
    table = _shape_by_id(slide, 2)

    if title:
        _set_text(title, "ПЛОЩАДКИ ОБСУЖДЕНИЙ", font_size=TITLE_FONT_SIZE)
    if insight:
        contract_text = slide07_sources_insight(data)
        data["insight_text"] = contract_text
        _set_insight(insight, contract_text, highlights=_insight_highlights(data))
    if type_caption:
        _set_text(type_caption, "Типы площадок", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, color=CAPTION_COLOR)
    if source_caption:
        _set_text(source_caption, "Топ источников упоминаний", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, color=CAPTION_COLOR)
    if source_note:
        _set_text(source_note, _source_note(data), font_size=SOURCE_FONT_SIZE)
    if table:
        table = _ensure_communities_table_rows(slide, table, body_rows=10)
        _set_communities_table(table, data.get("community_rows") or [])

    chart_updates = []
    image_outputs = []
    donut_label_layout = {}
    if donut:
        pos = (donut.left, donut.top, donut.width, donut.height)
        donut.left, donut.top, donut.width, donut.height = [int(value) for value in pos]
        donut_items = _donut_items(data)
        replace_pie_chart_data(
            donut,
            labels=[item["label"] for item in donut_items],
            values=[item["value"] for item in donut_items],
            colors=[item["color"] for item in donut_items],
            show_legend=False,
            label_position=XL_LABEL_POSITION.OUTSIDE_END,
            show_legend_key=True,
            doughnut=True,
        )
        chart_updates.append(chart_update("platform_types", "template_doughnut_replace_data", label_layout=donut_label_layout))
    if bars:
        pos = (bars.left, bars.top, bars.width, bars.height)
        bars.left, bars.top, bars.width, bars.height = [int(value) for value in pos]
        rows = sort_other_last(data.get("source_rows") or [], label_key="label", value_key="messages")[:11]
        labels = [_short_label(row.get("label", ""), 18) for row in rows] or ["Нет данных"]
        values = [int(row.get("messages", 0) or 0) for row in rows] or [0]
        replace_chart_data(
            bars,
            categories=labels,
            series=[("Сообщения", values, "#808080")],
            show_legend=False,
            reverse_categories=True,
        )
        chart_updates.append(chart_update("top_sources", "template_horizontal_bar_replace_data", categories=labels))

    for shape in [title, insight, type_caption, source_caption, source_note, table]:
        if shape:
            _bring_to_front(shape)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    text_after = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text.strip()
            if txt:
                text_after.append(txt)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_after.append(cell.text.strip())
    joined = "\n".join(text_after)
    forbidden = ["Пикторид", "декабрь 2025", "Sprosivracha.com", "345 сообщений", "43%"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]

    return {
        "output": str(output_path),
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "old_text_left": old_text_left,
        "source_note": source_note.text if source_note else "",
        "insight_text": data.get("insight_text", ""),
        "text_after": text_after,
        "visual_guardrails": {
            "donut_geometry": "uses inherited reference chart slot",
            "bar_geometry": "uses inherited reference chart slot",
            "table_geometry": "uses inherited reference table slot with at least 10 body rows",
            "community_table_body_rows": len(table.table.rows) - 1 if table and getattr(table, "has_table", False) else 0,
            "donut_label_layout": donut_label_layout,
        },
    }


def qa_slide07_sources(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    source_note = result.get("source_note", "")
    joined = "\n".join(result.get("text_after") or [])
    total = int(data.get("total_messages", 0) or 0)
    type_sum = sum(int(row.get("messages", 0) or 0) for row in data.get("platform_type_rows") or [])
    source_sum = sum(int(row.get("messages", 0) or 0) for row in data.get("source_rows") or [])
    if not total:
        blockers.append("Slide 07 total message count is missing.")
    if type_sum != total:
        blockers.append(f"Platform type totals mismatch: types={type_sum} total={total}.")
    if source_sum != total:
        blockers.append(f"Source row totals mismatch: sources={source_sum} total={total}.")
    if len(data.get("platform_type_rows") or []) < 4:
        blockers.append("Slide 07 needs at least four platform type rows.")
    if len(data.get("source_rows") or []) < 6:
        warnings.append("Slide 07 has fewer than six top source rows in the source data; chart uses all available rows.")
    if len(data.get("community_rows") or []) < 10:
        blockers.append("Slide 07 needs at least ten community rows.")
    table_body_rows = int(((result.get("visual_guardrails") or {}).get("community_table_body_rows") or 0))
    if table_body_rows < 10:
        blockers.append(f"Slide 07 community table must show ten rows, got {table_body_rows}.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 07: " + ", ".join(result["old_text_left"]))
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    if any(time in source_note for time in ["00:00", "23:59"]):
        blockers.append("Source note must not include hours/minutes.")
    if ";" in joined:
        blockers.append("Slide 07 text must not use semicolons.")
    if data.get("methodology_status") != "ready":
        warnings.append("Slide 07 data model is not marked ready.")
    message_filter = data.get("message_filter") or {}
    if message_filter.get("scope") == "all_messages":
        blockers.append("Slide 07 must use project-brand messages when brand/object/tag filters are available, not all messages.")
    duplicated_sources = [label for label in {str(row.get("label", "")).strip() for row in data.get("source_rows") or []} if label and sum(1 for row in data.get("source_rows") or [] if row.get("label") == label) > 1]
    if duplicated_sources:
        blockers.append("Slide 07 source labels must be deduplicated after normalization: " + ", ".join(duplicated_sources[:4]))
    label_layout = (result.get("visual_guardrails") or {}).get("donut_label_layout") or {}
    if label_layout:
        if label_layout.get("label_count") != len(data.get("platform_type_rows") or []):
            blockers.append("Platform type donut label count must match slice count.")
        if label_layout.get("overlap_count", 0):
            blockers.append(f"Platform type donut labels overlap: {label_layout.get('overlap_count')}.")
        if label_layout.get("out_of_bounds_count", 0):
            blockers.append(f"Platform type donut labels are out of bounds: {label_layout.get('out_of_bounds_count')}.")
        if label_layout.get("label_chart_overlap_count", 0):
            blockers.append(f"Platform type donut labels overlap the chart circle: {label_layout.get('label_chart_overlap_count')}.")
        if label_layout.get("color_mismatch_count", 0):
            blockers.append(f"Platform type donut label colors do not match slices: {label_layout.get('color_mismatch_count')}.")
        if label_layout.get("clipped_text_count", 0):
            blockers.append(f"Platform type donut labels may be clipped: {label_layout.get('clipped_text_count')}.")
        if label_layout.get("missing_nonzero_label_count", 0):
            blockers.append(f"Platform type donut has missing non-zero labels: {label_layout.get('missing_nonzero_label_count')}.")
    else:
        warnings.append("Platform type donut uses native PowerPoint legend labels; external label layout diagnostics are not applicable.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=7)
    text_qa = qa_text_contract(joined, slide_number=7)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "total_messages_present": bool(total),
            "platform_type_total_matches": type_sum == total,
            "source_total_matches": source_sum == total,
            "communities_present": len(data.get("community_rows") or []) >= 6,
            "no_stale_template_text": not result.get("old_text_left"),
            "no_semicolon": ";" not in joined,
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "source_note_without_time": not any(time in source_note for time in ["00:00", "23:59"]),
            "project_brand_filter_used": message_filter.get("scope") != "all_messages",
            "source_labels_deduplicated": not duplicated_sources,
            "donut_labels_match_slices": not label_layout or label_layout.get("label_count") == len(data.get("platform_type_rows") or []),
            "donut_labels_no_overlap": not label_layout or label_layout.get("overlap_count", 0) == 0,
            "donut_labels_in_bounds": not label_layout or label_layout.get("out_of_bounds_count", 0) == 0,
            "donut_labels_clear_chart": not label_layout or label_layout.get("label_chart_overlap_count", 0) == 0,
            "donut_label_colors_match": not label_layout or label_layout.get("color_mismatch_count", 0) == 0,
            "donut_labels_not_clipped": not label_layout or label_layout.get("clipped_text_count", 0) == 0,
            "donut_nonzero_labels_present": not label_layout or label_layout.get("missing_nonzero_label_count", 0) == 0,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
    }
