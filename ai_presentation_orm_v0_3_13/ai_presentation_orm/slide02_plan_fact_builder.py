from __future__ import annotations

from pathlib import Path
import re
import urllib.parse
import urllib.request
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

from .brand_styler import bold_project_brand_in_presentation
from .style_rules import TABLE_FONT_SIZE


def _set_text_preserve_style(text_container, text: str) -> None:
    tf = text_container.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
        p.runs[0].text = str(text)
        for run in p.runs[1:]:
            run.text = ""
    else:
        text_container.text = str(text)


def _fmt_value(value) -> str:
    if value in ("", None):
        return ""
    try:
        f = float(value)
        if f.is_integer():
            return f"{int(f):,}".replace(",", " ")
        return str(round(f, 2)).replace(".", ",")
    except Exception:
        return str(value)


def _set_cell_text_and_style(cell, text: str, *, font_size_pt: int = 9, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    cell.text = str(text)
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold


def _remove_template_pictures(slide) -> int:
    removed = 0
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape._element.getparent().remove(shape._element)
            removed += 1
    return removed


def _picture_slots(slide) -> list[tuple[int, int, int, int]]:
    return [
        (shape.left, shape.top, shape.width, shape.height)
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


def _fetch_apteka_product_image(brand: str, output_dir: Path) -> tuple[str, dict]:
    status = {"status": "skipped", "reason": "brand_missing"}
    clean_brand = re.sub(r"\s*[®Â®]\s*", "", str(brand or "")).strip()
    if not clean_brand:
        return "", status
    try:
        query = urllib.parse.quote(clean_brand)
        url = f"https://apteka.ru/search/?q={query}"
        request = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(request, timeout=8) as response:
            html = response.read().decode("utf-8", errors="ignore")
        image_match = re.search(r'<meta[^>]+property=["\']og:image["\'][^>]+content=["\']([^"\']+)', html, flags=re.IGNORECASE)
        if not image_match:
            image_match = re.search(r'https://[^"\']+\.(?:png|jpg|jpeg|webp)', html, flags=re.IGNORECASE)
            image_url = image_match.group(0) if image_match else ""
        else:
            image_url = image_match.group(1)
        if not image_url:
            return "", {"status": "skipped", "reason": "apteka_image_not_found", "search_url": url}
        if image_url.startswith("//"):
            image_url = "https:" + image_url
        if image_url.startswith("/"):
            image_url = "https://apteka.ru" + image_url
        suffix = Path(urllib.parse.urlparse(image_url).path).suffix.lower()
        if suffix not in {".png", ".jpg", ".jpeg", ".webp"}:
            suffix = ".jpg"
        out_dir = output_dir / "slide02_product_image"
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / f"{re.sub(r'[^A-Za-zА-Яа-я0-9]+', '_', clean_brand).strip('_')}{suffix}"
        request = urllib.request.Request(image_url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(request, timeout=8) as response:
            payload = response.read()
        if len(payload) < 1024:
            return "", {"status": "skipped", "reason": "downloaded_image_too_small", "image_url": image_url}
        out_path.write_bytes(payload)
        return str(out_path), {"status": "downloaded", "search_url": url, "image_url": image_url, "path": str(out_path)}
    except Exception as exc:
        return "", {"status": "skipped", "reason": "apteka_fetch_failed", "error": str(exc)}


def _trim_table_body_rows(tbl, body_rows_to_keep: int) -> int:
    body_capacity = len(tbl.rows) - 1
    keep = max(0, min(body_rows_to_keep, body_capacity))
    removed = 0
    for row_idx in range(body_capacity, keep, -1):
        tr = tbl._tbl.tr_lst[row_idx]
        tbl._tbl.remove(tr)
        removed += 1
    return removed


def _is_service_row(row: dict) -> bool:
    return bool(row.get("is_service_row")) or "Мониторинг" in str(row.get("name", ""))


def _fit_rows_to_body_capacity(rows: list[dict], body_capacity: int) -> list[dict]:
    if body_capacity <= 0 or len(rows) <= body_capacity:
        return rows
    service_rows = [row for row in rows if _is_service_row(row)]
    if not service_rows:
        return rows[:body_capacity]
    non_service = [row for row in rows if not _is_service_row(row)]
    return non_service[: max(body_capacity - 1, 0)] + service_rows[:1]


def build_slide02_from_media_plan(template_path: Path, output_path: Path, project: dict, media_plan_table: dict) -> dict:
    """Update slide 2 table using a media plan table model.

    v0.2.2 rules:
    - columns are previous month, current month, total;
    - third numeric header is always 'Всего';
    - total is computed as selected month 1 + selected month 2;
    - row labels are standardized for presentation readability;
    - no empty rows;
    - service row is always last;
    - last three headers are centered.
    """
    prs = Presentation(str(template_path))
    if len(prs.slides) < 2:
        raise ValueError("Template does not have slide 2")

    slide = prs.slides[1]
    period = project.get("period", "")
    brand = project.get("brand", "")
    product_image = project.get("product_image", "")

    replacements_done = 0
    pictures_removed = 0

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text or ""
            if "ПЕРИОД:" in txt.upper() and period:
                _set_text_preserve_style(shape, f"ПЕРИОД: {period.upper()}")
                replacements_done += 1
            if "ТЕКУЩИЙ ОТЧЕТ" in txt.upper():
                _set_text_preserve_style(shape, "ТЕКУЩИЙ ОТЧЕТ ОБ ИНФОРМАЦИОННОЙ КАМПАНИИ")
                replacements_done += 1

    picture_slots = _picture_slots(slide)
    image_status = {"status": "not_requested"}
    if not product_image:
        product_image, image_status = _fetch_apteka_product_image(brand, output_path.parent)
    if product_image and Path(product_image).exists() and picture_slots:
        pictures_removed = _remove_template_pictures(slide)
        slide.shapes.add_picture(str(product_image), *picture_slots[0])
    else:
        pictures_removed = _remove_template_pictures(slide)

    table_shape = next((shape for shape in slide.shapes if getattr(shape, "has_table", False)), None)
    if table_shape is None:
        raise ValueError("Slide 2 table not found")

    tbl = table_shape.table
    model = media_plan_table.get("slide_table", {})
    headers = model.get("headers", [])
    rows = model.get("rows", [])

    table_rows = len(tbl.rows)
    table_cols = len(tbl.columns)
    body_capacity = table_rows - 1
    trimmed_template_rows = 0

    if len(rows) < body_capacity:
        if not rows:
            raise ValueError("Slide 2 media plan table has no body rows")
        trimmed_template_rows = _trim_table_body_rows(tbl, len(rows))
        table_rows = len(tbl.rows)
        table_cols = len(tbl.columns)
        body_capacity = table_rows - 1
    rows = _fit_rows_to_body_capacity(rows, body_capacity)

    # Header: first column left, last three centered.
    for c in range(table_cols):
        value = headers[c] if c < len(headers) else ""
        _set_cell_text_and_style(
            tbl.cell(0, c),
            value,
            font_size_pt=TABLE_FONT_SIZE,
            bold=True,
            align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER,
        )
        replacements_done += 1

    for r in range(body_capacity):
        source_row = rows[r]
        values = [source_row.get("name", "")] + source_row.get("values", []) + [source_row.get("total", "")]
        for c in range(table_cols):
            val = values[c] if c < len(values) else ""
            if c == 0:
                _set_cell_text_and_style(tbl.cell(r + 1, c), val, font_size_pt=TABLE_FONT_SIZE, bold=False, align=PP_ALIGN.LEFT)
            else:
                _set_cell_text_and_style(
                    tbl.cell(r + 1, c),
                    _fmt_value(val),
                    font_size_pt=TABLE_FONT_SIZE,
                    bold=(c == table_cols - 1),
                    align=PP_ALIGN.CENTER,
                )
            replacements_done += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    brand_bold_updates = bold_project_brand_in_presentation(prs, brand)
    prs.save(str(output_path))

    scan_text = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            scan_text.append(shape.text)
    for r in range(table_rows):
        for c in range(table_cols):
            scan_text.append(tbl.cell(r, c).text)
    joined = "\n".join(scan_text)

    old_text_left = []
    for old in ["Всего 2025", "Всего факт", "Бактоблис", "БАКТОБЛИС", "Пикторид"]:
        if old.lower() in joined.lower() and old.lower() != brand.lower():
            old_text_left.append(old)

    empty_rows = []
    for r in range(1, table_rows):
        if not any(tbl.cell(r, c).text.strip() for c in range(table_cols)):
            empty_rows.append(r)

    return {
        "output": str(output_path),
        "slide_count": len(prs.slides),
        "brand": brand,
        "period": period,
        "replacement_count": replacements_done,
        "headers": headers,
        "rows_used": len(rows),
        "table_body_capacity": body_capacity,
        "trimmed_template_rows": trimmed_template_rows,
        "overflow_rows_count": media_plan_table.get("overflow_rows_count", 0),
        "pictures_removed": pictures_removed,
        "product_image": product_image,
        "product_image_status": image_status,
        "old_text_left": old_text_left,
        "empty_rows": empty_rows,
        "last_row": [tbl.cell(table_rows - 1, c).text for c in range(table_cols)],
        "media_plan_source": {
            "file": media_plan_table.get("source_file"),
            "sheet": media_plan_table.get("sheet"),
            "layout": "plan_fact_pairs" if media_plan_table.get("layout", {}).get("has_plan_fact_pairs") else "monthly_values",
            "total_mode": media_plan_table.get("total_mode"),
        },
        "risk_note": "No product image supplied: old template picture removed to avoid wrong-brand slide.",
        "brand_bold_updates": brand_bold_updates,
    }


def qa_slide02_plan_fact(media_plan_table: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    headers = result.get("headers") or []
    rows = (media_plan_table.get("slide_table") or {}).get("rows") or []
    if result.get("slide_count") != 12:
        blockers.append(f"Output PPTX must keep 12 slides, got {result.get('slide_count')}.")
    if len(headers) < 4 or headers[-1] != "Всего":
        blockers.append("Slide 02 table must have previous month, current month and 'Всего' columns.")
    if not rows:
        blockers.append("Slide 02 media-plan rows are missing.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 02: " + ", ".join(result["old_text_left"]))
    if result.get("empty_rows"):
        blockers.append("Slide 02 table contains empty rows: " + ", ".join(str(item) for item in result["empty_rows"]))
    last_row = result.get("last_row") or []
    if not last_row or "Мониторинг" not in str(last_row[0]):
        blockers.append("Slide 02 service row must remain last.")
    for row in rows:
        values = list(row.get("values") or [])
        numeric_values = [value for value in values if isinstance(value, (int, float))]
        if numeric_values and row.get("total") != sum(numeric_values):
            blockers.append(f"Slide 02 row total mismatch: {row.get('name')}.")
            break
    fact_overrides = media_plan_table.get("fact_overrides") or {}
    if fact_overrides and fact_overrides.get("status") not in {"applied", "not_needed"}:
        blockers.append("Slide 02 fact overrides did not apply cleanly: " + str(fact_overrides.get("reason") or fact_overrides.get("error") or fact_overrides.get("status")))
    product_image_status = result.get("product_image_status") or {}
    if product_image_status.get("status") in {"skipped", "not_requested"} and not result.get("product_image"):
        warnings.append("Slide 02 product image was not inserted; the slide keeps a manual image slot or no image to avoid wrong-brand assets.")
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "slide_count_12": result.get("slide_count") == 12,
            "headers_present": len(headers) >= 4 and headers[-1] == "Всего",
            "rows_present": bool(rows),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_empty_rows": not result.get("empty_rows"),
            "service_row_last": bool(last_row and "Мониторинг" in str(last_row[0])),
            "brand_bold_updates": result.get("brand_bold_updates", 0) > 0,
            "product_image_status": product_image_status.get("status"),
        },
    }
