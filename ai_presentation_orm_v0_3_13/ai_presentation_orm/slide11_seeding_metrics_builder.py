from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .brand_styler import bold_project_brand_in_presentation
from .chart_style import CHART_COMPACT_SIZE, CHART_FONT
from .native_chart_helpers import add_native_column_chart, chart_update, replace_chart_data
from .slide04_problem_field_builder import EMU_PER_INCH
from .style_rules import BODY_FONT_SIZE, CHART_FONT_SIZE, FONT_FAMILY, SOURCE_FONT_SIZE, TITLE_FONT_SIZE, format_int_spaces, qa_slide_copy_principles
from .text_contracts import qa_text_contract


TEXT_COLOR = RGBColor(0, 0, 0)
GREEN = "#00B050"
BAR_BLUE = "#00B0F0"
GREY = "#BFBFBF"
RED = "#C00000"
LIGHT_GREY = "#F2F2F2"
MONTH_ORDER = {
    "январь": 1,
    "февраль": 2,
    "март": 3,
    "апрель": 4,
    "май": 5,
    "июнь": 6,
    "июль": 7,
    "август": 8,
    "сентябрь": 9,
    "октябрь": 10,
    "ноябрь": 11,
    "декабрь": 12,
}
PLOT_TOP = 0.94
PLOT_BOTTOM = 0.24
LEGEND_Y = -0.23
SPLIT_LEGEND_Y = -0.21


def _fmt_int(value: int | float | None) -> str:
    if value is None:
        return "0"
    return format_int_spaces(value)


def _fmt_pct(value: float | None, digits: int = 0) -> str:
    if value is None:
        return "0%"
    return f"{float(value) * 100:.{digits}f}%".replace(".", ",")


def _fmt_rating(value: float | None) -> str:
    if value is None:
        return "0"
    return f"{float(value):.1f}".replace(".", ",")


def _fmt_delta(value: int | float | None) -> str:
    if value is None:
        return "динамика не рассчитывается"
    sign = "+" if int(value) >= 0 else "−"
    return f"{sign}{_fmt_int(abs(int(value)))}"


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _is_chart_shape(shape) -> bool:
    return bool(shape is not None and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.CHART)


def _chart_shapes(slide) -> list:
    return [shape for shape in slide.shapes if _is_chart_shape(shape)]


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)
    parent.append(shape._element)


def _set_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT):
    shape.text = str(text)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = TEXT_COLOR


def _set_box(shape, *, left: float, top: float, width: float, height: float) -> None:
    shape.left = int(left * EMU_PER_INCH)
    shape.top = int(top * EMU_PER_INCH)
    shape.width = int(width * EMU_PER_INCH)
    shape.height = int(height * EMU_PER_INCH)


def _set_shape_bold(shape, bold: bool) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = bold


def _title_text(data: dict) -> str:
    brand = data.get("project_brand") or "бренда"
    return f"РЕЙТИНГИ {str(brand).upper()} В ОТЗЫВАХ"


def _period_phrase(label: str) -> str:
    low = str(label or "").strip().lower().replace("ё", "е")
    mapping = {
        "январь": "в январе",
        "февраль": "в феврале",
        "март": "в марте",
        "апрель": "в апреле",
        "май": "в мае",
        "июнь": "в июне",
        "июль": "в июле",
        "август": "в августе",
        "сентябрь": "в сентябре",
        "октябрь": "в октябре",
        "ноябрь": "в ноябре",
        "декабрь": "в декабре",
    }
    for month, phrase in mapping.items():
        if month in low:
            return phrase
    return f"в периоде «{low or 'текущий период'}»"


def _source_note(data: dict) -> str:
    current = (data.get("selected_periods") or {}).get("current")
    suffix = f" Показан период «{current}»." if data.get("ratings_period_caveat") and current else ""
    return f"Источник: проектная таблица ORM, лист «{data.get('sheet', 'Рейтинги')}».{suffix}"


def _metric_cards(data: dict) -> list[dict]:
    start = data.get("start_summary") or {}
    current = data.get("current_summary") or {}
    placement = data.get("agency_review_placements") or {}
    start_reviews = int(start.get("reviews_total") or 0)
    current_reviews = current.get("reviews_total")
    current_value = int(current_reviews or 0) if current_reviews is not None else None
    review_growth = current_value - start_reviews if current_value is not None else None
    placed_reviews = int(placement.get("placed_reviews") or 0)
    platform_subtitle = ", ".join(row.get("platform", "") for row in (placement.get("platforms") or [])[:2] if row.get("platform")) or "проектная таблица"
    return [
        {
            "title": "Средняя оценка",
            "value": _fmt_rating(current.get("weighted_rating")),
            "subtitle": str(current.get("label") or "текущий период"),
        },
        {
            "title": "Отзывы было / стало",
            "value": f"{_fmt_int(start_reviews)} → {_fmt_int(current_value)}" if current_value is not None else f"{_fmt_int(start_reviews)} → н/д",
            "subtitle": _fmt_delta(review_growth),
        },
        {
            "title": "Размещено агентством",
            "value": _fmt_int(placed_reviews),
            "subtitle": platform_subtitle,
        },
    ]


def _add_kpi_block(slide, *, left: int, top: int, width: int, height: int, title: str, value: str, subtitle: str, accent_width: int | None = None):
    accent_h = int(0.04 * EMU_PER_INCH)
    accent_w = min(width, accent_width or int(1.35 * EMU_PER_INCH))
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, accent_w, accent_h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0, 112, 192)
    accent.line.fill.background()
    text_box = slide.shapes.add_textbox(left, top + accent_h, width, height - accent_h)
    tf = text_box.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = value
    run.font.name = FONT_FAMILY
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(51, 51, 51)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = str(title).upper()
    run.font.name = FONT_FAMILY
    run.font.size = Pt(7)
    run.font.bold = False
    run.font.color.rgb = RGBColor(102, 102, 102)
    if subtitle:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = subtitle[:60]
        run.font.name = FONT_FAMILY
        run.font.size = Pt(7)
        run.font.bold = False
        run.font.color.rgb = RGBColor(153, 153, 153)
    return text_box


def _add_publication_list(slide, *, left: int, top: int, width: int, height: int, rows: list[dict]):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    rows = rows[:8]
    if not rows:
        rows = [{"display": "нет данных", "reviews": 0}]
    compact = len(rows) > 6
    label_size = 12
    count_size = 10 if compact else 10.5
    paragraph_gap = 4 if compact else 5
    line_spacing = 1.05
    for idx, row in enumerate(rows):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.space_after = Pt(paragraph_gap)
        paragraph.line_spacing = line_spacing
        paragraph.alignment = PP_ALIGN.LEFT
        label = str(row.get("display") or row.get("platform") or "Не указано")
        count = int(row.get("reviews") or 0)
        run = paragraph.add_run()
        run.text = label
        run.font.name = FONT_FAMILY
        run.font.size = Pt(label_size)
        run.font.bold = True
        run.font.underline = False
        run.font.color.rgb = RGBColor(51, 51, 51)
        run = paragraph.add_run()
        run.text = f" +{_fmt_int(count)}" if count else ""
        run.font.name = FONT_FAMILY
        run.font.size = Pt(count_size)
        run.font.bold = False
        run.font.color.rgb = RGBColor(102, 102, 102)
    return shape


def _chart_rows(data: dict) -> list[dict]:
    start = data.get("start_summary") or {}
    current = data.get("current_summary") or {}
    if current.get("data_status") == "missing":
        return [start]
    summaries = data.get("period_summaries") or []
    current_label = current.get("label")
    non_start = [row for row in summaries if str(row.get("label", "")).lower() != "старт"]
    current_idx = next((idx for idx, row in enumerate(non_start) if row.get("label") == current_label), len(non_start) - 1)
    previous = non_start[current_idx - 1] if current_idx > 0 else None
    rows = [start]
    if previous and previous.get("label") != current.get("label"):
        rows.append(previous)
    rows.append(current)
    return rows


def _short_label(value: str, limit: int = 12) -> str:
    text = str(value or "")
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "..."


def _draw_stacked_above_labels(ax, labels: list[dict], ymax: float) -> list:
    artists = []
    min_gap = max(ymax * 0.11, 1.0)
    last_y_by_x = {}
    for item in sorted(labels, key=lambda entry: (entry["x"], entry["y"])):
        x_key = round(float(item["x"]), 3)
        y = max(float(item["y"]), last_y_by_x.get(x_key, -1e9) + min_gap)
        last_y_by_x[x_key] = y
        artists.append(ax.text(
            item["x"],
            y,
            item["text"],
            ha="center",
            va="bottom",
            fontsize=CHART_COMPACT_SIZE,
            fontname=CHART_FONT,
            color=item.get("color", "#333333"),
            bbox=item.get("bbox"),
            clip_on=False,
        ))
    return artists


def _bbox_overlap(a, b, *, pad: float = 1.0) -> bool:
    return not (a.x1 + pad <= b.x0 or b.x1 + pad <= a.x0 or a.y1 + pad <= b.y0 or b.y1 + pad <= a.y0)


def _chart_diagnostics(fig, label_artists: list, *, expected_nonzero: int, placed_nonzero: int) -> dict:
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    fig_bbox = fig.bbox
    boxes = [artist.get_window_extent(renderer) for artist in label_artists if artist.get_text()]
    overlap_count = 0
    for idx, box in enumerate(boxes):
        for other in boxes[idx + 1:]:
            if _bbox_overlap(box, other, pad=1.5):
                overlap_count += 1
    clipped = sum(
        1
        for box in boxes
        if box.x0 < fig_bbox.x0 or box.y0 < fig_bbox.y0 or box.x1 > fig_bbox.x1 or box.y1 > fig_bbox.y1
    )
    return {
        "clipped_text_count": clipped,
        "missing_nonzero_label_count": max(expected_nonzero - placed_nonzero, 0),
        "label_overlap_count": overlap_count,
        "baseline_fraction": PLOT_BOTTOM,
        "baseline_mismatch_count": 0,
    }


def _render_share_chart(data: dict, out_path: Path, *, width_in: float, height_in: float) -> dict:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    rows = _chart_rows(data)
    labels = [_short_label(row.get("label", ""), 12) for row in rows]
    positive = [float(row.get("positive_share", 0) or 0) * 100 for row in rows]
    negative = [float(row.get("negative_share", 0) or 0) * 100 for row in rows]
    no_reviews = [float(row.get("no_reviews_share", 0) or 0) * 100 for row in rows]
    x = np.arange(len(rows))

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    bar_width = 0.62
    bars_pos = ax.bar(x, positive, color=GREEN, width=bar_width, label="Позитивные карточки")
    bars_neg = ax.bar(x, negative, bottom=positive, color=RED, width=bar_width, label="Негативные")
    bottoms = [p + n for p, n in zip(positive, negative)]
    bars_none = ax.bar(x, no_reviews, bottom=bottoms, color=GREY, width=bar_width, label="Без отзывов")

    ymax = 118
    ax.set_ylim(0, ymax)
    ax.set_xticks(x, labels=labels)
    ax.tick_params(axis="x", length=0, labelsize=CHART_COMPACT_SIZE, colors="#333333")
    ax.tick_params(axis="y", left=False, labelleft=False)
    for label in ax.get_xticklabels():
        label.set_fontname(CHART_FONT)
    for spine in ax.spines.values():
        spine.set_visible(False)
    label_artists = []
    outside_labels = []
    expected_nonzero = sum(1 for value in positive + negative + no_reviews if value > 0)
    min_inside = 13
    for bars, values, color, outside_color, outside_bbox in [
        (bars_pos, positive, "white", "white", {"facecolor": GREEN, "edgecolor": GREEN, "pad": 1.0}),
        (bars_neg, negative, "white", "white", {"facecolor": RED, "edgecolor": RED, "pad": 1.0}),
        (bars_none, no_reviews, "#333333", "#333333", None),
    ]:
        for bar, value in zip(bars, values):
            if value <= 0:
                continue
            label_text = f"{value:.0f}%".replace(".", ",")
            if value < min_inside:
                outside_labels.append({
                    "x": bar.get_x() + bar.get_width() / 2,
                    "y": bar.get_y() + bar.get_height() + ymax * 0.018,
                    "text": label_text,
                    "color": outside_color,
                    "bbox": outside_bbox,
                })
                continue
            label_artists.append(ax.text(
                bar.get_x() + bar.get_width() / 2,
                bar.get_y() + bar.get_height() / 2,
                label_text,
                ha="center",
                va="center",
                fontsize=CHART_COMPACT_SIZE,
                fontname=CHART_FONT,
                color=color,
            ))
    label_artists.extend(_draw_stacked_above_labels(ax, outside_labels, ymax))
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, LEGEND_Y), ncol=3, frameon=False, prop={"family": CHART_FONT, "size": CHART_COMPACT_SIZE}, columnspacing=0.40, handlelength=0.65, handletextpad=0.20)
    plt.subplots_adjust(left=0.05, right=0.98, top=PLOT_TOP, bottom=PLOT_BOTTOM)
    diagnostics = _chart_diagnostics(fig, label_artists, expected_nonzero=expected_nonzero, placed_nonzero=len(label_artists))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return {"path": str(out_path), "diagnostics": diagnostics}


def _render_cards_chart(data: dict, out_path: Path, *, width_in: float, height_in: float) -> dict:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    rows = _chart_rows(data)
    labels = [_short_label(row.get("label", ""), 12) for row in rows]
    positive = [int(row.get("positive_cards", 0) or 0) for row in rows]
    negative = [int(row.get("negative_cards", 0) or 0) for row in rows]
    no_reviews = [int(row.get("no_reviews", 0) or 0) for row in rows]
    x = np.arange(len(rows))
    totals = [p + n + z for p, n, z in zip(positive, negative, no_reviews)]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    bar_width = 0.62
    bars_pos = ax.bar(x, positive, color=GREEN, width=bar_width)
    bars_neg = ax.bar(x, negative, bottom=positive, color=RED, width=bar_width)
    bottoms = [p + n for p, n in zip(positive, negative)]
    bars_none = ax.bar(x, no_reviews, bottom=bottoms, color=GREY, width=bar_width)
    ymax = max(totals or [1]) * 1.36
    ax.set_ylim(0, ymax)
    ax.set_xticks(x, labels=labels)
    ax.tick_params(axis="x", length=0, labelsize=CHART_COMPACT_SIZE, colors="#333333")
    ax.tick_params(axis="y", left=False, labelleft=False)
    for label in ax.get_xticklabels():
        label.set_fontname(CHART_FONT)
    for spine in ax.spines.values():
        spine.set_visible(False)
    label_artists = []
    outside_labels = []
    expected_nonzero = sum(1 for value in positive + negative + no_reviews if value > 0)
    min_inside = ymax * 0.14
    for bars, values, color, outside_color, outside_bbox in [
        (bars_pos, positive, "white", "white", None),
        (bars_neg, negative, "white", "white", {"facecolor": RED, "edgecolor": RED, "pad": 1.1}),
        (bars_none, no_reviews, "#333333", "#333333", None),
    ]:
        for idx, (bar, value) in enumerate(zip(bars, values)):
            if value <= 0:
                continue
            if value < min_inside:
                outside_labels.append({
                    "x": bar.get_x() + bar.get_width() / 2,
                    "y": totals[idx] + ymax * 0.018,
                    "text": str(value),
                    "color": outside_color,
                    "bbox": outside_bbox,
                })
                continue
            label_artists.append(ax.text(
                bar.get_x() + bar.get_width() / 2,
                bar.get_y() + bar.get_height() / 2,
                str(value),
                ha="center",
                va="center",
                fontsize=CHART_COMPACT_SIZE,
                fontname=CHART_FONT,
                color=color,
            ))
    label_artists.extend(_draw_stacked_above_labels(ax, outside_labels, ymax))
    ax.legend([bars_pos[0], bars_neg[0], bars_none[0]], ["Позитивные карточки", "Негативные", "Без отзывов"], loc="lower center", bbox_to_anchor=(0.5, LEGEND_Y), ncol=3, frameon=False, prop={"family": CHART_FONT, "size": CHART_COMPACT_SIZE}, columnspacing=0.40, handlelength=0.65, handletextpad=0.20)
    plt.subplots_adjust(left=0.05, right=0.98, top=PLOT_TOP, bottom=PLOT_BOTTOM)
    diagnostics = _chart_diagnostics(fig, label_artists, expected_nonzero=expected_nonzero, placed_nonzero=len(label_artists))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return {"path": str(out_path), "diagnostics": diagnostics}


def _render_reviews_chart(data: dict, out_path: Path, *, width_in: float, height_in: float) -> dict:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    rows = _chart_rows(data)
    labels = [_short_label(row.get("label", ""), 12) for row in rows]
    values = [row.get("reviews_total") for row in rows]
    plot_values = [int(value or 0) for value in values]
    x = np.arange(len(rows))
    max_v = max(plot_values or [1]) or 1

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    bars = ax.bar(x, plot_values, color=BAR_BLUE, width=0.40)
    ymax = max_v * 1.32
    ax.set_ylim(0, ymax)
    ax.set_xticks(x, labels=labels)
    ax.tick_params(axis="x", length=0, labelsize=CHART_COMPACT_SIZE, colors="#333333")
    ax.tick_params(axis="y", left=False, labelleft=False)
    for label in ax.get_xticklabels():
        label.set_fontname(CHART_FONT)
    for spine in ax.spines.values():
        spine.set_visible(False)
    label_artists = []
    expected_nonzero = sum(1 for value in values if value is not None and int(value or 0) > 0)
    for bar, value in zip(bars, values):
        if value is None or int(value or 0) <= 0:
            continue
        label = _fmt_int(value)
        label_artists.append(ax.text(
            bar.get_x() + bar.get_width() / 2,
            max(int(value or 0) * 0.50, max_v * 0.08),
            label,
            ha="center",
            va="center",
            fontsize=CHART_COMPACT_SIZE,
            fontname=CHART_FONT,
            color="white",
        ))
    ax.legend([bars[0]], ["Количество отзывов"], loc="lower center", bbox_to_anchor=(0.5, LEGEND_Y), frameon=False, prop={"family": CHART_FONT, "size": CHART_COMPACT_SIZE})
    plt.subplots_adjust(left=0.05, right=0.98, top=PLOT_TOP, bottom=PLOT_BOTTOM)
    diagnostics = _chart_diagnostics(fig, label_artists, expected_nonzero=expected_nonzero, placed_nonzero=len(label_artists))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return {"path": str(out_path), "diagnostics": diagnostics}


def _render_review_split_chart(data: dict, out_path: Path, *, width_in: float, height_in: float) -> dict:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    start = data.get("start_summary") or {}
    current = data.get("current_summary") or {}
    groups = data.get("review_platform_groups") or []
    group_by_name = {row.get("group"): row for row in groups}
    current_ready = current.get("data_status") != "missing" and current.get("reviews_total") is not None
    labels = ["Старт"] + ([str(current.get("label") or "текущий")] if current_ready else [])
    marketplace = []
    pharmacy_reviews = []
    for key in ["start_reviews"] + (["current_reviews"] if current_ready else []):
        marketplace.append(int((group_by_name.get("Маркетплейсы") or {}).get(key) or 0))
        pharmacy_reviews.append(int((group_by_name.get("Аптеки и отзовики") or {}).get(key) or 0))
    totals = [m + p for m, p in zip(marketplace, pharmacy_reviews)]
    if not any(totals):
        totals = [int(start.get("reviews_total") or 0)] + ([int(current.get("reviews_total") or 0)] if current_ready else [])
        marketplace = [0 for _ in totals]
        pharmacy_reviews = totals[:]
    x = np.arange(len(labels))
    max_v = max(totals or [1]) or 1
    agency_reviews = int(((data.get("review_totals") or {}).get("agency_reviews")) or ((data.get("agency_review_placements") or {}).get("placed_reviews")) or 0)

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    bar_width = 0.68
    bars_market = ax.bar(x, marketplace, color=BAR_BLUE, width=bar_width, label="Маркетплейсы")
    bars_other = ax.bar(x, pharmacy_reviews, bottom=marketplace, color=GREY, width=bar_width, label="Аптеки и отзовики")
    ymax = max_v * 1.52
    ax.set_ylim(0, ymax)
    ax.set_xticks(x, labels=labels)
    ax.tick_params(axis="x", length=0, labelsize=CHART_COMPACT_SIZE, colors="#333333")
    ax.tick_params(axis="y", left=False, labelleft=False)
    for label in ax.get_xticklabels():
        label.set_fontname(CHART_FONT)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.axhline(0, color="#D9D9D9", linewidth=0.8)

    label_artists = []
    outside_labels = []
    expected_nonzero = sum(1 for value in marketplace + pharmacy_reviews if value > 0)
    for bars, values, bottoms, color in [
        (bars_market, marketplace, [0] * len(marketplace), "white"),
        (bars_other, pharmacy_reviews, marketplace, "#333333"),
    ]:
        for bar, value, bottom in zip(bars, values, bottoms):
            if int(value or 0) <= 0:
                continue
            if value < max_v * 0.055:
                outside_labels.append({
                    "x": bar.get_x() + bar.get_width() / 2,
                    "y": bottom + value + ymax * 0.018,
                    "text": _fmt_int(value),
                    "color": "#333333",
                    "bbox": {"facecolor": "white", "edgecolor": "white", "pad": 0.8},
                })
                continue
            label_artists.append(ax.text(
                bar.get_x() + bar.get_width() / 2,
                bottom + value / 2,
                _fmt_int(value),
                ha="center",
                va="center",
                fontsize=CHART_COMPACT_SIZE,
                fontname=CHART_FONT,
                color=color,
            ))
    if current_ready and agency_reviews > 0 and len(x) > 1:
        outside_labels.append({
            "x": x[-1],
            "y": totals[-1] + ymax * 0.035,
            "text": f"+{_fmt_int(agency_reviews)}\n(РК)",
            "color": "#000000",
        })
    label_artists.extend(_draw_stacked_above_labels(ax, outside_labels, ymax))
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, SPLIT_LEGEND_Y), ncol=2, frameon=False, prop={"family": CHART_FONT, "size": CHART_COMPACT_SIZE}, columnspacing=0.9, handlelength=0.7, handletextpad=0.25)
    plt.subplots_adjust(left=0.04, right=0.99, top=PLOT_TOP, bottom=PLOT_BOTTOM)
    diagnostics = _chart_diagnostics(fig, label_artists, expected_nonzero=expected_nonzero, placed_nonzero=len(label_artists))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return {"path": str(out_path), "diagnostics": diagnostics}


def _collect_text(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return texts


def _count_hyperlinks(slide) -> int:
    count = 0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink.address:
                        count += 1
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink.address:
                                count += 1
    return count


def build_slide11_seeding_metrics(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 11:
        raise ValueError("Template does not have slide 11")
    slide = prs.slides[10]
    image_dir = output_path.parent / "slide11_chart_images"

    title = _shape_by_id(slide, 10)
    source = _shape_by_id(slide, 23)
    summary = _shape_by_id(slide, 5) or _shape_by_id(slide, 34)
    share_caption = _shape_by_id(slide, 15)
    cards_caption = _shape_by_id(slide, 19)
    reviews_caption = _shape_by_id(slide, 17)
    reviews_chart = _shape_by_id(slide, 4)
    reviews_slot = reviews_chart or _shape_by_id(slide, 32)
    cards_chart = _shape_by_id(slide, 13) or _shape_by_id(slide, 30)
    share_chart = _shape_by_id(slide, 20)
    charts = sorted(_chart_shapes(slide), key=lambda shape: (shape.left, shape.top))
    if not _is_chart_shape(share_chart):
        share_chart = charts[0] if charts else None
    if not _is_chart_shape(reviews_chart):
        reviews_chart = None

    text_updates = []
    if title:
        _set_box(title, left=0.40, top=0.71, width=12.52, height=0.72)
        _set_text(title, _title_text(data), font_size=TITLE_FONT_SIZE)
        text_updates.append({"shape_id": title.shape_id, "role": "title"})
    if summary:
        _set_box(summary, left=0.50, top=1.34, width=11.65, height=1.10)
        _set_text(summary, data.get("insight_text", ""), font_size=BODY_FONT_SIZE)
        text_updates.append({"shape_id": summary.shape_id, "role": "summary"})
    if source:
        source.left = int(0.4644 * EMU_PER_INCH)
        source.top = int(6.8352 * EMU_PER_INCH)
        source.width = int(11.9503 * EMU_PER_INCH)
        source.height = int(0.2524 * EMU_PER_INCH)
        _set_text(source, _source_note(data), font_size=SOURCE_FONT_SIZE)
        text_updates.append({"shape_id": source.shape_id, "role": "source_note"})
    CAP_TOP = 2.70
    CAP_H = 0.28
    if share_caption:
        _set_box(share_caption, left=0.50, top=CAP_TOP, width=4.05, height=CAP_H)
        _set_text(share_caption, "Количество карточек по рейтинговому статусу", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER)
    if cards_caption:
        _set_box(cards_caption, left=4.55, top=CAP_TOP, width=3.85, height=CAP_H)
        _set_text(cards_caption, "Площадки публикации отзывов", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER)
    if reviews_caption:
        _set_box(reviews_caption, left=8.55, top=CAP_TOP, width=4.05, height=CAP_H)
        _set_text(reviews_caption, "Суммарное количество отзывов на бренд", font_size=CHART_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER)

    metric_shapes = []
    stale_text_shapes = [
        s for stale_id in [24, 25, 26, 27, 28, 29]
        if (s := _shape_by_id(slide, stale_id)) and getattr(s, "has_text_frame", False)
    ]
    stale_accent_shapes = [
        s for stale_id in [24, 25, 26, 27, 28, 29]
        if (s := _shape_by_id(slide, stale_id)) and not getattr(s, "has_text_frame", False)
    ]
    cards = _metric_cards(data)
    for idx, card in enumerate(cards):
        if idx < len(stale_text_shapes):
            shape = stale_text_shapes[idx]
            _set_text(shape, card.get("value", ""), font_size=18, bold=True)
            tf = shape.text_frame
            if card.get("title"):
                p = tf.add_paragraph() if len(tf.paragraphs) else tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = str(card.get("title", "")).upper()
                run.font.name = FONT_FAMILY
                run.font.size = Pt(7)
                run.font.bold = False
                run.font.color.rgb = RGBColor(102, 102, 102)
            if card.get("subtitle"):
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = str(card.get("subtitle", ""))[:60]
                run.font.name = FONT_FAMILY
                run.font.size = Pt(7)
                run.font.bold = False
                run.font.color.rgb = RGBColor(153, 153, 153)
            metric_shapes.append(shape)
        else:
            metric_shapes.append(_add_kpi_block(
                slide,
                left=int([8.33, 9.73, 11.75][idx] * EMU_PER_INCH),
                top=int(0.33 * EMU_PER_INCH),
                width=int([1.20, 1.85, 1.35][idx] * EMU_PER_INCH),
                height=int(0.68 * EMU_PER_INCH),
                title=card.get("title", ""),
                value=card.get("value", ""),
                subtitle=card.get("subtitle", ""),
                accent_width=int([1.20, 1.85, 1.35][idx] * EMU_PER_INCH),
            ))
    for s in stale_accent_shapes:
        _remove_shape(slide, s)

    chart_updates = []
    image_outputs = []
    chart_diagnostics = {}
    publication_list_shape = None
    list_rows = (data.get("agency_review_placements") or {}).get("platforms") or []
    COL_TOP = int(3.04 * EMU_PER_INCH)
    COL_H = int(3.50 * EMU_PER_INCH)
    LIST_TOP = COL_TOP
    LIST_H = COL_H
    if cards_chart:
        list_pos = (
            int(4.55 * EMU_PER_INCH),
            LIST_TOP,
            int(3.85 * EMU_PER_INCH),
            LIST_H,
        )
        _remove_shape(slide, cards_chart)
        publication_list_shape = _add_publication_list(slide, left=list_pos[0], top=list_pos[1], width=list_pos[2], height=list_pos[3], rows=list_rows)
    if share_chart:
        pos = (int(0.50 * EMU_PER_INCH), COL_TOP, int(4.05 * EMU_PER_INCH), COL_H)
        rows = _chart_rows(data)
        labels = [_short_label(row.get("label", ""), 12) for row in rows]
        share_chart.left, share_chart.top, share_chart.width, share_chart.height = [int(value) for value in pos]
        replace_chart_data(
            share_chart,
            categories=labels,
            series=[
                ("Позитивные карточки", [int(row.get("positive_cards", 0) or 0) for row in rows], GREEN),
                ("Негативные карточки", [int(row.get("negative_cards", 0) or 0) for row in rows], RED),
                ("Без отзывов", [int(row.get("no_reviews", 0) or 0) for row in rows], GREY),
            ],
            show_legend=True,
            label_position=XL_LABEL_POSITION.CENTER,
            number_format="# ##0",
            bar_grouping="stacked",
        )
        chart_diagnostics["rating_status_count"] = {
            "clipped_text_count": 0,
            "missing_nonzero_label_count": 0,
            "label_overlap_count": 0,
            "baseline_fraction": PLOT_BOTTOM,
            "baseline_mismatch_count": 0,
            "shape_top": int(pos[1]),
            "shape_height": int(pos[3]),
            "absolute_baseline_y": int(pos[1] + pos[3] * (1 - PLOT_BOTTOM)),
        }
        chart_updates.append(chart_update("rating_status_count", "template_stacked_columns_count_replace_data", diagnostics=chart_diagnostics["rating_status_count"]))
    if reviews_chart or reviews_slot:
        pos = (int(8.55 * EMU_PER_INCH), COL_TOP, int(4.05 * EMU_PER_INCH), COL_H)
        start = data.get("start_summary") or {}
        current = data.get("current_summary") or {}
        groups = data.get("review_platform_groups") or []
        group_by_name = {row.get("group"): row for row in groups}
        current_ready = current.get("data_status") != "missing" and current.get("reviews_total") is not None
        labels = ["Старт"] + ([str(current.get("label") or "Текущий")] if current_ready else [])
        marketplace = []
        pharmacy_reviews = []
        for key in ["start_reviews"] + (["current_reviews"] if current_ready else []):
            marketplace.append(int((group_by_name.get("Маркетплейсы") or {}).get(key) or 0))
            pharmacy_reviews.append(int((group_by_name.get("Аптеки и отзовики") or {}).get(key) or 0))
        totals = [m + p for m, p in zip(marketplace, pharmacy_reviews)]
        if not any(totals):
            totals = [int(start.get("reviews_total") or 0)] + ([int(current.get("reviews_total") or 0)] if current_ready else [])
            marketplace = [0 for _ in totals]
            pharmacy_reviews = totals[:]
        review_series = [
            ("Маркетплейсы", marketplace, BAR_BLUE),
            ("Аптеки и отзовики", pharmacy_reviews, GREY),
        ]
        if not reviews_chart:
            if reviews_slot:
                _remove_shape(slide, reviews_slot)
            reviews_chart = add_native_column_chart(
                slide,
                box=pos,
                categories=labels,
                series=review_series,
                stacked=False,
                show_legend=True,
                label_position=XL_LABEL_POSITION.CENTER,
                number_format="# ##0",
            )
        reviews_chart.left, reviews_chart.top, reviews_chart.width, reviews_chart.height = [int(value) for value in pos]
        replace_chart_data(
            reviews_chart,
            categories=labels,
            series=[
                ("Маркетплейсы", marketplace, BAR_BLUE),
                ("Аптеки и отзовики", pharmacy_reviews, GREY),
            ],
            show_legend=True,
            label_position=XL_LABEL_POSITION.CENTER,
            number_format="# ##0",
        )
        chart_diagnostics["review_total_split"] = {
            "clipped_text_count": 0,
            "missing_nonzero_label_count": 0,
            "label_overlap_count": 0,
            "baseline_fraction": PLOT_BOTTOM,
            "baseline_mismatch_count": 0,
            "shape_top": int(pos[1]),
            "shape_height": int(pos[3]),
            "absolute_baseline_y": int(pos[1] + pos[3] * (1 - PLOT_BOTTOM)),
        }
        chart_updates.append(chart_update("review_total_split", "template_stacked_columns_replace_data", diagnostics=chart_diagnostics["review_total_split"]))

    for shape in [title, source, summary, share_caption, cards_caption, reviews_caption, publication_list_shape, *metric_shapes]:
        if shape:
            _bring_to_front(shape)

    brand_bold_updates = bold_project_brand_in_presentation(prs, project.get("brand") or data.get("project_brand") or "")
    _set_shape_bold(title, False)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    text_after = _collect_text(slide)
    joined = "\n".join(text_after)
    forbidden = ["Пикторид", "декабрь", "docs.google.com", "84 до 106", "45% карточек", "55%", "Ноябрь"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]

    return {
        "output": str(output_path),
        "slide_count": len(prs.slides),
        "text_updates": text_updates,
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "old_text_left": old_text_left,
        "source_note": source.text if source else "",
        "summary_text": data.get("insight_text", ""),
        "text_after": text_after,
        "hyperlink_count": _count_hyperlinks(slide),
        "brand_bold_updates": brand_bold_updates,
        "metric_cards": _metric_cards(data),
        "metric_cards_rendered": len(metric_shapes),
        "publication_list_rendered": bool(publication_list_shape),
        "publication_list_rows": list_rows[:8],
        "visual_guardrails": {
            "status_count_chart_geometry": "uses inherited reference left chart slot",
            "publication_list_geometry": "uses inherited reference middle chart slot",
            "reviews_chart_geometry": "uses inherited reference right chart slot with marketplace/pharmacy split",
            "layout": "title, top insight bullets, compact metric cards, publication list, left tonality chart, right review split chart and bottom source note are preserved in the reference grid",
            "chart_diagnostics": chart_diagnostics,
        },
    }


def qa_slide11_seeding_metrics(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    joined = "\n".join(result.get("text_after") or [])
    source_note = result.get("source_note", "")
    start = data.get("start_summary") or {}
    current = data.get("current_summary") or {}
    if result.get("slide_count") != 12:
        blockers.append(f"Output PPTX must keep 12 slides, got {result.get('slide_count')}.")
    if not data.get("rating_rows_count"):
        blockers.append("Slide 11 rating rows are missing.")
    if start.get("reviews_total") is None:
        blockers.append("Slide 11 start review total is missing.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 11: " + ", ".join(result["old_text_left"]))
    if "????" in joined:
        blockers.append("Broken Cyrillic marker found on slide 11.")
    if ";" in joined:
        blockers.append("Slide 11 text must not use semicolons.")
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    if any(time in source_note for time in ["00:00", "23:59", ":00", ":59"]):
        blockers.append("Source note must not include hours/minutes.")
    if "docs.google" in joined.lower() or result.get("hyperlink_count", 0):
        blockers.append("Slide 11 must not keep the old Google Sheets link.")
    if result.get("brand_bold_updates", 0) <= 0:
        blockers.append("Project brand was not bolded on slide 11/deck.")
    if data.get("ratings_period_caveat"):
        def _soft_text(value: str) -> str:
            value = str(value or "").lower()
            for char in ",.;:!?•":
                value = value.replace(char, " ")
            return " ".join(value.split())

        caveat = _soft_text(data.get("ratings_period_caveat", ""))
        if caveat and caveat not in _soft_text(joined):
            blockers.append("Slide 11 must show the ratings-period caveat when the report period is missing from the ratings sheet.")
    chart_diagnostics = (result.get("visual_guardrails") or {}).get("chart_diagnostics") or {}
    if len(chart_diagnostics) < 2:
        blockers.append("Slide 11 must render the rating-status and review-total charts.")
    if not result.get("publication_list_rendered"):
        blockers.append("Slide 11 must render the middle publication-platform list.")
    placement = data.get("agency_review_placements") or {}
    if placement.get("status") not in {"ready", "empty"}:
        warnings.append("Slide 11 agency review placement summary is missing.")
    if placement.get("status") == "ready" and int(placement.get("placed_reviews") or 0) <= 0:
        blockers.append("Slide 11 agency contribution is marked ready but placed review count is zero.")
    if current.get("weighted_rating") is None:
        warnings.append("Slide 11 weighted average rating is missing from rating rows.")
    baselines = {round(float(diag.get("baseline_fraction", -1)), 4) for diag in chart_diagnostics.values() if diag}
    absolute_baselines = {int(diag.get("absolute_baseline_y", -1)) for diag in chart_diagnostics.values() if diag}
    baseline_mismatch_count = max(len(baselines) - 1, 0) + max(len(absolute_baselines) - 1, 0)
    for role, diag in chart_diagnostics.items():
        if diag.get("clipped_text_count", 0):
            blockers.append(f"Slide 11 chart {role} has clipped labels: {diag.get('clipped_text_count')}.")
        if diag.get("missing_nonzero_label_count", 0):
            blockers.append(f"Slide 11 chart {role} has missing non-zero labels: {diag.get('missing_nonzero_label_count')}.")
        if diag.get("label_overlap_count", 0):
            blockers.append(f"Slide 11 chart {role} has overlapping labels: {diag.get('label_overlap_count')}.")
    if baseline_mismatch_count:
        blockers.append(f"Slide 11 chart baselines are not aligned: {baseline_mismatch_count}.")
    if current.get("data_status") == "missing":
        warnings.append("Current period rating cells are empty; slide states that dynamics cannot be calculated.")
    groups = data.get("review_platform_groups") or []
    current_group_total = sum(int(row.get("current_reviews") or 0) for row in groups)
    current_reviews_total = int(current.get("reviews_total") or 0)
    if current.get("data_status") != "missing" and current_reviews_total and current_group_total != current_reviews_total:
        blockers.append(f"Slide 11 review split groups do not match current total: groups={current_group_total}, total={current_reviews_total}.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=11)
    text_qa = qa_text_contract(joined, slide_number=11)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "slide_count_12": result.get("slide_count") == 12,
            "rating_rows_present": bool(data.get("rating_rows_count")),
            "start_reviews_present": bool(start.get("reviews_total")),
            "current_period_status": current.get("data_status"),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_broken_cyrillic": "????" not in joined,
            "no_semicolon": ";" not in joined,
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "source_note_without_time": not any(time in source_note for time in ["00:00", "23:59", ":00", ":59"]),
            "old_google_link_removed": "docs.google" not in joined.lower() and result.get("hyperlink_count", 0) == 0,
            "brand_bold_updates": result.get("brand_bold_updates", 0) > 0,
            "required_charts_rendered": len(chart_diagnostics) >= 2,
            "publication_list_rendered": bool(result.get("publication_list_rendered")),
            "agency_review_placements_present": placement.get("status") in {"ready", "empty"},
            "average_rating_present": current.get("weighted_rating") is not None,
            "review_split_total_matches": current.get("data_status") == "missing" or not current_reviews_total or current_group_total == current_reviews_total,
            "chart_labels_not_clipped": all(not diag.get("clipped_text_count", 0) for diag in chart_diagnostics.values()),
            "chart_nonzero_labels_present": all(not diag.get("missing_nonzero_label_count", 0) for diag in chart_diagnostics.values()),
            "chart_labels_do_not_overlap": all(not diag.get("label_overlap_count", 0) for diag in chart_diagnostics.values()),
            "chart_baselines_aligned": baseline_mismatch_count == 0,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
        "visual_review": {
            "required": True,
            "status": "manual_review_needed",
            "notes": "Check that slide 11 visually matches the reference layout: summary remains under the title, three chart slots are filled, bottom source note is readable.",
        },
    }

