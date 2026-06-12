from __future__ import annotations

import re

from pptx.util import Pt


def _iter_text_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape
        if getattr(shape, "shape_type", None) and str(shape.shape_type) == "GROUP (6)":
            yield from _iter_text_shapes(shape.shapes)


def _count_brand_in_paragraph(paragraph, pattern: re.Pattern) -> int:
    return sum(len(pattern.findall(run.text or "")) for run in paragraph.runs)


def _copy_font(source, target) -> None:
    target.font.name = source.font.name
    target.font.size = source.font.size
    target.font.bold = source.font.bold
    target.font.italic = source.font.italic
    target.font.underline = source.font.underline
    try:
        if source.font.color and source.font.color.rgb:
            target.font.color.rgb = source.font.color.rgb
    except Exception:
        pass


def _set_registered_superscript(run) -> None:
    r_pr = run._r.get_or_add_rPr()
    r_pr.set("baseline", "30000")


def _split_tokens(text: str, brand_pattern: re.Pattern) -> list[tuple[str, bool, bool]]:
    parts: list[tuple[str, bool, bool]] = []
    cursor = 0
    pattern = re.compile(f"({brand_pattern.pattern}|Â®|®)", flags=re.IGNORECASE)
    for match in pattern.finditer(text):
        if match.start() > cursor:
            parts.append((text[cursor:match.start()], False, False))
        value = match.group(0)
        parts.append(("®" if value in {"Â®", "®"} else value, bool(brand_pattern.fullmatch(value)), value in {"Â®", "®"}))
        cursor = match.end()
    if cursor < len(text):
        parts.append((text[cursor:], False, False))
    return [part for part in parts if part[0]]


def _style_paragraph(paragraph, brand_pattern: re.Pattern) -> int:
    source_runs = list(paragraph.runs)
    if not source_runs:
        return 0
    mentions = sum(len(brand_pattern.findall(run.text or "")) for run in source_runs)
    has_registered = any("®" in (run.text or "") or "Â®" in (run.text or "") for run in source_runs)
    if not mentions and not has_registered:
        return 0
    payload = [(run, _split_tokens(run.text or "", brand_pattern)) for run in source_runs]
    paragraph.clear()
    for source_run, tokens in payload:
        for text, is_brand, is_registered in tokens:
            run = paragraph.add_run()
            run.text = text
            _copy_font(source_run, run)
            if is_brand:
                run.font.bold = True
            if is_registered:
                _set_registered_superscript(run)
    return mentions


def bold_project_brand_in_presentation(prs, brand: str) -> int:
    """Bold project-brand mentions and raise registered marks in editable text."""
    clean_brand = str(brand or "").strip()
    if not clean_brand:
        return 0
    pattern = re.compile(re.escape(clean_brand), flags=re.IGNORECASE)
    changes = 0
    for slide in prs.slides:
        for shape in _iter_text_shapes(slide.shapes):
            for paragraph in shape.text_frame.paragraphs:
                changes += _style_paragraph(paragraph, pattern)
    return changes
