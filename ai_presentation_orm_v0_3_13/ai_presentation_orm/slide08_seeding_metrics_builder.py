from __future__ import annotations

from pathlib import Path
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .brand_styler import bold_project_brand_in_presentation
from .chart_style import CHART_COMPACT_SIZE, CHART_FONT, CHART_LABEL_SIZE
from .native_chart_helpers import chart_update, replace_chart_data
from .slide04_problem_field_builder import EMU_PER_INCH, _format_source_period
from .style_rules import BODY_FONT_SIZE, CHART_FONT_SIZE, FONT_FAMILY, SOURCE_FONT_SIZE, TABLE_FONT_SIZE, TITLE_FONT_SIZE, format_int_spaces, qa_slide_copy_principles
from .text_contracts import slide08_seeding_summary, qa_text_contract


TEXT_COLOR = RGBColor(0, 0, 0)
LINE_BLUE = RGBColor(0, 176, 240)
WHITE = RGBColor(255, 255, 255)
BLUE = "#0070C0"
LIGHT_BLUE = "#1F9BEA"
GREY = "#D9D9D9"


def _fmt_int(value: int | float) -> str:
    return format_int_spaces(value)


def _fmt_pct(value: float, digits: int = 1) -> str:
    return f"{float(value) * 100:.{digits}f}%".replace(".", ",")


def _message_word(value: int | float) -> str:
    number = abs(int(round(float(value or 0))))
    if 11 <= number % 100 <= 14:
        return "сообщений"
    last = number % 10
    if last == 1:
        return "сообщение"
    if 2 <= last <= 4:
        return "сообщения"
    return "сообщений"


def _month_sentence_start(data: dict) -> str:
    month = str(data.get("month_prepositional") or "отчетном периоде").strip()
    if not month:
        return "В отчетном периоде"
    if month.lower().startswith(("в ", "во ")):
        return month[:1].upper() + month[1:]
    return f"В {month}"


def _materials_word(value: int | float) -> str:
    number = abs(int(round(float(value))))
    if 11 <= number % 100 <= 14:
        return "материалов"
    last = number % 10
    if last == 1:
        return "материал"
    if 2 <= last <= 4:
        return "материала"
    return "материалов"


def _shape_by_id(slide, shape_id: int):
    return next((shape for shape in slide.shapes if shape.shape_id == shape_id), None)


def _find_summary_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip().lower()
        if not text or text.startswith("источник:") or "показатели посева" in text:
            continue
        if (
            "комментар" in text
            or "отзыв" in text
            or "просмотр" in text
            or len(text) > 120
        ):
            if shape.top < int(2.90 * EMU_PER_INCH):
                candidates.append(shape)
    return sorted(candidates, key=lambda shape: (shape.top, shape.left))[0] if candidates else None


def _find_source_shape(slide):
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and (shape.text or "").strip().lower().startswith("источник:")
        and shape.top > int(5.70 * EMU_PER_INCH)
    ]
    return sorted(candidates, key=lambda shape: shape.shape_id)[0] if candidates else None


def _remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def _bring_to_front(shape):
    parent = shape._element.getparent()
    if parent is None:
        return
    parent.remove(shape._element)
    parent.append(shape._element)


def _set_text(shape, text: str, *, font_size: float, bold: bool = False, align=PP_ALIGN.LEFT, color: RGBColor = TEXT_COLOR):
    shape.text = str(text)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def _add_run(paragraph, text: str, *, bold: bool = False, color: RGBColor = TEXT_COLOR):
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = FONT_FAMILY
    run.font.size = Pt(BODY_FONT_SIZE)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def _seeding_focus_text(data: dict, project: dict) -> str:
    for key in ("seeding_focus", "category_focus", "category_context"):
        value = str(data.get(key) or project.get(key) or "").strip()
        if value:
            return value
    category = str(project.get("category") or data.get("category") or "").strip()
    if category:
        return category
    return "ключевыми контекстами категории"


def _set_summary(shape, data: dict, project: dict):
    total = int(data.get("total_materials", 0) or 0)
    metrics = data.get("views_metrics", {}) or {}
    views = int(metrics.get("selected_views_total", metrics.get("received_views_total", 0)) or 0)
    focus = _seeding_focus_text(data, project)
    month = _month_sentence_start(data)
    bullet_rows = [
        [
            (f"{_fmt_int(total)} {_message_word(total)}", True),
            (
                " разместили авторы агентства в активных обсуждениях в соцсетях, форумах, чатах и топе выдачи "
                f"по запросам тематики. {month} наибольший упор в посеве был сделан на сообщества и контент, "
                f"связанные с обсуждением {focus}.",
                False,
            ),
        ],
        [
            (
                f"Комментарии получили {_fmt_int(views)} просмотров."
                if views
                else "Просмотры по заполненным полям не зафиксированы.",
                True,
            )
        ],
    ]
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for idx, runs in enumerate(bullet_rows):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10 if idx == 0 else 0)
        paragraph.line_spacing = 1.0
        _add_run(paragraph, "•  ")
        for text, bold in runs:
            _add_run(paragraph, text, bold=bold)


def _set_chart_title(chart_shape, text: str) -> None:
    try:
        chart_shape.chart.has_title = True
        tf = chart_shape.chart.chart_title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = str(text)
        run.font.name = FONT_FAMILY
        run.font.size = Pt(CHART_FONT_SIZE)
        run.font.bold = True
        run.font.color.rgb = RGBColor(102, 102, 102)
    except Exception:
        pass


def _source_note(data: dict) -> str:
    period = _format_source_period(data.get("period_raw"))
    return f"Источник: проектная таблица ORM, лист «{data.get('sheet', 'Май')}», {period}."


def _title_text(data: dict) -> str:
    return "ПОКАЗАТЕЛИ ПОСЕВА СООБЩЕНИЙ"


def _summary_text(data: dict) -> str:
    total = int(data.get("total_materials", 0) or 0)
    metrics = data.get("views_metrics", {}) or {}
    views = int(metrics.get("selected_views_total", metrics.get("received_views_total", 0)) or 0)
    return (
        f"• {_fmt_int(total)} {_message_word(total)} разместили авторы агентства в активных обсуждениях.\n"
        f"• Комментарии получили {_fmt_int(views)} просмотров." if views else
        f"• {_fmt_int(total)} {_message_word(total)} разместили авторы агентства в активных обсуждениях.\n"
        "• Просмотры по заполненным полям не зафиксированы."
    )


def _short_label(value: str, limit: int = 28) -> str:
    text = str(value or "")
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _render_views_chart(data: dict, out_path: Path, *, width_in: float, height_in: float) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    views_available = bool(data.get("views_metrics", {}).get("has_views"))
    chart_rows = list(data.get("views_chart", {}).get("rows") or [])
    if not chart_rows:
        chart_rows = [{"label": "Период", "value": int(data.get("views_metrics", {}).get("selected_views_total", data.get("views_metrics", {}).get("received_views_total", 0)) or 0)}]
    values = [int(row.get("value", 0) or 0) for row in chart_rows] if views_available else [0 for _row in chart_rows]
    labels = [str(row.get("label") or "Период") for row in chart_rows]
    display_max = max(values + [1])
    if display_max <= 1:
        display_max = 1
    x = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=300)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    bars = ax.bar(x, values, color="#13AEE2", width=0.32)
    ax.set_xticks(x, labels=labels)
    ax.set_ylim(0, display_max * 1.35)
    ax.tick_params(axis="x", length=0, labelsize=CHART_LABEL_SIZE, colors="#404040")
    ax.tick_params(axis="y", left=False, labelleft=False)
    for label in ax.get_xticklabels():
        label.set_fontname(CHART_FONT)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.spines["bottom"].set_visible(True)
    ax.spines["bottom"].set_color("#D9D9D9")
    ax.spines["bottom"].set_linewidth(0.8)
    if views_available:
        for bar, value in zip(bars, values):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                max(value + display_max * 0.035, display_max * 0.035),
                _fmt_int(value),
                va="bottom",
                ha="center",
                fontsize=CHART_LABEL_SIZE,
                fontname=CHART_FONT,
                color="#404040",
            )
    else:
        ax.text(
            0.5,
            0.52,
            "Данные по просмотрам\nне заполнены",
            transform=ax.transAxes,
            va="center",
            ha="center",
            fontsize=CHART_LABEL_SIZE,
            fontname=CHART_FONT,
            color="#6F6F6F",
        )
    fig.text(
        0.50,
        0.90,
        str(data.get("views_metrics", {}).get("metric_label") or "Количество просмотров"),
        ha="center",
        va="top",
        fontsize=CHART_LABEL_SIZE,
        fontweight="bold",
        fontname=CHART_FONT,
        color="#6F6F6F",
    )
    if views_available:
        ax.legend([bars[-1]], [str(data.get("views_metrics", {}).get("metric_label") or "Количество прочтений")], loc="lower center", bbox_to_anchor=(0.5, -0.30), frameon=False, prop={"family": CHART_FONT, "size": CHART_COMPACT_SIZE})
    plt.subplots_adjust(left=0.06, right=0.96, top=0.74, bottom=0.25)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, transparent=False, facecolor="white", dpi=300, pad_inches=0)
    plt.close(fig)
    return str(out_path)


def _style_cell(cell, *, text: str, font_size: float = 9.0, bold: bool = False, align=PP_ALIGN.LEFT, header: bool = False, hyperlink: str = ""):
    cell.text = str(text)
    cell.margin_left = Pt(2)
    cell.margin_right = Pt(2)
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    try:
        for side in [cell._tc.tcPr.lnL, cell._tc.tcPr.lnR, cell._tc.tcPr.lnT, cell._tc.tcPr.lnB]:
            side.solidFill = LINE_BLUE
    except Exception:
        pass
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = TEXT_COLOR
            if hyperlink:
                run.hyperlink.address = hyperlink
                run.font.underline = False


def _set_platform_table(table_shape, data: dict):
    table = table_shape.table
    headers = ["Площадка", "Сообщения", "Просмотры"]
    for col, header in enumerate(headers):
        _style_cell(table.cell(0, col), text=header, font_size=TABLE_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, header=True)
    rows = list(data.get("platform_rows") or [])
    views_available = bool(data.get("views_metrics", {}).get("has_views"))
    for idx in range(1, len(table.rows)):
        item = rows[idx - 1] if idx - 1 < len(rows) else {}
        values = [
            item.get("platform", ""),
            _fmt_int(item.get("materials", 0)) if item else "",
            _fmt_int(item.get("received_views", 0)) if item and views_available else ("н/д" if item else ""),
        ]
        for col, value in enumerate(values):
            _style_cell(
                table.cell(idx, col),
                text=value,
                font_size=TABLE_FONT_SIZE,
                bold=False,
                align=PP_ALIGN.CENTER,
                hyperlink=item.get("sample_link", "") if col == 0 and item else "",
            )


def _collect_text(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


def _count_hyperlinks(slide) -> int:
    count = 0
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink.address:
                                count += 1
    return count


def build_slide08_seeding_metrics(template_path: Path, output_path: Path, project: dict, data: dict) -> dict:
    prs = Presentation(str(template_path))
    if len(prs.slides) < 8:
        raise ValueError("Template does not have slide 8")
    slide = prs.slides[7]
    image_dir = output_path.parent / "slide08_chart_images"

    title = _shape_by_id(slide, 10)
    summary = _shape_by_id(slide, 25) or _find_summary_shape(slide)
    chart = next((shape for shape in slide.shapes if getattr(shape, "has_chart", False)), None)
    table = next((shape for shape in slide.shapes if getattr(shape, "has_table", False)), None)

    text_updates = []
    if title:
        _set_text(title, _title_text(data), font_size=TITLE_FONT_SIZE)
        text_updates.append({"shape_id": title.shape_id, "role": "title"})
    if summary:
        _set_summary(summary, data, project)
        text_updates.append({"shape_id": summary.shape_id, "role": "summary"})
    if table:
        _set_platform_table(table, data)
        text_updates.append({"shape_id": table.shape_id, "role": "platform_table"})

    chart_updates = []
    image_outputs = []
    if chart:
        pos = (chart.left, chart.top, chart.width, chart.height)
        chart.left, chart.top, chart.width, chart.height = [int(value) for value in pos]
        views_available = bool(data.get("views_metrics", {}).get("has_views"))
        chart_rows = list(data.get("views_chart", {}).get("rows") or [])
        if not chart_rows:
            chart_rows = [{"label": "Период", "value": int(data.get("views_metrics", {}).get("selected_views_total", data.get("views_metrics", {}).get("received_views_total", 0)) or 0)}]
        labels = [str(row.get("label") or "Период") for row in chart_rows]
        values = [int(row.get("value", 0) or 0) for row in chart_rows] if views_available else [0 for _row in chart_rows]
        if len(chart_rows) == 1 and views_available:
            series = [
                ("Сообщения", [int(data.get("total_materials", 0) or 0)], "#808080"),
                ("Просмотры", values, "#13AEE2"),
            ]
            chart_type = "template_clustered_columns_messages_views"
        else:
            series = [(str(data.get("views_metrics", {}).get("metric_label") or "Просмотры"), values, "#13AEE2")]
            chart_type = "template_clustered_columns_views"
        replace_chart_data(
            chart,
            categories=labels,
            series=series,
            show_legend=views_available,
            label_position=XL_LABEL_POSITION.OUTSIDE_END,
            number_format="# ##0",
        )
        _set_chart_title(chart, str(data.get("views_metrics", {}).get("metric_label") or "Просмотры"))
        if not views_available:
            notice = slide.shapes.add_textbox(pos[0], pos[1] + int(pos[3] * 0.40), pos[2], int(0.45 * EMU_PER_INCH))
            _set_text(notice, "Данные по просмотрам не заполнены", font_size=CHART_FONT_SIZE, align=PP_ALIGN.CENTER, color=RGBColor(111, 111, 111))
        chart_updates.append(chart_update("views", chart_type, categories=labels, series=[item[0] for item in series]))

    source_box = None
    if table:
        existing_source = _find_source_shape(slide)
        if existing_source:
            _set_text(existing_source, _source_note(data), font_size=SOURCE_FONT_SIZE)
            source_box = existing_source
        else:
            left = int(0.4644 * EMU_PER_INCH)
            top = int(6.8352 * EMU_PER_INCH)
            width = int(11.9503 * EMU_PER_INCH)
            height = int(0.2524 * EMU_PER_INCH)
            source_box = slide.shapes.add_textbox(left, top, width, height)
            _set_text(source_box, _source_note(data), font_size=SOURCE_FONT_SIZE)
        text_updates.append({"shape_id": source_box.shape_id, "role": "source_note"})

    for shape in [title, summary, table, source_box]:
        if shape:
            _bring_to_front(shape)

    brand_bold_updates = bold_project_brand_in_presentation(prs, project.get("brand") or data.get("project_brand") or "")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    text_after = _collect_text(slide)
    joined = "\n".join(text_after)
    forbidden = ["Пикторид", "декабрь", "714 248", "170 сообщений", "Ноябрь"]
    old_text_left = [item for item in forbidden if item.lower() in joined.lower()]

    return {
        "output": str(output_path),
        "slide_count": len(prs.slides),
        "text_updates": text_updates,
        "chart_updates": chart_updates,
        "chart_image_outputs": image_outputs,
        "old_text_left": old_text_left,
        "source_note": source_box.text if source_box else "",
        "summary_text": _summary_text(data),
        "text_after": text_after,
        "hyperlink_count": _count_hyperlinks(slide),
        "brand_bold_updates": brand_bold_updates,
        "visual_guardrails": {
            "chart_geometry": "uses inherited reference chart slot",
            "table_geometry": "uses inherited reference table slot",
            "layout": "title, summary, left chart and right table are preserved from reference slide 8",
        },
    }


def qa_slide08_seeding_metrics(data: dict, result: dict) -> dict:
    blockers = []
    warnings = []
    joined = "\n".join(result.get("text_after") or [])
    source_note = result.get("source_note", "")
    total = int(data.get("total_materials", 0) or 0)
    platform_total = sum(int(row.get("materials", 0) or 0) for row in data.get("platform_rows") or [])
    type_total = sum(int(row.get("materials", 0) or 0) for row in data.get("material_type_rows") or [])
    if result.get("slide_count") != 12:
        blockers.append(f"Output PPTX must keep 12 slides, got {result.get('slide_count')}.")
    if not total:
        blockers.append("Slide 08 material count is missing.")
    if type_total != total:
        blockers.append(f"Material type totals mismatch: types={type_total} total={total}.")
    platform_data_total = int(data.get("tabular_rows_with_platform_data", 0) or 0)
    if platform_total != total:
        if platform_data_total and platform_total == platform_data_total and platform_total < total:
            warnings.append(
                f"Platform rows cover {platform_total} materials with platform/link data; authoritative publication fact is {total}."
            )
        else:
            blockers.append(f"Platform totals mismatch: platforms={platform_total} total={total}.")
    if result.get("old_text_left"):
        blockers.append("Stale template text remains on slide 08: " + ", ".join(result["old_text_left"]))
    if "????" in joined:
        blockers.append("Broken Cyrillic marker found on slide 08.")
    if ";" in joined:
        blockers.append("Slide 08 text must not use semicolons.")
    if ".xlsx" in source_note.lower() or ".pptx" in source_note.lower():
        blockers.append("Source note must not expose local filenames.")
    if any(time in source_note for time in ["00:00", "23:59", ":00", ":59"]):
        blockers.append("Source note must not include hours/minutes.")
    sample_link_count = sum(1 for row in data.get("platform_rows") or [] if row.get("sample_link"))
    if sample_link_count and result.get("hyperlink_count", 0) <= 0:
        blockers.append("Slide 08 must preserve useful source links for platform rows.")
    views_metrics = data.get("views_metrics", {}) or {}
    views_chart = data.get("views_chart", {}) or {}
    if views_metrics.get("metric_label") != "Просмотры" or views_chart.get("metric_label") != "Просмотры":
        blockers.append("Slide 08 views metric label must be 'Просмотры'.")
    if views_metrics.get("has_views"):
        rows = data.get("platform_rows") or []
        view_values = [int(row.get("received_views", 0) or 0) for row in rows if row.get("platform") != "Другие"]
        if view_values != sorted(view_values, reverse=True):
            blockers.append("Slide 08 platform rows must be sorted by views when views are available.")
    if not data.get("views_metrics", {}).get("has_views"):
        warnings.append("Views fields are empty in the project table; slide marks views as unavailable instead of estimating them.")
    if views_metrics.get("has_views"):
        chart_rows = list(views_chart.get("rows") or [])
        chart_total = sum(int(row.get("value", 0) or 0) for row in chart_rows) if chart_rows else 0
        text_views = int(views_metrics.get("selected_views_total", views_metrics.get("received_views_total", 0)) or 0)
        if chart_total > 0 and text_views > 0 and chart_total != text_views:
            blockers.append(
                f"Slide 08 text says {text_views} views, but chart row sum is {chart_total}."
            )
    if data.get("views_metrics", {}).get("metric_kind") != "received":
        if not data.get("views_metrics", {}).get("metric_caveat") and not data.get("views_metrics", {}).get("is_canonical_override"):
            blockers.append("Slide 08 must include a caveat when it uses a non-received views metric.")
        if len((data.get("views_chart") or {}).get("rows") or []) > 1:
            blockers.append("Slide 08 must not show month-over-month comparison for non-received views metrics.")
        if "поле «просмотров получено» не заполнено" not in joined.lower():
            blockers.append("Slide 08 text must state that received views are not filled when using a fallback views field.")
    copy_qa = qa_slide_copy_principles(joined, slide_number=8)
    text_qa = qa_text_contract(joined, slide_number=8)
    blockers.extend(copy_qa["blockers"])
    blockers.extend(text_qa["blockers"])
    warnings.extend(copy_qa["warnings"])
    warnings.extend(text_qa["warnings"])
    return {
        "status": "ready_for_visual_review" if not blockers else "blocked",
        "blockers": blockers,
        "warnings": warnings,
        "checks": {
            "slide_count_12": result.get("slide_count") == 12,
            "total_materials_present": bool(total),
            "type_total_matches": type_total == total,
            "platform_total_matches": platform_total == total or (platform_data_total and platform_total == platform_data_total and platform_total < total),
            "no_stale_template_text": not result.get("old_text_left"),
            "no_broken_cyrillic": "????" not in joined,
            "no_semicolon": ";" not in joined,
            "source_note_without_filename": ".xlsx" not in source_note.lower() and ".pptx" not in source_note.lower(),
            "source_note_without_time": not any(time in source_note for time in ["00:00", "23:59", ":00", ":59"]),
            "source_links_preserved": not sample_link_count or result.get("hyperlink_count", 0) > 0,
            "views_metric_label": views_metrics.get("metric_label") == "Просмотры" and views_chart.get("metric_label") == "Просмотры",
            "platforms_sorted_by_views": not views_metrics.get("has_views") or [int(row.get("received_views", 0) or 0) for row in data.get("platform_rows") or [] if row.get("platform") != "Другие"] == sorted([int(row.get("received_views", 0) or 0) for row in data.get("platform_rows") or [] if row.get("platform") != "Другие"], reverse=True),
            "brand_bold_updates": result.get("brand_bold_updates", 0) > 0,
            "copy_principles": copy_qa["checks"],
            "text_contract": text_qa["checks"],
        },
        "visual_review": {
            "required": True,
            "status": "manual_review_needed",
            "notes": "Check that slide 8 visually matches the reference layout: title, summary block, left chart and right table remain in their inherited slots.",
        },
    }
